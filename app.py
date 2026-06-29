import os
import sys

# ── Confiança na CA corporativa (proxy de inspeção TLS da rede DHL) ──────────
# No venv isso vem do pip_system_certs via arquivo .pth, mas o PyInstaller NÃO
# executa .pth — então no EXE o httpx (usado pelo Supabase) não confia no cert
# do proxy e as chamadas HTTPS travam. Ativamos o truststore explicitamente aqui
# (usa o cofre de certificados do Windows) — vale no dev e no EXE. Precisa rodar
# ANTES de criar o cliente Supabase, pois o httpx monta o contexto SSL na criação.
try:
    import truststore as _truststore
    _truststore.inject_into_ssl()
except Exception:
    pass

import re
import base64
import textwrap
import smtplib
import threading
from collections import Counter
import random
import string
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from io import BytesIO
from datetime import datetime
from functools import wraps
import oracledb
import httpx
from supabase import create_client, ClientOptions
from flask import (
    Flask, render_template, request, redirect, url_for,
    flash, session, send_file, current_app, jsonify
)

from flask_sqlalchemy import SQLAlchemy
from sqlalchemy.orm import defer
from sqlalchemy.orm.attributes import set_committed_value
from sqlalchemy import func, case, text, or_
from datetime import date
from werkzeug.security import generate_password_hash, check_password_hash

from PIL import Image as _PILImage

# ── Importações pesadas — carregadas só na primeira chamada ───────────────────
def _import_docx():
    """Lazy-load python-docx no namespace global (no-op após a primeira chamada)."""
    global Document, Cm, Inches, Pt, RGBColor
    global WD_ALIGN_PARAGRAPH, WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
    global OxmlElement, qn
    from docx import Document
    from docx.shared import Cm, Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    globals()['_import_docx'] = lambda: None


def _import_openpyxl():
    """Lazy-load openpyxl no namespace global (no-op após a primeira chamada)."""
    global Workbook, Font, PatternFill, XLAlignment, _xl_col_letter
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment as XLAlignment
    from openpyxl.utils import get_column_letter as _xl_col_letter
    globals()['_import_openpyxl'] = lambda: None


def _import_reportlab():
    """Lazy-load reportlab no namespace global (no-op após a primeira chamada)."""
    global colors, A4, getSampleStyleSheet, ParagraphStyle, rcm
    global TA_CENTER, TA_RIGHT, ImageReader
    global SimpleDocTemplate, Paragraph, Spacer, RLImage, Table, TableStyle, PageBreak
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm as rcm
    from reportlab.lib.enums import TA_CENTER, TA_RIGHT
    from reportlab.lib.utils import ImageReader
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer,
        Image as RLImage, Table, TableStyle, PageBreak
    )
    globals()['_import_reportlab'] = lambda: None

if getattr(sys, 'frozen', False):
    BASE_DIR = sys._MEIPASS
    UPLOADS_ROOT = os.path.dirname(sys.executable)
else:
    BASE_DIR = os.path.abspath(os.path.dirname(__file__))
    UPLOADS_ROOT = os.path.join(BASE_DIR, 'static')

app = Flask(
    __name__,
    template_folder=os.path.join(BASE_DIR, 'templates'),
    static_folder=os.path.join(BASE_DIR, 'static'),
)
app.config["SECRET_KEY"] = "controle-ocorrencia-executivo"
app.config["SQLALCHEMY_DATABASE_URI"] = (
    "oracle+oracledb://SECPANEL:SEC003q2w3e4r2026"
    "@usqasap023-scan.phx-dc.dhl.com:1521"
    "/?service_name=SECPANEL"
)
app.config["SQLALCHEMY_TRACK_MODIFICATIONS"] = False
app.config["MAX_CONTENT_LENGTH"]   = 200 * 1024 * 1024  # 200 MB — limite de upload
app.config["MAX_FORM_MEMORY_SIZE"] = 16  * 1024 * 1024  # 16 MB  — formulários com muitos campos
app.config["SESSION_COOKIE_HTTPONLY"] = False
# Sessão permanente — usuário não precisa re-logar a cada abertura do app
from datetime import timedelta
app.config["PERMANENT_SESSION_LIFETIME"] = timedelta(days=30)
# Pool configurável por ambiente: o app desktop (1 usuário) usa o padrão enxuto;
# o modo servidor (serve.py, multi-cliente) eleva CCTV_POOL_SIZE para manter
# conexões "quentes" iguais ao nº de threads — evita abrir conexão nova no Oracle
# a cada pico (caro via VPN).
_POOL_SIZE    = int(os.environ.get("CCTV_POOL_SIZE", "5"))
_MAX_OVERFLOW = int(os.environ.get("CCTV_MAX_OVERFLOW", "5"))
app.config["SQLALCHEMY_ENGINE_OPTIONS"] = {
    "pool_size":    _POOL_SIZE,   # conexões permanentes no pool
    "max_overflow": _MAX_OVERFLOW,# conexões extras permitidas em pico
    "pool_recycle": 1800,    # recicla conexão a cada 30 min
    "pool_pre_ping": True,   # detecta conexões mortas antes de usar
    "pool_timeout": 60,      # aguarda até 60s por uma conexão (VPN tem latência maior)
    "connect_args": {
        "tcp_connect_timeout": 20,   # timeout TCP ao conectar no Oracle via VPN
    },
}

db = SQLAlchemy(app)

# ── Garante devolução de conexões ao pool mesmo em erros ─────────────────────
@app.teardown_appcontext
def _fechar_sessao_db(exc):
    if exc is not None:
        db.session.rollback()
    db.session.remove()

# ── Handler global para erros de banco (pool esgotado, VPN caída, etc.) ──────
from sqlalchemy.exc import OperationalError, TimeoutError as SATimeout
@app.errorhandler(OperationalError)
@app.errorhandler(SATimeout)
def _erro_banco(exc):
    import logging
    logging.error("Erro de banco de dados: %s", exc)
    db.session.rollback()
    from flask import flash, redirect, url_for, request as _req
    flash(
        "Erro de conexão com o banco de dados. "
        "Se estiver usando VPN, aguarde alguns segundos e tente novamente.",
        "danger"
    )
    destino = url_for("login") if not session.get("user_id") else url_for("ocorrencias")
    return redirect(destino)

@app.errorhandler(413)
def _arquivo_muito_grande(exc):
    """Exibe mensagem amigável quando o upload ultrapassa MAX_CONTENT_LENGTH (200 MB)."""
    from flask import flash, redirect, request as _req
    flash("O arquivo enviado é muito grande. O limite máximo é 200 MB.", "danger")
    ref = _req.referrer
    return redirect(ref) if ref else redirect(url_for("ocorrencias"))

# =========================
# PERFORMANCE — cache de estáticos + gzip
# =========================
import gzip as _gzip

@app.after_request
def _performance_headers(response):
    # Cache de 1 hora para arquivos estáticos (JS, CSS, imagens)
    if request.path.startswith('/static/'):
        response.headers['Cache-Control'] = 'public, max-age=3600, stale-while-revalidate=600'
        return response

    # Gzip para respostas HTML/JSON (só quando o cliente aceita)
    if (response.status_code in (200, 201)
            and not response.direct_passthrough
            and 'Content-Encoding' not in response.headers
            and 'gzip' in request.headers.get('Accept-Encoding', '').lower()):
        ct = response.content_type or ''
        if any(x in ct for x in ('text/html', 'application/json', 'text/css', 'javascript')):
            data = response.get_data()
            if len(data) > 800:          # não comprime respostas muito pequenas
                compressed = _gzip.compress(data, compresslevel=5)
                if len(compressed) < len(data):
                    response.set_data(compressed)
                    response.headers['Content-Encoding'] = 'gzip'
                    response.headers['Content-Length']   = len(compressed)
                    response.headers['Vary'] = 'Accept-Encoding'
    return response

# =========================
# LOGGING DE ERROS
# =========================
import logging, traceback as _tb

# No empacotamento one-folder o exe fica em <install>\app\ e a pasta 'app' é
# substituída a cada atualização. Grava o log no install root (pai de 'app') para
# que ele persista entre updates. Fallback para a pasta do exe se algo falhar.
if getattr(sys, 'frozen', False):
    _exe_dir = os.path.dirname(sys.executable)
    _parent  = os.path.dirname(_exe_dir)
    _log_dir = _parent if os.path.basename(_exe_dir).lower() == "app" else _exe_dir
else:
    _log_dir = BASE_DIR
try:
    os.makedirs(_log_dir, exist_ok=True)
except Exception:
    _log_dir = os.path.dirname(sys.executable) if getattr(sys, 'frozen', False) else BASE_DIR
logging.basicConfig(
    filename=os.path.join(_log_dir, 'cctv_error.log'),
    level=logging.ERROR,
    format='%(asctime)s %(levelname)s: %(message)s'
)

@app.errorhandler(500)
def _erro_500(e):
    trace = _tb.format_exc()
    logging.error(trace)
    return (
        f"<pre style='font-family:monospace;padding:20px;color:#b91c1c'>"
        f"Erro interno — reporte ao suporte:\n\n{trace}</pre>"
    ), 500


# =========================
# CONTROLE DE VERSÃO
# =========================
APP_VERSION = "4.2"

SMTP_HOST     = "smtp.dhl.com"
SMTP_PORT     = 25
EMAIL_FROM    = "Security.processassistant@dhl.com"
EMAIL_PASSWORD= "L0sspr3v3ntion@D3VT3AML4TAM"
EMAIL_DEVS    = [
    "deivid.martinsl@dhl.com",
    "Gilmar.SantosGJ@dhl.com",
    "anderson.rodriguesd@dhl.com",
]
EMAIL_BCC = "deivid.martinsl@dhl.com"   # cópia oculta em todos os e-mails automáticos

# =========================
# CONTROLE DE ATIVOS — usuários do app Ativos vivem no Supabase, não no Oracle
# =========================
def _proxy_corporativo():
    """Detecta o proxy local do Zscaler (DHL) lendo o AutoConfigURL (PAC) que o
    Windows usa para o navegador — só ISSO é configurado na rede DHL (não há
    proxy nem em variável de ambiente, nem no WinHTTP do sistema). httpx/supabase
    não sabem avaliar PAC, então extraímos host:porta direto da URL do PAC (ele
    sempre é servido pelo próprio proxy local, no mesmo host:porta). Sem isso, o
    httpx tenta conectar direto à internet e trava (ConnectTimeout) — funciona
    em dev porque o ambiente tinha HTTPS_PROXY exportado manualmente, mas o EXE
    real (aberto por atalho/Explorer) nunca tem essa variável.
    Retorna None em qualquer outro ambiente (sem Zscaler) — comportamento padrão
    do httpx (conexão direta) é preservado."""
    try:
        import winreg, re as _re
        key = winreg.OpenKey(winreg.HKEY_CURRENT_USER,
            r"Software\Microsoft\Windows\CurrentVersion\Internet Settings")
        url, _ = winreg.QueryValueEx(key, "AutoConfigURL")
        winreg.CloseKey(key)
        m = _re.match(r"https?://([\d.]+):(\d+)/", url or "")
        if m:
            return f"http://{m.group(1)}:{m.group(2)}"
    except Exception:
        pass
    return None


SUPABASE_URL = "https://zmjzptwwxixitqiiqspy.supabase.co"
SUPABASE_KEY = "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJpc3MiOiJzdXBhYmFzZSIsInJlZiI6InptanpwdHd3eGl4aXRxaWlxc3B5Iiwicm9sZSI6InNlcnZpY2Vfcm9sZSIsImlhdCI6MTc2MDM3MTU0MSwiZXhwIjoyMDc1OTQ3NTQxfQ.sOENYPMqOk4fEth6heMw3-SM2QluqIdd-YCFPuoU73Y"
_sb_proxy = _proxy_corporativo()
if _sb_proxy:
    sb = create_client(SUPABASE_URL, SUPABASE_KEY,
        options=ClientOptions(httpx_client=httpx.Client(proxy=_sb_proxy, timeout=15)))
else:
    sb = create_client(SUPABASE_URL, SUPABASE_KEY)

# =========================
# CONTROLE DE VERSÃO — verifica SISTEMA_CONFIG.VERSAO_EXIGIDA no banco
# Para bloquear uma versão antiga: UPDATE SISTEMA_CONFIG SET VERSAO_EXIGIDA = 'X.X'
# =========================
_v_cache: dict = {}

def _get_versao_banco():
    """Retorna (ok, versao_banco). Usa cache após primeira consulta bem-sucedida."""
    if _v_cache:
        return _v_cache["ok"], _v_cache["db_ver"]
    result = {}
    def _query():
        try:
            row = db.session.execute(
                db.text("SELECT VERSAO_EXIGIDA FROM SISTEMA_CONFIG WHERE ROWNUM = 1")
            ).fetchone()
            result["db_ver"] = (row[0] or "").strip() if row else "?"
            result["ok"]     = result["db_ver"] == APP_VERSION
        except Exception:
            result["ok"]     = True   # fail-open: se banco inacessível, não bloqueia
            result["db_ver"] = "?"
    import threading as _th
    t = _th.Thread(target=_query, daemon=True)
    t.start()
    t.join(timeout=10)
    ok     = result.get("ok",     True)
    db_ver = result.get("db_ver", "?")
    if db_ver != "?":            # só armazena cache se o banco respondeu
        _v_cache["ok"]     = ok
        _v_cache["db_ver"] = db_ver
    return ok, db_ver

ROTAS_LIVRES = {"static", "versao_bloqueada"}

@app.before_request
def _verificar_versao():
    if request.endpoint in ROTAS_LIVRES or request.endpoint is None:
        return
    ok, _ = _get_versao_banco()
    if not ok:
        return redirect(url_for("versao_bloqueada"))

@app.route("/versao-desatualizada")
def versao_bloqueada():
    _, db_ver = _get_versao_banco()
    return render_template("versao_bloqueada.html",
                           v_atual=APP_VERSION, v_nova=db_ver)


def _parse_valor(s):
    """Converte string monetária para float. Ex: 'R$ 1.234,56' → 1234.56"""
    if not s:
        return 0.0
    s = s.strip()
    # Remove prefixos (R$, $, €, etc.)
    s = re.sub(r'[^\d.,]', '', s)
    # Padrão BR: 1.234,56  →  remove ponto de milhar, troca vírgula por ponto
    if ',' in s and '.' in s:
        if s.rfind(',') > s.rfind('.'):
            s = s.replace('.', '').replace(',', '.')
        else:
            s = s.replace(',', '')
    elif ',' in s:
        s = s.replace(',', '.')
    try:
        return float(s)
    except ValueError:
        return 0.0


def _formatar_valor(v):
    """Formata float como moeda BRL. Ex: 1234.56 → 'R$ 1.234,56'"""
    if v == 0:
        return "—"
    return f"R$ {v:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")


def _sites_do_usuario(user_id=None, user_site=None):
    """
    Retorna a lista de sites que o usuário logado pode visualizar.

    Regras (em ordem de prioridade):
      1. ADMIN             → [] (sem filtro — vê tudo)
      2. Usuário com vínculos em USUARIO_SITES → apenas os sites vinculados
      3. Usuário sem vínculos → [user_site] (apenas o próprio site)
    """
    if user_id is None:
        user_id   = session.get("user_id")
    if user_site is None:
        user_site = session.get("user_site", "")
    perfil = (session.get("user_perfil") or "").upper()

    if perfil == "ADMIN":
        return []   # admin não filtra por site

    # Verifica se o usuário tem sites vinculados explicitamente
    if user_id:
        vinculos = UsuarioSite.query.filter_by(usuario_id=user_id).all()
        if vinculos:
            return [v.site_nome for v in vinculos]

    # Sem vínculos: usa o site padrão do usuário
    return [user_site] if user_site else []


def _query_filtrar_sites(query, model_class, user_id=None, user_site=None):
    """
    Aplica filtro de site(s) em uma query SQLAlchemy.
    Retorna a query filtrada (ou sem filtro para ADMIN).
    """
    sites = _sites_do_usuario(user_id=user_id, user_site=user_site)
    if not sites:          # admin ou sem site definido
        return query
    if len(sites) == 1:
        return query.filter(model_class.site == sites[0])
    return query.filter(model_class.site.in_(sites))



ALLOWED_EXTENSIONS = {"png", "jpg", "jpeg", "webp"}
def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS

EXTENSOES_PERMITIDAS_IMAGEM = {"png", "jpg", "jpeg", "webp"}
EXTENSOES_PERMITIDAS_POST = {"png", "jpg", "jpeg", "webp", "pdf", "doc", "docx", "xlsx"}


# =========================
# MODELS
# =========================
class Usuario(db.Model):
    __tablename__ = "USERS_LIVRO"

    id = db.Column("ID", db.Integer, db.Identity(start=1), primary_key=True)
    nome = db.Column("NOME", db.String(120), nullable=False)
    email = db.Column("EMAIL", db.String(120), unique=True, nullable=False, index=True)
    password_hash = db.Column("PASSWORD_HASH", db.String(255), nullable=False)
    perfil = db.Column("ROLE", db.String(30), nullable=False, default="OPERACIONAL")
    site = db.Column("SITE", db.String(80), nullable=True)
    is_active = db.Column("IS_ACTIVE", db.Boolean, nullable=False, default=True)
    created_at = db.Column("CREATED_AT", db.DateTime, nullable=False, default=datetime.utcnow)
    foto_perfil  = db.Column("FOTO_PERFIL",  db.Text,       nullable=True)
    tem_foto     = db.Column("TEM_FOTO",     db.String(1),  nullable=True, default="N")
    lgpd_aceito  = db.Column("LGPD_ACEITO",  db.String(3),  nullable=True, default=None)
    lgpd_aceito_em = db.Column("LGPD_ACEITO_EM", db.DateTime, nullable=True, default=None)
    cargo        = db.Column("CARGO",        db.String(120), nullable=True)

    def set_password(self, senha: str):
        self.password_hash = generate_password_hash(senha)

    def check_password(self, senha: str) -> bool:
        return check_password_hash(self.password_hash, senha)


class SyncUsuariosLedger(db.Model):
    """Estado da última sincronização de logins entre o Oracle (USERS_LIVRO) e o
    Supabase (usuarios). Guarda, por e-mail, o hash de senha já "acordado" entre os
    dois lados — é o que permite o sync BIDIRECIONAL saber qual lado mudou desde a
    última vez (last-write-wins) e detectar exclusões com segurança."""
    __tablename__ = "SYNC_USUARIOS_LEDGER"

    email         = db.Column("EMAIL",         db.String(120), primary_key=True)
    password_hash = db.Column("PASSWORD_HASH", db.String(255), nullable=True)
    synced_at     = db.Column("SYNCED_AT",     db.DateTime,    nullable=True, default=datetime.utcnow)


class ResetToken(db.Model):
    __tablename__ = "RESET_TOKENS"
    id        = db.Column(db.Integer, db.Identity(start=1), primary_key=True)
    user_id   = db.Column(db.Integer, nullable=False)
    token     = db.Column(db.String(6), nullable=False)
    expira_em = db.Column(db.DateTime, nullable=False)
    usado     = db.Column(db.Integer, nullable=False, default=0)


class SolicitacaoCadastro(db.Model):
    __tablename__ = "SOLICITACOES_CADASTRO"
    id         = db.Column(db.Integer, db.Identity(start=1), primary_key=True)
    nome       = db.Column(db.String(120), nullable=False)
    email      = db.Column(db.String(120), nullable=False)
    site       = db.Column(db.String(128), nullable=True)
    status     = db.Column(db.String(20),  nullable=False, default="PENDENTE")
    criado_em  = db.Column(db.DateTime,    nullable=False, default=datetime.utcnow)


class UsuarioSite(db.Model):
    """Vínculo entre usuário OVERHEAD e os sites que ele pode acessar."""
    __tablename__ = "USUARIO_SITES"
    id          = db.Column(db.Integer, db.Identity(start=1), primary_key=True)
    usuario_id  = db.Column(db.Integer, db.ForeignKey("USERS_LIVRO.ID"), nullable=False, index=True)
    site_nome   = db.Column(db.String(128), nullable=False)


class AppRelease(db.Model):
    """Cada linha é uma versão publicada do CCTV_ControlPanel.exe."""
    __tablename__ = "APP_RELEASES"
    id            = db.Column(db.Integer, db.Identity(start=1), primary_key=True)
    versao        = db.Column(db.String(20),  nullable=False)
    nome_arquivo  = db.Column(db.String(255), nullable=True)
    tamanho_bytes = db.Column(db.Integer,     nullable=True)
    exe_blob      = db.Column(db.LargeBinary, nullable=True)
    publicado_em  = db.Column(db.DateTime,    nullable=False, default=datetime.utcnow)
    publicado_por = db.Column(db.String(120), nullable=True)
    ativo         = db.Column(db.String(1),   nullable=False, default='N')


def _fmt_size(n):
    """Formata bytes em KB / MB."""
    if n is None:
        return "—"
    if n < 1024:
        return f"{n} B"
    if n < 1024 ** 2:
        return f"{n / 1024:.1f} KB"
    return f"{n / 1024 ** 2:.1f} MB"


class SiteCompleto(db.Model):
    __tablename__ = "SITES_COMPLETO"

    nome_do_site = db.Column("NOME_DO_SITE", db.String(128), primary_key=True)
    endereco     = db.Column("ENDEREÇO", db.String(255), nullable=True)
    cidade       = db.Column("CIDADE", db.String(50), nullable=True)
    estado       = db.Column("ESTADO", db.String(2), nullable=True)
    pais         = db.Column("PAÍS", db.String(26), nullable=True)
    responsavel_security = db.Column("RESPONSÁVEL_SECURITY", db.String(128), nullable=True)
    coordenador  = db.Column("COORDENADOR", db.String(26), nullable=True)
    latitude     = db.Column("LATIDUDE", db.Numeric(38, 0), nullable=True)
    longitude    = db.Column("LONGITUDE", db.Numeric(38, 0), nullable=True)
    sector       = db.Column("SECTOR", db.String(26), nullable=True)
    security_responsible = db.Column("SECURITY_RESPONSIBLE", db.String(26), nullable=True)


class AnaliseInvestigativa(db.Model):
    __tablename__ = "analises_investigativas"

    id = db.Column(db.Integer, db.Identity(start=1), primary_key=True)
    codigo = db.Column(db.String(30), nullable=True, unique=True)
    numero_site = db.Column(db.Integer, nullable=True)
    site = db.Column(db.String(128), nullable=True)

    id_relatorio = db.Column(db.String(30), nullable=True)
    empresa = db.Column(db.String(120), nullable=True)
    unidade = db.Column(db.String(180), nullable=True)
    endereco = db.Column(db.String(255), nullable=True)
    classificacao = db.Column(db.String(80), nullable=True)
    produtos_segmento = db.Column(db.String(120), nullable=True)
    clientes = db.Column(db.String(120), nullable=True)

    objetivo = db.Column(db.Text, nullable=True)
    responsavel = db.Column(db.String(150), nullable=True)
    nome_funcao_data = db.Column(db.String(255), nullable=True)   # legado — mantido para registros antigos
    funcao_levantamento = db.Column(db.String(255), nullable=True)
    data_levantamento = db.Column(db.String(100), nullable=True)

    descricao_registro = db.Column(db.Text, nullable=True)
    conclusao = db.Column(db.Text, nullable=True)
    sugestao = db.Column(db.Text, nullable=True)

    criado_por = db.Column(db.String(120), nullable=True)
    criado_em = db.Column(db.DateTime, default=datetime.utcnow)
    docx_arquivo = db.Column(db.LargeBinary, nullable=True)
    valor = db.Column(db.String(50), nullable=True)

    # ── Fechamento ──────────────────────────────────────────────
    status_analise = db.Column(db.String(30), nullable=True, default="EM ANDAMENTO")
    texto_fechamento = db.Column(db.Text, nullable=True)
    fechado_por = db.Column(db.String(120), nullable=True)
    fechado_em = db.Column(db.DateTime, nullable=True)
    anexo_fechamento_nome = db.Column(db.String(255), nullable=True)
    anexo_fechamento = db.Column(db.LargeBinary, nullable=True)

    imagens = db.relationship(
        "ImagemAnaliseInvestigativa",
        backref="analise",
        cascade="all, delete-orphan",
        lazy=True
    )


class ImagemAnaliseInvestigativa(db.Model):
    __tablename__ = "imagens_analises_investigativas"

    id = db.Column(db.Integer, db.Identity(start=1), primary_key=True)
    analise_id = db.Column(db.Integer, db.ForeignKey("analises_investigativas.id"), nullable=False)

    arquivo = db.Column(db.String(255), nullable=True)       # legado (não usado em novos registros)
    arquivo_b64 = db.Column(db.Text, nullable=True)          # base64 completo da imagem (CLOB)
    descricao = db.Column(db.Text, nullable=True)

    criado_em = db.Column(db.DateTime, default=datetime.utcnow)

    @property
    def b64(self):
        """Retorna o base64 da imagem como string.
        O Oracle pode devolver CLOBs como objeto LOB (thin driver) — lê com .read() se necessário."""
        raw = self.arquivo_b64 or self.arquivo or ""
        if hasattr(raw, "read"):   # cx_Oracle / oracledb LOB object
            raw = raw.read()
        return raw or ""


class NaturezaOcorrencia(db.Model):
    __tablename__ = "NATUREZAS_OCORRENCIA"
    id   = db.Column(db.Integer, db.Identity(start=1), primary_key=True)
    nome = db.Column(db.String(120), nullable=False, unique=True)

    def __repr__(self):
        return f"<NaturezaOcorrencia {self.nome}>"


_NATUREZAS_PADRAO = [
    "Acidente", "Agressão", "Armazenamento incorreto", "Assédio",
    "Ato inseguro", "Avaria de produto", "Colisão", "Condição insegura",
    "Dano ao Patrimônio", "Desinteligência", "Desvio de conduta",
    "Desvio de processo", "Falta de volume", "Furto",
    "Pacote no descarte", "Pacote violado", "Roubo", "Vandalismo",
]


# =========================
# CONFIG. DE CAMPOS — MODELS
# =========================
class NaturezaConfig(db.Model):
    """Naturezas configuradas por site."""
    __tablename__ = "NATUREZA_CONFIG"
    id              = db.Column(db.Integer, db.Identity(start=1), primary_key=True)
    natureza        = db.Column(db.String(200), nullable=False)
    site            = db.Column(db.String(128), nullable=False, index=True)
    data_criacao    = db.Column(db.DateTime, default=datetime.utcnow)
    usuario_criador = db.Column(db.String(120), nullable=True)

class LocalConfig(db.Model):
    """Locais configurados por site."""
    __tablename__ = "LOCAL_CONFIG"
    id              = db.Column(db.Integer, db.Identity(start=1), primary_key=True)
    local           = db.Column(db.String(200), nullable=False)
    site            = db.Column(db.String(128), nullable=False, index=True)
    data_criacao    = db.Column(db.DateTime, default=datetime.utcnow)
    usuario_criador = db.Column(db.String(120), nullable=True)

class SetorConfig(db.Model):
    """Setores configurados por site."""
    __tablename__ = "SETOR_CONFIG"
    id              = db.Column(db.Integer, db.Identity(start=1), primary_key=True)
    setor           = db.Column(db.String(200), nullable=False)
    site            = db.Column(db.String(128), nullable=False, index=True)
    data_criacao    = db.Column(db.DateTime, default=datetime.utcnow)
    usuario_criador = db.Column(db.String(120), nullable=True)


def _get_naturezas(site=None):
    """Retorna lista de naturezas: padrão (hardcode) + adições criadas pelos gestores/keyusers.
    Os valores padrão são sempre mantidos; o DB só adiciona, nunca substitui."""
    try:
        rows = db.session.query(NaturezaOcorrencia.nome).all()
        db_nomes = {r[0] for r in rows}
        return sorted(set(_NATUREZAS_PADRAO) | db_nomes)
    except Exception:
        return sorted(_NATUREZAS_PADRAO)

def _get_locais(site=None):
    """Retorna lista de locais configurados para o site."""
    try:
        q = LocalConfig.query
        if site:
            q = q.filter_by(site=site)
        return [c.local for c in q.order_by(LocalConfig.local).all()]
    except Exception:
        return []

def _get_setores(site=None):
    """Retorna lista de setores configurados para o site."""
    try:
        q = SetorConfig.query
        if site:
            q = q.filter_by(site=site)
        return [c.setor for c in q.order_by(SetorConfig.setor).all()]
    except Exception:
        return []


# =========================
# CHECKLIST DE CÂMERAS — MODELS
# =========================

class Camera(db.Model):
    __tablename__ = "CAMERAS"
    id        = db.Column(db.Integer, db.Identity(start=1), primary_key=True)
    site      = db.Column(db.String(100), nullable=False, index=True)
    numero    = db.Column(db.String(50),  nullable=False)   # identificador: "CAM-001", "01", etc.
    nome      = db.Column(db.String(120), nullable=True)    # nome amigável opcional
    ativo     = db.Column(db.Integer,     default=1)        # 1=ativa  0=desativada
    criado_em = db.Column(db.DateTime,    default=datetime.utcnow)


class ChecklistCamera(db.Model):
    __tablename__ = "CHECKLIST_CAMERAS"
    id              = db.Column(db.Integer, db.Identity(start=1), primary_key=True)
    site            = db.Column(db.String(100), nullable=False)
    data_checklist  = db.Column(db.String(10),  nullable=False)   # YYYY-MM-DD
    hora_checklist  = db.Column(db.String(5))                     # HH:MM
    operador        = db.Column(db.String(120))
    total           = db.Column(db.Integer, default=0)
    funcionais      = db.Column(db.Integer, default=0)
    inoperantes     = db.Column(db.Integer, default=0)
    tratativa       = db.Column(db.Text)
    anexo_nome      = db.Column(db.String(255))
    anexo_dados     = db.Column(db.LargeBinary)
    criado_em       = db.Column(db.DateTime, default=datetime.utcnow)
    itens           = db.relationship("ChecklistCameraItem", backref="checklist",
                                      lazy="dynamic", cascade="all, delete-orphan")


class ChecklistCameraItem(db.Model):
    __tablename__ = "CHECKLIST_CAMERA_ITEM"
    id           = db.Column(db.Integer, db.Identity(start=1), primary_key=True)
    checklist_id = db.Column(db.Integer, db.ForeignKey("CHECKLIST_CAMERAS.id"), nullable=False)
    camera_id    = db.Column(db.Integer, db.ForeignKey("CAMERAS.id"), nullable=False)
    status       = db.Column(db.String(20), nullable=False, default="FUNCIONAL")
    motivo       = db.Column(db.Text)
    camera       = db.relationship("Camera")


# =========================
# CHECKLIST DE CÂMERAS — PDF
# =========================

def _gerar_pdf_checklist_cameras(chk_id):
    """Gera PDF profissional do checklist de câmeras.
    Suporta 1000+ câmeras com layout compacto multi-coluna."""
    _import_reportlab()

    chk    = ChecklistCamera.query.get_or_404(chk_id)
    itens  = chk.itens.order_by(ChecklistCameraItem.id).all()

    buf = BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=A4,
        leftMargin=1.5*rcm, rightMargin=1.5*rcm,
        topMargin=2*rcm, bottomMargin=2*rcm,
        title=f"Checklist Câmeras — {chk.site} — {chk.data_checklist}",
    )

    W, H = A4
    usable_w = W - 3*rcm   # ~174mm

    # ── paleta ──────────────────────────────────────────────────────────────
    C_RED    = colors.HexColor("#d40511")
    C_DARK   = colors.HexColor("#1f2937")
    C_GREEN  = colors.HexColor("#16a34a")
    C_AMBER  = colors.HexColor("#d97706")
    C_YELW   = colors.HexColor("#ffcc00")
    C_BGRED  = colors.HexColor("#fee2e2")
    C_BGGRN  = colors.HexColor("#dcfce7")
    C_BGLT   = colors.HexColor("#f8fafc")
    C_LINE   = colors.HexColor("#e5e7eb")

    styles = getSampleStyleSheet()
    sN = ParagraphStyle("N",  fontName="Helvetica",      fontSize=8,  leading=10, textColor=C_DARK)
    sB = ParagraphStyle("B",  fontName="Helvetica-Bold", fontSize=8,  leading=10, textColor=C_DARK)
    sR = ParagraphStyle("R",  fontName="Helvetica-Bold", fontSize=8,  leading=10, textColor=C_RED)
    sG = ParagraphStyle("G",  fontName="Helvetica-Bold", fontSize=8,  leading=10, textColor=C_GREEN)
    sSm= ParagraphStyle("Sm", fontName="Helvetica",      fontSize=7,  leading=9,  textColor=colors.HexColor("#6b7280"))
    sIt= ParagraphStyle("It", fontName="Helvetica-Oblique", fontSize=7, leading=9, textColor=colors.HexColor("#6b7280"))

    story = []

    # ═══════════════════════════════════════════════════════════════
    # CABEÇALHO / CAPA
    # ═══════════════════════════════════════════════════════════════
    pct_func = round(chk.funcionais / chk.total * 100, 1) if chk.total else 0
    pct_inop = round(chk.inoperantes / chk.total * 100, 1) if chk.total else 0

    hdr_data = [[
        Paragraph("<b><font color='#d40511' size=14>CHECKLIST DE CÂMERAS</font></b><br/>"
                  f"<font size=9 color='#6b7280'>{chk.site} &nbsp;|&nbsp; "
                  f"{chk.data_checklist} &nbsp;|&nbsp; {chk.hora_checklist or ''}</font>", styles["Normal"]),
        Paragraph(f"<b>Operador:</b> {chk.operador or '—'}", sN),
    ]]
    hdr_tbl = Table(hdr_data, colWidths=[usable_w*0.65, usable_w*0.35])
    hdr_tbl.setStyle(TableStyle([
        ("BACKGROUND",  (0,0), (-1,-1), C_DARK),
        ("TEXTCOLOR",   (0,0), (-1,-1), colors.white),
        ("PADDING",     (0,0), (-1,-1), 10),
        ("VALIGN",      (0,0), (-1,-1), "MIDDLE"),
        ("ROUNDEDCORNERS", [6,6,6,6]),
    ]))
    story.append(hdr_tbl)
    story.append(Spacer(1, 0.4*rcm))

    # ── KPI Summary bar ──────────────────────────────────────────────────────
    kpi_data = [[
        Paragraph(f"<b>{chk.total}</b><br/><font size=7 color='#6b7280'>Total</font>", sB),
        Paragraph(f"<b><font color='#16a34a'>{chk.funcionais}</font></b><br/>"
                  f"<font size=7 color='#6b7280'>Funcionais ({pct_func}%)</font>", sB),
        Paragraph(f"<b><font color='#d40511'>{chk.inoperantes}</font></b><br/>"
                  f"<font size=7 color='#6b7280'>Inoperantes ({pct_inop}%)</font>", sB),
        Paragraph(f"<b><font color='{'#16a34a' if pct_func>=90 else '#d97706' if pct_func>=70 else '#d40511'}'>"
                  f"{'Ótimo' if pct_func>=90 else 'Atenção' if pct_func>=70 else 'Crítico'}</font></b><br/>"
                  f"<font size=7 color='#6b7280'>Status Geral</font>", sB),
    ]]
    kpi_tbl = Table(kpi_data, colWidths=[usable_w/4]*4)
    kpi_tbl.setStyle(TableStyle([
        ("BACKGROUND",  (0,0), (-1,-1), C_BGLT),
        ("BOX",         (0,0), (-1,-1), 0.5, C_LINE),
        ("INNERGRID",   (0,0), (-1,-1), 0.5, C_LINE),
        ("ALIGN",       (0,0), (-1,-1), "CENTER"),
        ("VALIGN",      (0,0), (-1,-1), "MIDDLE"),
        ("PADDING",     (0,0), (-1,-1), 8),
    ]))
    story.append(kpi_tbl)
    story.append(Spacer(1, 0.35*rcm))

    # ── Seção de inoperantes (com motivo) ────────────────────────────────────
    inoperantes = [i for i in itens if i.status == "INOPERANTE"]
    if inoperantes:
        story.append(Paragraph(
            "<b><font color='#d40511'>▸ CÂMERAS INOPERANTES</font></b>",
            ParagraphStyle("SI", fontName="Helvetica-Bold", fontSize=9, leading=12,
                           textColor=C_RED, spaceBefore=4, spaceAfter=4)
        ))
        inop_rows = [[
            Paragraph("<b>Câmera</b>", sB),
            Paragraph("<b>Nome</b>", sB),
            Paragraph("<b>Motivo / Observação</b>", sB),
        ]]
        for item in inoperantes:
            cam = item.camera
            inop_rows.append([
                Paragraph(cam.numero, sR),
                Paragraph(cam.nome or "—", sSm),
                Paragraph(item.motivo or "—", sIt),
            ])
        inop_tbl = Table(inop_rows,
                         colWidths=[usable_w*0.14, usable_w*0.24, usable_w*0.62],
                         repeatRows=1)
        inop_tbl.setStyle(TableStyle([
            ("BACKGROUND",  (0,0), (-1,0),  C_RED),
            ("TEXTCOLOR",   (0,0), (-1,0),  colors.white),
            ("BACKGROUND",  (0,1), (-1,-1), C_BGRED),
            ("ROWBACKGROUNDS", (0,1), (-1,-1), [C_BGRED, colors.HexColor("#fff5f5")]),
            ("BOX",         (0,0), (-1,-1), 0.5, C_RED),
            ("INNERGRID",   (0,0), (-1,-1), 0.3, colors.HexColor("#fecaca")),
            ("ALIGN",       (0,0), (-1,-1), "LEFT"),
            ("VALIGN",      (0,0), (-1,-1), "MIDDLE"),
            ("PADDING",     (0,0), (-1,-1), 5),
            ("WORDWRAP",    (2,1), (2,-1),  True),
        ]))
        story.append(inop_tbl)
        story.append(Spacer(1, 0.35*rcm))

    # ── Tratativa ────────────────────────────────────────────────────────────
    if chk.tratativa:
        story.append(Paragraph(
            "<b><font color='#1f2937'>▸ TRATATIVA / PLANO DE AÇÃO</font></b>",
            ParagraphStyle("ST", fontName="Helvetica-Bold", fontSize=9, leading=12,
                           textColor=C_DARK, spaceBefore=4, spaceAfter=4)
        ))
        trat_data = [[Paragraph(chk.tratativa, sN)]]
        trat_tbl  = Table(trat_data, colWidths=[usable_w])
        trat_tbl.setStyle(TableStyle([
            ("BACKGROUND", (0,0), (-1,-1), C_BGLT),
            ("BOX",        (0,0), (-1,-1), 0.5, C_LINE),
            ("PADDING",    (0,0), (-1,-1), 8),
        ]))
        story.append(trat_tbl)
        story.append(Spacer(1, 0.35*rcm))

    if chk.anexo_nome:
        story.append(Paragraph(f"📎 Anexo: {chk.anexo_nome}", sSm))
        story.append(Spacer(1, 0.2*rcm))

    # ═══════════════════════════════════════════════════════════════
    # LISTAGEM COMPLETA — layout 3 colunas (câmera por câmera)
    # ═══════════════════════════════════════════════════════════════
    story.append(PageBreak())
    story.append(Paragraph(
        "<b><font color='#1f2937'>▸ LISTAGEM COMPLETA DAS CÂMERAS</font></b>",
        ParagraphStyle("SL", fontName="Helvetica-Bold", fontSize=9, leading=12,
                       textColor=C_DARK, spaceBefore=0, spaceAfter=6)
    ))

    # Cabeçalho da tabela multi-coluna (3 câmeras por linha)
    COLS = 3
    # colWidths por bloco: Num(18mm) | Nome(30mm) | Status(17mm) = 65mm  × 3 = 195mm → ajusta
    bw_num  = usable_w * 0.11
    bw_nome = usable_w * 0.21
    bw_st   = usable_w * 0.13
    col_ws  = ([bw_num, bw_nome, bw_st] * COLS)

    cam_hdr = (
        [Paragraph("<b>Câm.</b>", sB), Paragraph("<b>Nome</b>", sB), Paragraph("<b>Status</b>", sB)] * COLS
    )
    cam_rows = [cam_hdr]

    chunk = []
    for item in itens:
        cam   = item.camera
        s_par = Paragraph("● FUNCIONAL" if item.status == "FUNCIONAL" else "● INOP.", sG if item.status == "FUNCIONAL" else sR)
        chunk.append([
            Paragraph(cam.numero, sN),
            Paragraph(cam.nome or "—", sSm),
            s_par,
        ])
        if len(chunk) == COLS:
            cam_rows.append([c for cell in chunk for c in cell])
            chunk = []

    if chunk:  # última linha incompleta → preenche células vazias
        while len(chunk) < COLS:
            chunk.append([Paragraph("", sN), Paragraph("", sN), Paragraph("", sN)])
        cam_rows.append([c for cell in chunk for c in cell])

    cam_tbl = Table(cam_rows, colWidths=col_ws, repeatRows=1)

    # estilos base
    ts = [
        ("BACKGROUND",     (0,0), (-1,0),  C_DARK),
        ("TEXTCOLOR",      (0,0), (-1,0),  colors.white),
        ("FONT",           (0,0), (-1,0),  "Helvetica-Bold", 7),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.white, C_BGLT]),
        ("BOX",            (0,0), (-1,-1), 0.4, C_LINE),
        ("INNERGRID",      (0,0), (-1,-1), 0.3, C_LINE),
        ("ALIGN",          (0,0), (-1,-1), "LEFT"),
        ("VALIGN",         (0,0), (-1,-1), "MIDDLE"),
        ("PADDING",        (0,0), (-1,-1), 4),
        ("FONTSIZE",       (0,1), (-1,-1), 7),
    ]

    # destacar linhas com câmera INOPERANTE — vermelho suave
    for row_idx, item in enumerate(itens, start=1):
        if item.status == "INOPERANTE":
            # cada item ocupa 1/COLS de uma linha — localiza a linha e coluna correta
            tbl_row = (row_idx - 1) // COLS + 1
            col_off = ((row_idx - 1) % COLS) * 3
            ts.append(("BACKGROUND", (col_off, tbl_row), (col_off+2, tbl_row), C_BGRED))

    cam_tbl.setStyle(TableStyle(ts))
    story.append(cam_tbl)

    # ── rodapé de assinatura ─────────────────────────────────────────────────
    story.append(Spacer(1, 0.8*rcm))
    sig_data = [[
        Paragraph(f"<b>Operador:</b> {chk.operador or '—'}", sN),
        Paragraph("Assinatura: ______________________________", sN),
        Paragraph(f"<b>Data/Hora:</b> {chk.data_checklist} {chk.hora_checklist or ''}", sN),
    ]]
    sig_tbl = Table(sig_data, colWidths=[usable_w*0.34, usable_w*0.4, usable_w*0.26])
    sig_tbl.setStyle(TableStyle([
        ("BOX",     (0,0), (-1,-1), 0.5, C_LINE),
        ("INNERGRID",(0,0),(-1,-1), 0.5, C_LINE),
        ("PADDING", (0,0), (-1,-1), 8),
        ("VALIGN",  (0,0), (-1,-1), "MIDDLE"),
    ]))
    story.append(sig_tbl)

    def _footer(canvas, doc):
        canvas.saveState()
        canvas.setFont("Helvetica", 7)
        canvas.setFillColor(colors.HexColor("#9ca3af"))
        canvas.drawString(1.5*rcm, 1.2*rcm, f"DHL Security — CCTV Control Panel — Checklist: {chk.site} / {chk.data_checklist}")
        canvas.drawRightString(W - 1.5*rcm, 1.2*rcm, f"Pág. {doc.page}")
        canvas.restoreState()

    doc.build(story, onFirstPage=_footer, onLaterPages=_footer)
    buf.seek(0)
    return buf


class ANC(db.Model):
    __tablename__ = "ancs"

    id = db.Column(db.Integer, db.Identity(start=1), primary_key=True)
    numero_anc = db.Column(db.String(50), nullable=True, unique=True)
    data_nc = db.Column(db.String(10), nullable=True)
    hora_nc = db.Column(db.String(5), nullable=True)
    site = db.Column(db.String(120), nullable=True)
    setor = db.Column(db.String(120), nullable=True)
    tipo_ocorrencia = db.Column(db.String(120), nullable=True)
    gravidade = db.Column(db.String(30), nullable=True)
    natureza = db.Column(db.String(255), nullable=True)
    responsavel = db.Column(db.String(120), nullable=True)
    gestor_responsavel = db.Column(db.String(120), nullable=True)
    local = db.Column(db.String(120), nullable=True)
    envolvido = db.Column(db.String(255), nullable=True)
    tipo = db.Column(db.String(80), nullable=True)
    cargo = db.Column(db.String(120), nullable=True)
    turno = db.Column(db.String(20), nullable=True)
    status = db.Column(db.String(30), nullable=False, default="ABERTO")
    descricao = db.Column(db.Text, nullable=True)
    inicio_investigacao = db.Column(db.String(16), nullable=True)
    fim_investigacao = db.Column(db.String(16), nullable=True)
    imagem_1 = db.Column(db.Text, nullable=True)
    imagem_2 = db.Column(db.Text, nullable=True)
    imagem_3 = db.Column(db.Text, nullable=True)
    imagem_4 = db.Column(db.Text, nullable=True)
    imagem_5 = db.Column(db.Text, nullable=True)
    imagem_6 = db.Column(db.Text, nullable=True)
    numero_site = db.Column(db.Integer, nullable=True)
    criado_por = db.Column(db.String(80), nullable=True)
    criado_em = db.Column(db.DateTime, default=datetime.utcnow)
    atualizado_em = db.Column(db.DateTime, default=datetime.utcnow, onupdate=datetime.utcnow)

    valor = db.Column(db.String(50), nullable=True)

    # ── Fechamento / Plano de Ação ────────────────────────────────
    plano_acao_texto      = db.Column(db.Text, nullable=True)
    fechado_por           = db.Column(db.String(120), nullable=True)
    fechado_em            = db.Column(db.DateTime, nullable=True)
    anexo_fechamento_nome = db.Column(db.String(255), nullable=True)
    anexo_fechamento      = db.Column(db.LargeBinary, nullable=True)

    # ── Soft-delete / Solicitação de Exclusão ─────────────────────
    excluido              = db.Column(db.String(1),   nullable=True, default="N")
    excl_status           = db.Column(db.String(20),  nullable=True)          # PENDENTE | APROVADO | REJEITADO
    excl_solicitado_por   = db.Column(db.String(120), nullable=True)
    excl_solicitado_em    = db.Column(db.DateTime,    nullable=True)
    excl_solicitante_email= db.Column(db.String(120), nullable=True)
    excl_motivo           = db.Column(db.String(500), nullable=True)
    excl_admin_por        = db.Column(db.String(120), nullable=True)
    excl_admin_em         = db.Column(db.DateTime,    nullable=True)
    excl_admin_motivo     = db.Column(db.String(500), nullable=True)


class Ocorrencia(db.Model):
    __tablename__ = "ocorrencias"

    id = db.Column(db.Integer, db.Identity(start=1), primary_key=True)
    codigo = db.Column(db.String(30), nullable=True, unique=True)
    numero_site = db.Column(db.Integer, nullable=True)
    data_hora = db.Column(db.String(30), nullable=False)
    hora_ocorrencia = db.Column(db.String(10), nullable=False)
    natureza = db.Column(db.String(120), nullable=False)
    descricao = db.Column(db.Text, nullable=False)
    site = db.Column(db.String(128), nullable=True)
    local = db.Column(db.String(120), nullable=False)
    operador = db.Column(db.String(120), nullable=False)
    gc = db.Column(db.Text, nullable=False)
    envolvido = db.Column(db.String(120), nullable=True)
    foto = db.Column(db.Text, nullable=True)

    prioridade = db.Column(db.String(20), nullable=False, default="MEDIA")
    status = db.Column(db.String(30), nullable=False, default="PENDENTE")

    situacao_investigacao = db.Column(db.String(30), nullable=True)
    conclusao_investigacao = db.Column(db.Text, nullable=True)
    anexo_post = db.Column(db.Text, nullable=True)
    anexo_post_nome = db.Column(db.String(255), nullable=True)

    criado_por = db.Column(db.String(120), nullable=True)
    criado_em = db.Column(db.DateTime, default=datetime.utcnow)
    atualizado_em = db.Column(db.DateTime, default=datetime.utcnow, onupdate=datetime.utcnow)
    boletim_ocorrencia = db.Column(db.Boolean, default=False, nullable=True)
    custo = db.Column(db.String(50), nullable=True)
    valor_recuperado = db.Column(db.String(80), nullable=True)
    valor_preventivo = db.Column(db.String(80), nullable=True)
    prejuizo         = db.Column(db.String(80), nullable=True)
    responsavel_fechamento = db.Column(db.String(120), nullable=True)
    anexo_post_2      = db.Column(db.Text, nullable=True)
    anexo_post_nome_2 = db.Column(db.String(255), nullable=True)
    anexo_post_3      = db.Column(db.Text, nullable=True)
    anexo_post_nome_3 = db.Column(db.String(255), nullable=True)

    # Vínculos com outros documentos
    anc_vinculada_id     = db.Column(db.Integer, nullable=True)
    analise_vinculada_id = db.Column(db.Integer, nullable=True)


# =========================
# HELPERS
# =========================
MIME_MAP = {
    "jpg": "image/jpeg", "jpeg": "image/jpeg",
    "png": "image/png", "webp": "image/webp",
    "gif": "image/gif", "pdf": "application/pdf",
    "doc": "application/msword",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xls": "application/vnd.ms-excel",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
}

def gerar_codigo(model_class, site):
    clean = re.sub(r'[^A-Z0-9]', '', (site or "SITE").upper())[:8] or "SITE"
    ano = datetime.now().year
    max_seq = db.session.query(func.max(model_class.numero_site)).filter(
        model_class.site == site
    ).scalar() or 0
    seq = max_seq + 1
    return f"{clean}-{ano}-{seq:04d}", seq


def gerar_numero_anc(site):
    clean = re.sub(r'[^A-Z0-9]', '', (site or "SITE").upper())[:8] or "SITE"
    ano = datetime.now().year
    max_seq = db.session.query(func.max(ANC.numero_site)).filter(
        ANC.site == site
    ).scalar() or 0
    seq = max_seq + 1
    return f"ANC-{clean}-{ano}-{seq:04d}", seq


def _b64_decode(s: str) -> bytes:
    """Decodifica base64 de forma robusta.
    Remove prefixo data URI (data:mime;base64,) e whitespace que o Oracle
    pode inserir em CLOBs longos antes de decodificar."""
    s = (s or "").strip()
    if "," in s:
        s = s.split(",", 1)[1]
    s = "".join(s.split())   # remove \n, \r, espaços
    return base64.b64decode(s)


def fmt_data_br(data_str):
    """Converte YYYY-MM-DD para DD/MM/YYYY."""
    if data_str and len(data_str) >= 10:
        try:
            return datetime.strptime(data_str[:10], "%Y-%m-%d").strftime("%d/%m/%Y")
        except ValueError:
            pass
    return data_str or "—"


def aplicar_filtros_anc(query, mostrar_excluidas=False):
    data_inicial = (request.args.get("data_inicial") or "").strip()
    data_final   = (request.args.get("data_final")   or "").strip()
    status       = (request.args.get("status")        or "").strip().upper()
    turno        = (request.args.get("turno")         or "").strip().upper()
    setor        = (request.args.get("setor")         or "").strip()
    natureza     = (request.args.get("natureza")      or "").strip()
    site_filtro  = (request.args.get("site_filtro")   or "").strip()

    # Oculta ANCs com soft-delete (excluido='S') — a menos que explicitamente pedido
    if not mostrar_excluidas:
        query = query.filter(ANC.excluido != 'S')

    # Empurra todos os filtros para o SQL — sem Python-level filtering
    if data_inicial:
        query = query.filter(ANC.data_nc >= data_inicial)
    if data_final:
        query = query.filter(ANC.data_nc <= data_final)
    if status:
        query = query.filter(func.upper(ANC.status) == status)
    if turno:
        query = query.filter(func.upper(ANC.turno) == turno)
    if setor:
        query = query.filter(func.lower(ANC.setor).contains(setor.lower()))
    if natureza:
        query = query.filter(func.lower(ANC.natureza) == natureza.lower())
    if site_filtro:
        query = query.filter(ANC.site == site_filtro)

    registros = query.all()

    filtros = {
        "data_inicial": data_inicial, "data_final": data_final,
        "status": status, "turno": turno, "setor": setor,
        "natureza": natureza, "site_filtro": site_filtro,
    }
    return registros, filtros


_LGPD_TEXT = (
    "Este documento e seus anexos podem conter dados pessoais protegidos pela Lei Geral de Proteção de Dados "
    "(LGPD - Lei nº 13.709/2018). A base legal para o tratamento de dados pessoais aqui realizada é a execução "
    "das políticas internas, conforme previsto no artigo 7º da LGPD. As informações são para uso exclusivo do(s) "
    "destinatário(s) original(is), e qualquer uso não autorizado pode violar a LGPD. Se você recebeu este "
    "documento por engano, por favor, informe o remetente e destrua-o imediatamente."
)

def _adicionar_lgpd_excel(ws, num_cols):
    """Adiciona aviso LGPD como última linha mesclada na planilha."""
    from openpyxl.styles import Font as _Font, Alignment as _Align, PatternFill as _PFill
    from openpyxl.utils import get_column_letter
    ws.append([""] * num_cols)                      # linha em branco
    ws.append([_LGPD_TEXT] + [""] * (num_cols - 1))
    lgpd_row = ws.max_row
    last_col  = get_column_letter(num_cols)
    ws.merge_cells(f"A{lgpd_row}:{last_col}{lgpd_row}")
    cell = ws.cell(row=lgpd_row, column=1)
    cell.font      = _Font(size=7, color="888888", italic=True)
    cell.alignment = _Align(wrap_text=True, horizontal="left", vertical="top")
    cell.fill      = _PFill("solid", fgColor="F9FAFB")
    ws.row_dimensions[lgpd_row].height = 72


def _desenhar_lgpd(canvas, x_ini, y_ini, font_size=5.5, leading=7.5):
    """Desenha o aviso LGPD em múltiplas linhas a partir de (x_ini, y_ini) para baixo."""
    from reportlab.lib import colors as _colors
    canvas.setFont("Helvetica", font_size)
    canvas.setFillColor(_colors.HexColor("#9ca3af"))
    linhas = textwrap.wrap(_LGPD_TEXT, width=170)
    y = y_ini
    for linha in linhas:
        canvas.drawString(x_ini, y, linha)
        y -= leading


def gerar_pdf_anc_bytes(anc):
    """Gera PDF do ANC no formato oficial DHL Security."""
    _import_reportlab()
    buffer  = BytesIO()
    BLACK   = colors.black
    YELLOW  = colors.HexColor("#FFCC00")

    pw = A4[0] - 2.6 * rcm
    doc_pdf = SimpleDocTemplate(
        buffer, pagesize=A4,
        leftMargin=1.3*rcm, rightMargin=1.3*rcm,
        topMargin=0.8*rcm, bottomMargin=2.0*rcm,
    )

    s_normal = ParagraphStyle("an", fontName="Helvetica",      fontSize=8,  textColor=BLACK, leading=10)
    s_th     = ParagraphStyle("ath",fontName="Helvetica-Bold", fontSize=8,  textColor=BLACK, alignment=TA_CENTER, leading=10)
    s_td     = ParagraphStyle("atd",fontName="Helvetica",      fontSize=8,  textColor=BLACK, alignment=TA_CENTER, leading=10)
    s_h3     = ParagraphStyle("ah3",fontName="Helvetica-Bold", fontSize=9,  textColor=BLACK)
    s_title  = ParagraphStyle("ati",fontName="Helvetica-Bold", fontSize=11, textColor=BLACK, alignment=TA_CENTER)
    s_foot   = ParagraphStyle("afo",fontName="Helvetica",      fontSize=8,  textColor=BLACK)

    def _fit_img(source, max_w, max_h):
        """Retorna RLImage escalada proporcionalmente para caber em max_w × max_h."""
        if isinstance(source, str):
            bio = BytesIO(_b64_decode(source))
        else:
            bio = source
        bio.seek(0)
        iw, ih = ImageReader(bio).getSize()
        scale = min(max_w / iw, max_h / ih, 1.0)
        bio.seek(0)
        return RLImage(bio, width=iw * scale, height=ih * scale)

    def yellow_bar(text):
        t = Table([[Paragraph(text, s_h3)]], colWidths=[pw])
        t.setStyle(TableStyle([
            ("BACKGROUND",    (0,0),(-1,-1), YELLOW),
            ("BOX",           (0,0),(-1,-1), 0.5, BLACK),
            ("TOPPADDING",    (0,0),(-1,-1), 3),
            ("BOTTOMPADDING", (0,0),(-1,-1), 3),
            ("LEFTPADDING",   (0,0),(-1,-1), 6),
        ]))
        return t

    story = []

    # ── 1. CABEÇALHO ──────────────────────────────────────────────
    logo_path = os.path.join(app.root_path, "static", "logo.png")
    if os.path.exists(logo_path):
        with open(logo_path, "rb") as _lf:
            logo_cell = _fit_img(BytesIO(_lf.read()), max_w=pw * 0.22, max_h=1.2*rcm)
    else:
        logo_cell = Paragraph("<b>DHL</b>", s_normal)

    hdr = Table(
        [[logo_cell, Paragraph("AVISO DE NÃO CONFORMIDADE   ", s_title)]],
        colWidths=[pw * 0.25, pw * 0.75],
    )
    hdr.setStyle(TableStyle([
        ("BACKGROUND",    (0,0),(-1,-1), YELLOW),
        ("BOX",           (0,0),(-1,-1), 0.5, BLACK),
        ("INNERGRID",     (0,0),(-1,-1), 0.5, BLACK),
        ("VALIGN",        (0,0),(-1,-1), "MIDDLE"),
        ("TOPPADDING",    (0,0),(-1,-1), 4),
        ("BOTTOMPADDING", (0,0),(-1,-1), 4),
        ("LEFTPADDING",   (0,0),(-1,-1), 6),
        ("RIGHTPADDING",  (0,0),(-1,-1), 6),
    ]))
    story += [hdr, Spacer(1, 0.2*rcm)]

    # ── 2. TABELA DE IDENTIFICAÇÃO ────────────────────────────────
    id_heads = ["DATA", "HORA DA\nOCORRÊNCIA", "NATUREZA", "LOCAL",
                "PESSOAS\nENVOLVIDAS", "Nº ANC"]
    id_vals  = [fmt_data_br(anc.data_nc), anc.hora_nc or "—", anc.natureza or "—",
                anc.local or "—", anc.envolvido or "—",
                str(anc.numero_site or anc.id)]
    cw = pw / 6
    id_tbl = Table(
        [[Paragraph(h, s_th) for h in id_heads],
         [Paragraph(v, s_td) for v in id_vals]],
        colWidths=[cw] * 6,
    )
    id_tbl.setStyle(TableStyle([
        ("BACKGROUND",    (0,0),(-1,0), YELLOW),
        ("BOX",           (0,0),(-1,-1), 0.5, BLACK),
        ("INNERGRID",     (0,0),(-1,-1), 0.5, BLACK),
        ("VALIGN",        (0,0),(-1,-1), "MIDDLE"),
        ("TOPPADDING",    (0,0),(-1,-1), 3),
        ("BOTTOMPADDING", (0,0),(-1,-1), 3),
    ]))
    story += [id_tbl, Spacer(1, 0.2*rcm)]

    # ── 3. DESCRIÇÃO ──────────────────────────────────────────────
    story.append(yellow_bar("Descrição da ocorrência:"))
    desc_tbl = Table([[Paragraph(anc.descricao or "—", s_normal)]], colWidths=[pw])
    desc_tbl.setStyle(TableStyle([
        ("BOX",           (0,0),(-1,-1), 0.5, BLACK),
        ("TOPPADDING",    (0,0),(-1,-1), 5),
        ("BOTTOMPADDING", (0,0),(-1,-1), 8),
        ("LEFTPADDING",   (0,0),(-1,-1), 8),
        ("RIGHTPADDING",  (0,0),(-1,-1), 8),
    ]))
    story += [desc_tbl, Spacer(1, 0.2*rcm)]

    # ── 4. GRAVIDADE / RESPONSÁVEL / PLANO DE AÇÃO ───────────────
    grav_tbl = Table(
        [[Paragraph(h, s_th) for h in ["GRAVIDADE", "RESPONSÁVEL", "PLANO DE AÇÃO", "VALOR"]],
         [Paragraph(v, s_td) for v in [anc.gravidade or "—",
                                        anc.responsavel or "—",
                                        anc.status or "—",
                                        anc.valor or "—"]]],
        colWidths=[pw/4, pw/4, pw/4, pw/4],
    )
    grav_tbl.setStyle(TableStyle([
        ("BACKGROUND",    (0,0),(-1,0), YELLOW),
        ("BOX",           (0,0),(-1,-1), 0.5, BLACK),
        ("INNERGRID",     (0,0),(-1,-1), 0.5, BLACK),
        ("VALIGN",        (0,0),(-1,-1), "MIDDLE"),
        ("TOPPADDING",    (0,0),(-1,-1), 3),
        ("BOTTOMPADDING", (0,0),(-1,-1), 3),
    ]))
    story += [grav_tbl, Spacer(1, 0.2*rcm)]

    # ── 5. REGISTROS FOTOGRÁFICOS (até 6 imagens, 2 por linha) ───
    story.append(yellow_bar("Registros Fotográficos:"))

    imgs_b64 = [x for x in [
        anc.imagem_1, anc.imagem_2, anc.imagem_3,
        anc.imagem_4, anc.imagem_5, anc.imagem_6,
    ] if x]

    img_col_w = pw / 3
    img_h     = 3.8 * rcm

    foto_rows = []
    for i in range(0, max(len(imgs_b64), 3), 3):
        row = []
        for j in range(3):
            idx = i + j
            cell = ""
            if idx < len(imgs_b64):
                try:
                    cell = _fit_img(imgs_b64[idx], img_col_w - 0.6*rcm, img_h)
                except Exception:
                    pass
            row.append(cell)
        foto_rows.append(row)

    foto_tbl = Table(foto_rows, colWidths=[img_col_w, img_col_w, img_col_w])
    foto_tbl.setStyle(TableStyle([
        ("BOX",           (0,0),(-1,-1), 0.5, BLACK),
        ("INNERGRID",     (0,0),(-1,-1), 0.5, BLACK),
        ("ALIGN",         (0,0),(-1,-1), "CENTER"),
        ("VALIGN",        (0,0),(-1,-1), "MIDDLE"),
        ("TOPPADDING",    (0,0),(-1,-1), 4),
        ("BOTTOMPADDING", (0,0),(-1,-1), 4),
        ("MINROWHEIGHT",  (0,0),(-1,-1), img_h),
    ]))
    story += [foto_tbl, Spacer(1, 0.2*rcm)]

    # ── 6. RESPONSÁVEL PELO LEVANTAMENTO ──────────────────────────
    resp_tbl = Table(
        [[Paragraph("RESPONSÁVEL PELO LEVANTAMENTO", s_th), Paragraph("CARGO", s_th)],
         [Paragraph(anc.gestor_responsavel or anc.responsavel or "—", s_td),
          Paragraph(anc.cargo or "Segurança Patrimonial", s_td)]],
        colWidths=[pw * 0.60, pw * 0.40],
    )
    resp_tbl.setStyle(TableStyle([
        ("BACKGROUND",    (0,0),(-1,0), YELLOW),
        ("BOX",           (0,0),(-1,-1), 0.5, BLACK),
        ("INNERGRID",     (0,0),(-1,-1), 0.5, BLACK),
        ("VALIGN",        (0,0),(-1,-1), "MIDDLE"),
        ("TOPPADDING",    (0,0),(-1,-1), 3),
        ("BOTTOMPADDING", (0,0),(-1,-1), 5),
        ("LEFTPADDING",   (0,0),(-1,-1), 6),
        ("RIGHTPADDING",  (0,0),(-1,-1), 6),
    ]))
    story += [resp_tbl, Spacer(1, 0.2*rcm)]

    # ── 7. PLANO DE AÇÃO / FECHAMENTO (se houver) ─────────────────
    if anc.plano_acao_texto:
        story.append(Spacer(1, 0.5*rcm))
        story.append(yellow_bar("Plano de Ação / Fechamento:"))
        pla_tbl = Table([[Paragraph(anc.plano_acao_texto or "—", s_normal)]], colWidths=[pw])
        pla_tbl.setStyle(TableStyle([
            ("BOX",           (0,0),(-1,-1), 0.5, BLACK),
            ("TOPPADDING",    (0,0),(-1,-1), 10),
            ("BOTTOMPADDING", (0,0),(-1,-1), 20),
            ("LEFTPADDING",   (0,0),(-1,-1), 10),
            ("RIGHTPADDING",  (0,0),(-1,-1), 10),
        ]))
        story.append(pla_tbl)
        if anc.fechado_por or anc.fechado_em:
            fechado_em_fmt = (
                anc.fechado_em.strftime("%d/%m/%Y %H:%M") if anc.fechado_em else "—"
            )
            story.append(Spacer(1, 0.2*rcm))
            story.append(Paragraph(
                f"<b>FECHADO POR:</b> {anc.fechado_por or '—'}   "
                f"<b>EM:</b> {fechado_em_fmt}",
                s_foot,
            ))

    def _footer_anc(canvas, doc):
        canvas.saveState()
        x0, x1 = 1.5*rcm, A4[0] - 1.5*rcm
        canvas.setStrokeColor(BLACK)
        canvas.setLineWidth(0.5)
        canvas.line(x0, 1.9*rcm, x1, 1.9*rcm)
        canvas.setFont("Helvetica", 7)
        canvas.setFillColor(colors.HexColor("#6b7280"))
        canvas.drawString(x0, 1.5*rcm, "DHL Security — Aviso de Não Conformidade")
        canvas.drawRightString(x1, 1.5*rcm, f"Página {doc.page}")
        _desenhar_lgpd(canvas, x0, 1.05*rcm)
        canvas.restoreState()

    doc_pdf.build(story, onFirstPage=_footer_anc, onLaterPages=_footer_anc)
    buffer.seek(0)
    return buffer


def gerar_docx_de_registro(registro):
    """Gera DOCX a partir dos campos salvos no banco (sem imagens)."""
    _import_docx()
    doc = Document()
    for sec in doc.sections:
        sec.left_margin  = Cm(1.7)
        sec.right_margin = Cm(1.7)
        sec.top_margin   = Cm(1.5)
        sec.bottom_margin = Cm(1.5)

    logo_path = os.path.join(app.root_path, "static", "logo.png")
    for sec in doc.sections:
        header = sec.header
        for p in header.paragraphs:
            p.text = ""
        avail = sec.page_width - sec.left_margin - sec.right_margin
        ht = header.add_table(rows=1, cols=2, width=avail)
        ht.alignment = WD_TABLE_ALIGNMENT.CENTER
        ht.autofit = False
        ht.columns[0].width = Cm(4.5)
        ht.columns[1].width = Cm(13.1)
        c_logo, c_info = ht.rows[0].cells[0], ht.rows[0].cells[1]
        for c in [c_logo, c_info]:
            set_cell_border(c, color="FFFFFF")
            c.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        p_logo = c_logo.paragraphs[0]
        p_logo.alignment = WD_ALIGN_PARAGRAPH.LEFT
        if os.path.exists(logo_path):
            p_logo.add_run().add_picture(logo_path, width=Cm(3.8))
        else:
            r = p_logo.add_run("DHL")
            r.bold = True; r.font.name = "Arial"
            r.font.size = Pt(16); r.font.color.rgb = RGBColor(212, 5, 17)
        p_info = c_info.paragraphs[0]
        p_info.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        r1 = p_info.add_run("DHL SECURITY\n")
        r1.bold = True; r1.font.name = "Arial"
        r1.font.size = Pt(10); r1.font.color.rgb = RGBColor(212, 5, 17)
        r2 = p_info.add_run("Relatório de Análise Investigativa")
        r2.font.name = "Arial"; r2.font.size = Pt(8)
        r2.font.color.rgb = RGBColor(90, 90, 90)

        # ── Rodapé LGPD ──
        footer = sec.footer
        for fp in footer.paragraphs:
            fp.text = ""
        fp_lgpd = footer.paragraphs[0]
        fp_lgpd.alignment = WD_ALIGN_PARAGRAPH.CENTER
        fp_lgpd.paragraph_format.space_before = Pt(2)
        pf = fp_lgpd.paragraph_format
        pf.left_indent  = Cm(0)
        pf.right_indent = Cm(0)
        # linha separadora via borda superior
        from docx.oxml.ns import qn as _qn
        from docx.oxml import OxmlElement as _OE
        _pPr = fp_lgpd._p.get_or_add_pPr()
        _pBdr = _OE("w:pBdr")
        _top  = _OE("w:top")
        _top.set(_qn("w:val"), "single")
        _top.set(_qn("w:sz"), "6")
        _top.set(_qn("w:space"), "1")
        _top.set(_qn("w:color"), "D40511")
        _pBdr.append(_top)
        _pPr.append(_pBdr)
        rLgpd = fp_lgpd.add_run(_LGPD_TEXT)
        rLgpd.font.name = "Arial"
        rLgpd.font.size = Pt(6)
        rLgpd.font.color.rgb = RGBColor(120, 120, 120)

    doc.add_paragraph()

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(8)
    p.paragraph_format.space_after  = Pt(4)
    run = p.add_run("ANÁLISE INVESTIGATIVA")
    run.bold = True; run.font.name = "Arial"
    run.font.size = Pt(17); run.font.color.rgb = RGBColor(31, 41, 55)

    faixa = doc.add_table(rows=1, cols=1)
    faixa.alignment = WD_TABLE_ALIGNMENT.CENTER
    cell = faixa.rows[0].cells[0]
    cell.text = "DHL SECURITY"
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
    set_cell_bg(cell, "FFCC00")
    set_cell_border(cell, color="D40511", size="12")
    p = cell.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(6)
    p.paragraph_format.space_after  = Pt(6)
    for run in p.runs:
        run.bold = True; run.font.name = "Arial"
        run.font.size = Pt(10); run.font.color.rgb = RGBColor(120, 0, 0)
    doc.add_paragraph()

    add_section_title(doc, "Dados da Operação")
    add_professional_table(doc, [
        ["Nº do Relatório (ID):", str(registro.numero_site or registro.id)],
        ["Empresa:",              registro.empresa              or ""],
        ["Unidade:",              registro.unidade              or ""],
        ["Endereço:",             registro.endereco             or ""],
        ["Classificação do Site:",registro.classificacao        or ""],
        ["Produtos Segmento:",    registro.produtos_segmento    or ""],
        ["Cliente(s):",           registro.clientes             or ""],
    ], col_widths=[6.0, 13.5])

    add_section_title(doc, "Dados do Levantamento")
    add_professional_table(doc, [
        ["Objetivo:",                      registro.objetivo             or ""],
        ["Responsável pelo Levantamento:", registro.responsavel          or ""],
        ["Função:",                        registro.funcao_levantamento  or ""],
        ["Data:",                          registro.data_levantamento    or ""],
        ["Valor estimado:",                registro.valor                or "—"],
    ], col_widths=[6.0, 13.5])

    add_section_title(doc, "Descrição do Registro")
    add_text_box(doc, registro.descricao_registro or "")

    if registro.imagens:
        add_section_title(doc, "Evidências")
        for idx, img_obj in enumerate(registro.imagens, start=1):
            try:
                _raw = img_obj.b64
                if not _raw:
                    continue
                bio = BytesIO(_b64_decode(_raw))
                bio.seek(0)
                tbl = doc.add_table(rows=1, cols=2)
                tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
                tbl.autofit = False
                img_cell  = tbl.rows[0].cells[0]
                desc_cell = tbl.rows[0].cells[1]
                img_cell.width  = Cm(9.0)
                desc_cell.width = Cm(8.6)
                p_img = img_cell.paragraphs[0]
                p_img.alignment = WD_ALIGN_PARAGRAPH.CENTER
                p_img.add_run().add_picture(bio, width=Cm(8.5))
                p_num = img_cell.add_paragraph(f"Imagem {idx}")
                p_num.alignment = WD_ALIGN_PARAGRAPH.CENTER
                for rn in p_num.runs:
                    rn.font.name = "Arial"; rn.font.size = Pt(8)
                    rn.italic = True; rn.font.color.rgb = RGBColor(100, 100, 100)
                format_cell(img_cell, bg="F5F5F5", align="center")
                desc_cell.text = img_obj.descricao or "Sem descrição"
                format_cell(desc_cell, bg="FFFFFF", align="left")
                p = doc.add_paragraph()
                p.paragraph_format.space_after = Pt(6)
                if idx % 5 == 0 and idx < len(registro.imagens):
                    doc.add_page_break()
            except Exception:
                pass

    add_section_title(doc, "Conclusão")
    add_text_box(doc, registro.conclusao or "")

    add_section_title(doc, "Recomendação")
    add_text_box(doc, registro.sugestao or "")

    doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    rod = p.add_run(
        f"Código: {registro.codigo or registro.id}"
        + (f" | Criado em: {registro.criado_em.strftime('%d/%m/%Y %H:%M')}"
           if registro.criado_em else "")
    )
    rod.font.name = "Arial"; rod.font.size = Pt(8)
    rod.font.color.rgb = RGBColor(100, 100, 100)

    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def gerar_pdf_analise_bytes(form_data, evidencias_bytes):
    """PDF corporativo DHL da análise investigativa."""
    _import_reportlab()
    buffer = BytesIO()

    DHL_RED    = colors.HexColor("#D40511")
    DHL_YELLOW = colors.HexColor("#FFCC00")
    DHL_DARK   = colors.HexColor("#1F2937")
    DHL_MUTED  = colors.HexColor("#6B7280")
    LABEL_BG   = colors.HexColor("#FFF2CC")

    pw = A4[0] - 3.4 * rcm

    doc_pdf = SimpleDocTemplate(
        buffer, pagesize=A4,
        leftMargin=1.7*rcm, rightMargin=1.7*rcm,
        topMargin=2.5*rcm, bottomMargin=2.8*rcm,
    )

    s_title   = ParagraphStyle("s_title",   fontName="Helvetica-Bold", fontSize=17,
                                textColor=DHL_DARK, alignment=TA_CENTER, spaceAfter=4)
    s_section = ParagraphStyle("s_section", fontName="Helvetica-Bold", fontSize=11,
                                textColor=DHL_RED, spaceBefore=14, spaceAfter=6)
    s_body    = ParagraphStyle("s_body",    fontName="Helvetica", fontSize=9,
                                textColor=DHL_DARK, leading=13, spaceAfter=3)
    s_label   = ParagraphStyle("s_label",   fontName="Helvetica-Bold", fontSize=9,
                                textColor=DHL_DARK)
    s_hdr_r   = ParagraphStyle("s_hdr_r",   fontName="Helvetica", fontSize=8,
                                alignment=TA_RIGHT, leading=13)

    def info_table(rows):
        data = [[Paragraph(f"<b>{k}</b>", s_label), Paragraph(str(v or "—"), s_body)]
                for k, v in rows]
        t = Table(data, colWidths=[5.5*rcm, pw - 5.5*rcm])
        t.setStyle(TableStyle([
            ("BACKGROUND",     (0,0), (0,-1),  LABEL_BG),
            ("BACKGROUND",     (1,0), (1,-1),  colors.white),
            ("GRID",           (0,0), (-1,-1), 0.5, colors.HexColor("#D9D9D9")),
            ("VALIGN",         (0,0), (-1,-1), "TOP"),
            ("TOPPADDING",     (0,0), (-1,-1), 5),
            ("BOTTOMPADDING",  (0,0), (-1,-1), 5),
            ("LEFTPADDING",    (0,0), (-1,-1), 8),
        ]))
        return t

    def text_box(text):
        t = Table([[Paragraph(str(text or "—"), s_body)]], colWidths=[pw])
        t.setStyle(TableStyle([
            ("BOX",            (0,0), (-1,-1), 0.5, colors.HexColor("#D9D9D9")),
            ("BACKGROUND",     (0,0), (-1,-1), colors.white),
            ("TOPPADDING",     (0,0), (-1,-1), 8),
            ("BOTTOMPADDING",  (0,0), (-1,-1), 8),
            ("LEFTPADDING",    (0,0), (-1,-1), 10),
            ("RIGHTPADDING",   (0,0), (-1,-1), 10),
        ]))
        return t

    story = []

    # ── cabeçalho ──
    logo_path = os.path.join(app.root_path, "static", "logo.png")
    if os.path.exists(logo_path):
        _lbio = BytesIO(open(logo_path, "rb").read())
        _liw, _lih = ImageReader(_lbio).getSize()
        _lscale = min(3.8*rcm / _liw, 1.4*rcm / _lih)
        logo_cell = RLImage(logo_path, width=_liw*_lscale, height=_lih*_lscale)
    else:
        logo_cell = Paragraph('<b><font color="#D40511" size="14">DHL</font></b>',
                              ParagraphStyle("tmp", fontName="Helvetica"))

    hdr = Table(
        [[logo_cell, Paragraph(
            '<font color="#D40511"><b>DHL SECURITY</b></font><br/>'
            '<font color="#6B7280">Relatório de Análise Investigativa</font>',
            s_hdr_r)]],
        colWidths=[pw * 0.35, pw * 0.65]
    )
    hdr.setStyle(TableStyle([
        ("VALIGN",       (0,0), (-1,-1), "MIDDLE"),
        ("LINEBELOW",    (0,0), (-1,-1), 1.5, DHL_RED),
        ("BOTTOMPADDING",(0,0), (-1,-1), 8),
    ]))
    story += [hdr, Spacer(1, 0.4*rcm)]

    # ── título ──
    story.append(Paragraph("ANÁLISE INVESTIGATIVA", s_title))
    story.append(Spacer(1, 0.35*rcm))

    # ── faixa amarela ──
    banner = Table([["DHL SECURITY"]], colWidths=[pw])
    banner.setStyle(TableStyle([
        ("BACKGROUND",     (0,0), (-1,-1), DHL_YELLOW),
        ("TEXTCOLOR",      (0,0), (-1,-1), colors.HexColor("#7A0000")),
        ("FONTNAME",       (0,0), (-1,-1), "Helvetica-Bold"),
        ("FONTSIZE",       (0,0), (-1,-1), 10),
        ("ALIGN",          (0,0), (-1,-1), "CENTER"),
        ("TOPPADDING",     (0,0), (-1,-1), 7),
        ("BOTTOMPADDING",  (0,0), (-1,-1), 7),
        ("BOX",            (0,0), (-1,-1), 1.5, DHL_RED),
    ]))
    story += [banner, Spacer(1, 0.5*rcm)]

    # ── dados da operação ──
    story.append(Paragraph("DADOS DA OPERAÇÃO", s_section))
    story.append(info_table([
        ("Nº do Relatório (ID):",  form_data.get("id_relatorio", "")),
        ("Empresa:",               form_data.get("empresa", "")),
        ("Unidade:",               form_data.get("unidade", "")),
        ("Endereço:",              form_data.get("endereco", "")),
        ("Classificação do Site:", form_data.get("classificacao", "")),
        ("Produtos / Segmento:",   form_data.get("produtos_segmento", "")),
        ("Cliente(s):",            form_data.get("clientes", "")),
    ]))
    story.append(Spacer(1, 0.3*rcm))

    # ── dados do levantamento ──
    story.append(Paragraph("DADOS DO LEVANTAMENTO", s_section))
    story.append(info_table([
        ("Objetivo:",    form_data.get("objetivo", "")),
        ("Responsável:", form_data.get("responsavel", "")),
        ("Função:",      form_data.get("funcao_levantamento", "")),
        ("Data:",        form_data.get("data_levantamento", "")),
    ]))
    story.append(Spacer(1, 0.3*rcm))

    # ── descrição ──
    story.append(Paragraph("DESCRIÇÃO DO REGISTRO", s_section))
    story.append(text_box(form_data.get("descricao_registro", "")))
    story.append(Spacer(1, 0.3*rcm))

    # ── evidências: imagem esquerda | descrição direita | 5 por página ──
    if evidencias_bytes:
        story.append(Paragraph("EVIDÊNCIAS", s_section))
        for idx, (img_bytes, desc) in enumerate(evidencias_bytes, start=1):
            if idx > 1 and (idx - 1) % 5 == 0:
                story.append(PageBreak())
                story.append(Paragraph("EVIDÊNCIAS (continuação)", s_section))
            try:
                rl_img = RLImage(BytesIO(img_bytes), width=7.5*rcm, height=5.5*rcm)
                num_p  = Paragraph(f"<i>Img {idx}</i>",
                                   ParagraphStyle("cap", fontName="Helvetica-Oblique",
                                                  fontSize=7, textColor=DHL_MUTED,
                                                  alignment=TA_CENTER))
                ev = Table(
                    [[rl_img, Paragraph(desc or "Sem descrição", s_body)],
                     [num_p,  ""]],
                    colWidths=[8.0*rcm, pw - 8.0*rcm],
                    rowHeights=[None, 0.5*rcm]
                )
                ev.setStyle(TableStyle([
                    ("BOX",           (0,0), (-1,-1), 0.5, colors.HexColor("#D9D9D9")),
                    ("INNERGRID",     (0,0), (-1,-1), 0.5, colors.HexColor("#D9D9D9")),
                    ("BACKGROUND",    (0,0), (0,-1),  colors.HexColor("#F5F5F5")),
                    ("VALIGN",        (0,0), (-1,-1), "MIDDLE"),
                    ("ALIGN",         (0,0), (0,-1),  "CENTER"),
                    ("SPAN",          (1,0), (1,1)),
                    ("TOPPADDING",    (0,0), (-1,-1), 6),
                    ("BOTTOMPADDING", (0,0), (-1,-1), 6),
                    ("LEFTPADDING",   (0,0), (-1,-1), 6),
                    ("RIGHTPADDING",  (0,0), (-1,-1), 8),
                ]))
                story += [ev, Spacer(1, 0.2*rcm)]
            except Exception:
                pass
        story.append(Spacer(1, 0.2*rcm))

    # ── conclusão ──
    story.append(Paragraph("CONCLUSÃO", s_section))
    story.append(text_box(form_data.get("conclusao", "")))
    story.append(Spacer(1, 0.3*rcm))

    # ── recomendação ──
    story.append(Paragraph("RECOMENDAÇÃO", s_section))
    story.append(text_box(form_data.get("sugestao", "")))

    # ── rodapé ──
    id_rel = form_data.get("id_relatorio", "")

    def footer(canvas, doc):
        canvas.saveState()
        x0, x1 = 1.7*rcm, A4[0] - 1.7*rcm
        canvas.setStrokeColor(DHL_RED)
        canvas.setLineWidth(0.8)
        canvas.line(x0, 1.9*rcm, x1, 1.9*rcm)
        canvas.setFont("Helvetica", 7)
        canvas.setFillColor(DHL_MUTED)
        canvas.drawString(x0, 1.5*rcm,
                          f"DHL Security — Análise Investigativa{' | ID: ' + id_rel if id_rel else ''}")
        canvas.drawRightString(x1, 1.5*rcm, f"Página {doc.page}")
        _desenhar_lgpd(canvas, x0, 1.05*rcm)
        canvas.restoreState()

    doc_pdf.build(story, onFirstPage=footer, onLaterPages=footer)
    buffer.seek(0)
    return buffer


def arquivo_para_base64(arquivo, extensoes):
    if not arquivo or not arquivo.filename:
        return None, None
    ext = arquivo.filename.rsplit(".", 1)[-1].lower() if "." in arquivo.filename else ""
    if ext not in extensoes:
        return None, None
    mime = MIME_MAP.get(ext, "application/octet-stream")
    dados = base64.b64encode(arquivo.read()).decode("utf-8")
    return f"data:{mime};base64,{dados}", arquivo.filename


@app.route('/data-upload/<path:filename>')
def serve_upload(filename):
    """Serve arquivos enviados pelo usuário a partir do diretório persistente.
    foto_path no banco é 'uploads/achados_perdidos/fname', então filename aqui
    já inclui 'uploads/', basta servir diretamente de UPLOADS_ROOT.
    Ex: UPLOADS_ROOT/uploads/achados_perdidos/fname"""
    from flask import send_from_directory as _sfd
    return _sfd(UPLOADS_ROOT, filename)


def login_required(func):
    @wraps(func)
    def wrapper(*args, **kwargs):
        if not session.get("user_id"):
            flash("Faça login para continuar.", "warning")
            return redirect(url_for("login"))
        # Bloqueia acesso enquanto LGPD não for aceita
        if session.get("user_lgpd_aceito") != "sim":
            return redirect(url_for("lgpd_aceite"))
        return func(*args, **kwargs)
    return wrapper


def perfil_required(*perfis):
    def decorator(func):
        @wraps(func)
        def wrapper(*args, **kwargs):
            if not session.get("user_id"):
                flash("Faça login para continuar.", "warning")
                return redirect(url_for("login"))

            perfil = (session.get("user_perfil") or "").upper()
            if perfil not in perfis:
                flash("Você não tem permissão para acessar esta área.", "danger")
                return redirect(url_for("ocorrencias"))
            return func(*args, **kwargs)
        return wrapper
    return decorator


# =========================
# CHECKLIST DE CÂMERAS — ROTAS
# =========================

@app.route("/cameras")
@login_required
def cameras_index():
    site     = session.get("user_site") or ""
    is_admin = _is_privileged()

    if is_admin:
        cameras    = Camera.query.order_by(Camera.site, Camera.ativo.desc(), Camera.numero).all()
        checklists = (ChecklistCamera.query
                      .order_by(ChecklistCamera.criado_em.desc())
                      .limit(50).all())
    else:
        cameras    = _query_filtrar_sites(Camera.query, Camera).order_by(Camera.ativo.desc(), Camera.numero).all()
        checklists = (_query_filtrar_sites(ChecklistCamera.query, ChecklistCamera)
                      .order_by(ChecklistCamera.criado_em.desc())
                      .limit(50).all())

    ativas   = sum(1 for c in cameras if c.ativo)
    inativas = sum(1 for c in cameras if not c.ativo)

    # Sites distintos para os dropdowns de filtro (só relevante para admin)
    sites_cam = sorted({c.site for c in cameras if c.site}) if is_admin else []
    sites_chk = sorted({c.site for c in checklists if c.site}) if is_admin else []

    return render_template("cameras.html",
        cameras=cameras, checklists=checklists,
        ativas=ativas, inativas=inativas, site=site,
        is_admin=is_admin, can_manage=_is_can_manage(),
        sites_cam=sites_cam, sites_chk=sites_chk)


@app.route("/cameras/cadastrar", methods=["GET", "POST"])
@login_required
def cameras_cadastrar():
    if not _is_can_manage():
        flash("Acesso restrito a administradores, gestores e key users.", "danger")
        return redirect(url_for("cameras_index"))
    site = session.get("user_site") or ""
    if request.method == "POST":
        modo = request.form.get("modo", "individual")
        criados = 0
        erros   = []
        if modo == "lote":
            rng     = (request.form.get("range_lote") or "").strip()
            prefixo = (request.form.get("prefixo_lote") or "CAM").strip()
            try:
                partes = rng.split("-")
                ini, fim = int(partes[0]), int(partes[1])
                pad = max(len(partes[0]), len(partes[1]))
                for n in range(ini, fim + 1):
                    num = f"{prefixo}-{str(n).zfill(pad)}"
                    if not Camera.query.filter_by(site=site, numero=num).first():
                        db.session.add(Camera(site=site, numero=num))
                        criados += 1
                db.session.commit()
            except Exception as exc:
                db.session.rollback()
                flash(f"Erro ao criar lote: {exc}", "danger")
                return redirect(url_for("cameras_cadastrar"))
        else:
            numeros = request.form.getlist("numero[]")
            nomes   = request.form.getlist("nome[]")
            for num, nome in zip(numeros, nomes):
                num  = (num or "").strip()
                nome = (nome or "").strip() or None
                if not num:
                    continue
                if Camera.query.filter_by(site=site, numero=num).first():
                    erros.append(f"{num} já existe")
                    continue
                db.session.add(Camera(site=site, numero=num, nome=nome))
                criados += 1
            try:
                db.session.commit()
            except Exception as exc:
                db.session.rollback()
                flash(f"Erro ao salvar: {exc}", "danger")
                return redirect(url_for("cameras_cadastrar"))

        flash(f"✓ {criados} câmera(s) cadastrada(s) com sucesso!", "success")
        if erros:
            flash(f"⚠️ {len(erros)} câmera(s) ignorada(s) (já existiam): {', '.join(erros[:5])}", "warning")
        return redirect(url_for("cameras_index"))

    return render_template("cameras_cadastro.html", site=site)


@app.route("/cameras/<int:cam_id>/toggle", methods=["POST"])
@login_required
def cameras_toggle(cam_id):
    if not _is_can_manage():
        flash("Acesso restrito a administradores, gestores e key users.", "danger")
        return redirect(url_for("cameras_index"))
    if _is_privileged():
        cam = Camera.query.filter_by(id=cam_id).first_or_404()
    else:
        site = session.get("user_site") or ""
        cam = Camera.query.filter_by(id=cam_id, site=site).first_or_404()
    cam.ativo = 0 if cam.ativo else 1
    db.session.commit()
    flash(f"Câmera {cam.numero} {'ativada' if cam.ativo else 'desativada'}.", "success")
    return redirect(url_for("cameras_index"))


@app.route("/cameras/<int:cam_id>/editar", methods=["POST"])
@login_required
def cameras_editar(cam_id):
    if not _is_can_manage():
        flash("Acesso restrito a administradores, gestores e key users.", "danger")
        return redirect(url_for("cameras_index"))
    if _is_privileged():
        cam = Camera.query.filter_by(id=cam_id).first_or_404()
    else:
        site = session.get("user_site") or ""
        cam = Camera.query.filter_by(id=cam_id, site=site).first_or_404()
    cam.numero = (request.form.get("numero") or cam.numero).strip()
    cam.nome   = (request.form.get("nome") or "").strip() or None
    db.session.commit()
    flash(f"Câmera {cam.numero} atualizada.", "success")
    return redirect(url_for("cameras_index"))


@app.route("/cameras/excluir-lote", methods=["POST"])
@login_required
def cameras_excluir_lote():
    if not _is_can_manage():
        flash("Acesso restrito a administradores, gestores e key users.", "danger")
        return redirect(url_for("cameras_index"))
    ids  = request.form.getlist("ids")
    site = session.get("user_site") or ""
    is_admin = _is_privileged()
    if not ids:
        flash("Nenhuma câmera selecionada.", "warning")
        return redirect(url_for("cameras_index"))
    removidas = 0
    for cam_id in ids:
        q = Camera.query.filter_by(id=cam_id)
        if not is_admin:
            q = q.filter_by(site=site)
        cam = q.first()
        if cam:
            cam.ativo = False
            removidas += 1
    try:
        db.session.commit()
        flash(f"{removidas} câmera(s) desativada(s) com sucesso.", "success")
    except Exception as e:
        db.session.rollback()
        flash(f"Erro ao desativar câmeras: {e}", "danger")
    return redirect(url_for("cameras_index"))


@app.route("/cameras/checklist/novo")
@login_required
def cameras_checklist_novo():
    site    = session.get("user_site") or ""
    cameras = Camera.query.filter_by(site=site, ativo=1).order_by(Camera.numero).all()
    if not cameras:
        flash("Nenhuma câmera ativa cadastrada para este site. Cadastre câmeras primeiro.", "warning")
        return redirect(url_for("cameras_index"))
    agora = datetime.now()
    return render_template("cameras_checklist.html",
        cameras=cameras, site=site,
        data_hoje=agora.strftime("%Y-%m-%d"),
        hora_agora=agora.strftime("%H:%M"),
        operador=session.get("user_nome", ""))


@app.route("/cameras/checklist/salvar", methods=["POST"])
@login_required
def cameras_checklist_salvar():
    site           = session.get("user_site") or ""
    cameras_ativas = Camera.query.filter_by(site=site, ativo=1).all()
    data      = (request.form.get("data_checklist") or "").strip()
    hora      = (request.form.get("hora_checklist") or "").strip()
    operador  = (request.form.get("operador") or "").strip()
    tratativa = (request.form.get("tratativa") or "").strip()
    total       = len(cameras_ativas)
    inoperantes = 0
    itens_dados = []
    for cam in cameras_ativas:
        status = request.form.get(f"status_{cam.id}", "FUNCIONAL")
        motivo = (request.form.get(f"motivo_{cam.id}") or "").strip() if status == "INOPERANTE" else ""
        if status == "INOPERANTE":
            inoperantes += 1
        itens_dados.append((cam.id, status, motivo))
    funcionais  = total - inoperantes
    anexo_nome  = None
    anexo_dados = None
    f = request.files.get("anexo")
    if f and f.filename:
        from werkzeug.utils import secure_filename
        anexo_nome  = secure_filename(f.filename)
        anexo_dados = f.read()
    chk = ChecklistCamera(
        site=site, data_checklist=data, hora_checklist=hora,
        operador=operador, tratativa=tratativa,
        total=total, funcionais=funcionais, inoperantes=inoperantes,
        anexo_nome=anexo_nome, anexo_dados=anexo_dados,
    )
    try:
        db.session.add(chk)
        db.session.flush()
        for cam_id_item, status, motivo in itens_dados:
            db.session.add(ChecklistCameraItem(
                checklist_id=chk.id, camera_id=cam_id_item,
                status=status, motivo=motivo or None
            ))
        db.session.commit()
    except Exception as exc:
        db.session.rollback()
        flash(f"Erro ao salvar checklist: {exc}", "danger")
        return redirect(url_for("cameras_checklist_novo"))
    flash("✓ Checklist salvo com sucesso!", "success")
    return redirect(url_for("cameras_checklist_ver", chk_id=chk.id))


@app.route("/cameras/checklist/<int:chk_id>")
@login_required
def cameras_checklist_ver(chk_id):
    site = session.get("user_site") or ""
    chk  = ChecklistCamera.query.filter_by(id=chk_id, site=site).first_or_404()
    itens       = chk.itens.order_by(ChecklistCameraItem.id).all()
    inoperantes = [i for i in itens if i.status == "INOPERANTE"]
    return render_template("cameras_checklist_ver.html",
        chk=chk, itens=itens, inoperantes=inoperantes, site=site)


@app.route("/cameras/checklist/<int:chk_id>/pdf")
@login_required
def cameras_checklist_pdf(chk_id):
    site = session.get("user_site") or ""
    ChecklistCamera.query.filter_by(id=chk_id, site=site).first_or_404()
    buf  = _gerar_pdf_checklist_cameras(chk_id)
    nome = f"checklist_cameras_{chk_id}.pdf"
    return send_file(buf, as_attachment=True, download_name=nome, mimetype="application/pdf")


@app.route("/cameras/checklist/<int:chk_id>/anexo")
@login_required
def cameras_checklist_anexo(chk_id):
    site = session.get("user_site") or ""
    chk  = ChecklistCamera.query.filter_by(id=chk_id, site=site).first_or_404()
    if not chk.anexo_dados:
        flash("Sem anexo.", "warning")
        return redirect(url_for("cameras_checklist_ver", chk_id=chk_id))
    return send_file(BytesIO(chk.anexo_dados), as_attachment=True,
                     download_name=chk.anexo_nome or "anexo")


def _is_privileged():
    """Retorna True se o usuário logado for ADMIN (único perfil com acesso cross-site irrestrito).
    MULTISITES, GESTOR, KEYUSER e OPERACIONAL são filtrados pelos sites vinculados (UsuarioSite)."""
    return (session.get("user_perfil") or "").upper() == "ADMIN"


def _is_can_manage():
    """Retorna True se o usuário pode gerenciar recursos (criar/editar).
    Inclui KEYUSER além dos perfis de _is_privileged().
    KEYUSER opera apenas dentro do próprio site (sem acesso cross-site).
    """
    return (session.get("user_perfil") or "").upper() in ("ADMIN", "GESTOR", "MULTISITES", "KEYUSER")


@app.context_processor
def _ctx_multisites():
    """Injeta multisites_lista (sites autorizados) em todos os templates."""
    lista = []
    if (session.get("user_perfil") or "").upper() == "MULTISITES":
        uid = session.get("user_id")
        if uid:
            lista = [v.site_nome for v in
                     UsuarioSite.query.filter_by(usuario_id=uid)
                     .order_by(UsuarioSite.site_nome).all()]
    return {"multisites_lista": lista}

def normalizar_status(valor):
    return (valor or "").strip().upper()


def normalizar_prioridade(valor):
    return (valor or "").strip().upper()


def aplicar_filtros(query):
    data_inicial = (request.args.get("data_inicial") or "").strip()
    data_final   = (request.args.get("data_final")   or "").strip()
    local        = (request.args.get("local")         or "").strip()
    natureza     = (request.args.get("natureza")      or "").strip()
    status       = (request.args.get("status")        or "").strip().upper()
    operador     = (request.args.get("operador")      or "").strip()
    site_f       = (request.args.get("site_filtro")   or "").strip()

    # Empurra todos os filtros para o SQL — sem carregar todos os registros em Python.
    # data_hora armazenada como "YYYY-MM-DDTHH:MM" — SUBSTR(1,10) extrai só a data.
    if data_inicial:
        query = query.filter(func.substr(Ocorrencia.data_hora, 1, 10) >= data_inicial)
    if data_final:
        query = query.filter(func.substr(Ocorrencia.data_hora, 1, 10) <= data_final)
    if local:
        query = query.filter(func.lower(Ocorrencia.local).contains(local.lower()))
    if natureza:
        query = query.filter(func.lower(Ocorrencia.natureza) == natureza.lower())
    if status:
        query = query.filter(func.upper(Ocorrencia.status) == status)
    if operador:
        query = query.filter(func.lower(Ocorrencia.operador).contains(operador.lower()))
    if site_f:
        query = query.filter(Ocorrencia.site == site_f)

    registros = query.all()

    filtros = {
        "data_inicial": data_inicial,
        "data_final":   data_final,
        "local":        local,
        "natureza":     natureza,
        "status":       status,
        "operador":     operador,
        "site_filtro":  site_f,
    }
    return registros, filtros


# =========================
# AUTH
# =========================
@app.route("/")
def index():
    if session.get("user_id"):
        return redirect(url_for("ocorrencias"))
    return redirect(url_for("login"))



@app.route("/analise-investigativa", methods=["GET"])
@login_required
def analise_investigativa():
    site_usuario = session.get("user_site") or ""
    site_atual = SiteCompleto.query.filter_by(nome_do_site=site_usuario).first() if site_usuario else None
    proximo_numero = AnaliseInvestigativa.query.filter_by(site=site_usuario or None).count() + 1
    return render_template(
        "analise_investigativa.html",
        site_atual=site_atual,
        proximo_numero=proximo_numero,
        user_nome=session.get("user_nome", ""),
    )


# =========================
# DASHBOARD — ANÁLISES
# =========================
@app.route("/analises/dashboard")
@login_required
def dashboard_analise():
    from datetime import datetime as _dt
    from collections import Counter

    is_admin    = (session.get("user_perfil") or "").upper() == "ADMIN"
    site_usuario = session.get("user_site") or None
    _hoje = _dt.now()

    # Dashboard só conta/agrega e mostra recentes[:10] (template não acessa os
    # CLOBs de texto), então defere todos os campos pesados — blobs e textos.
    _sem_blob = [
        defer(AnaliseInvestigativa.docx_arquivo),
        defer(AnaliseInvestigativa.anexo_fechamento),
        defer(AnaliseInvestigativa.objetivo),
        defer(AnaliseInvestigativa.descricao_registro),
        defer(AnaliseInvestigativa.conclusao),
        defer(AnaliseInvestigativa.sugestao),
        defer(AnaliseInvestigativa.texto_fechamento),
    ]
    if is_admin:
        registros = AnaliseInvestigativa.query.options(*_sem_blob).order_by(AnaliseInvestigativa.id.desc()).all()
    else:
        registros = _query_filtrar_sites(
            AnaliseInvestigativa.query.options(*_sem_blob), AnaliseInvestigativa
        ).order_by(AnaliseInvestigativa.id.desc()).all()

    # Filtros simples por querystring
    f_data_ini   = request.args.get("data_inicial", "")
    f_data_fim   = request.args.get("data_final", "")
    f_status     = request.args.get("status", "")
    f_classif    = request.args.get("classificacao", "")
    f_site       = request.args.get("site_filtro", "") if is_admin else ""

    filtrados = registros
    if f_data_ini:
        try:
            _di = _dt.strptime(f_data_ini, "%Y-%m-%d")
            filtrados = [r for r in filtrados if r.criado_em and r.criado_em >= _di]
        except Exception: pass
    if f_data_fim:
        try:
            _df = _dt.strptime(f_data_fim, "%Y-%m-%d").replace(hour=23, minute=59, second=59)
            filtrados = [r for r in filtrados if r.criado_em and r.criado_em <= _df]
        except Exception: pass
    if f_status:
        filtrados = [r for r in filtrados if (r.status_analise or "").upper() == f_status.upper()]
    if f_classif:
        filtrados = [r for r in filtrados if f_classif.lower() in (r.classificacao or "").lower()]
    if f_site:
        filtrados = [r for r in filtrados if (r.site or "") == f_site]

    # Sites disponíveis para o filtro
    if is_admin:
        todos_sites_dash = _get_distinct_cached(AnaliseInvestigativa.site, "ai_sites")
    else:
        todos_sites_dash = sorted(s for s in _sites_do_usuario() if s)

    total      = len(filtrados)
    andamento  = len([r for r in filtrados if (r.status_analise or "").upper() == "EM ANDAMENTO"])
    fechadas   = len([r for r in filtrados if (r.status_analise or "").upper() == "FECHADA"])
    taxa_resolucao = round(fechadas / total * 100) if total > 0 else 0

    def _mesmo_mes(r):
        try:
            return r.criado_em and r.criado_em.month == _hoje.month and r.criado_em.year == _hoje.year
        except Exception:
            return False
    registros_mes = len([r for r in filtrados if _mesmo_mes(r)])

    # Contagens para gráficos
    status_count   = Counter((r.status_analise or "Não informado").upper() for r in filtrados)
    classif_count  = Counter(r.classificacao or "Não informado" for r in filtrados)
    criador_count  = Counter(r.criado_por or "Não informado" for r in filtrados)
    resp_count     = Counter(r.responsavel or "Não informado" for r in filtrados)

    # Status em ordem fixa
    _STATUS_ORDER  = ["EM ANDAMENTO", "FECHADA"]
    labels_status  = []
    valores_status = []
    for _s in _STATUS_ORDER:
        if _s in status_count:
            labels_status.append(_s)
            valores_status.append(status_count[_s])
    for _s, _v in status_count.items():
        if _s not in _STATUS_ORDER:
            labels_status.append(_s)
            valores_status.append(_v)

    # Resto ordenado crescente
    classif_sorted = sorted(classif_count.items(), key=lambda x: x[1])
    criador_sorted = sorted(criador_count.items(), key=lambda x: x[1])
    resp_sorted    = sorted(resp_count.items(),    key=lambda x: x[1])

    # Todos os status e classificações distintos para filtros
    todas_classif = sorted(set(r.classificacao for r in registros if r.classificacao))

    return render_template(
        "dashboard_analise.html",
        is_admin=is_admin,
        todos_sites_dash=todos_sites_dash,
        recentes=filtrados[:10],
        filtros={"data_inicial": f_data_ini, "data_final": f_data_fim,
                 "status": f_status, "classificacao": f_classif, "site_filtro": f_site},
        todas_classif=todas_classif,
        resumo={
            "total":          total,
            "andamento":      andamento,
            "fechadas":       fechadas,
            "taxa_resolucao": taxa_resolucao,
            "registros_mes":  registros_mes,
            "valor_total":    _formatar_valor(sum(_parse_valor(r.valor) for r in filtrados if r.valor)),
        },
        labels_status=labels_status,
        valores_status=valores_status,
        labels_classif=[x[0] for x in classif_sorted],
        valores_classif=[x[1] for x in classif_sorted],
        labels_criador=[x[0] for x in criador_sorted],
        valores_criador=[x[1] for x in criador_sorted],
        labels_resp=[x[0] for x in resp_sorted],
        valores_resp=[x[1] for x in resp_sorted],
    )


@app.route("/analises")
@login_required
def analises():
    is_admin = (session.get("user_perfil") or "").upper() == "ADMIN"
    site_usuario = session.get("user_site") or None

    # Defere blobs + CLOBs não exibidos na lista. objetivo/conclusao/texto_fechamento
    # são lidos por linha no template (data-attrs do modal), então permanecem carregados;
    # descricao_registro e sugestao não são usados → deferidos.
    _sem_blob = [
        defer(AnaliseInvestigativa.docx_arquivo),
        defer(AnaliseInvestigativa.anexo_fechamento),
        defer(AnaliseInvestigativa.descricao_registro),
        defer(AnaliseInvestigativa.sugestao),
    ]
    if is_admin:
        query = AnaliseInvestigativa.query.options(*_sem_blob).order_by(AnaliseInvestigativa.id.desc())
    else:
        query = _query_filtrar_sites(
            AnaliseInvestigativa.query.options(*_sem_blob), AnaliseInvestigativa
        ).order_by(AnaliseInvestigativa.id.desc())

    # Filtros server-side
    data_inicial  = (request.args.get("data_inicial")  or "").strip()
    data_final    = (request.args.get("data_final")    or "").strip()
    status_f      = (request.args.get("status")        or "").strip().upper()
    empresa_f     = (request.args.get("empresa")       or "").strip().lower()
    responsavel_f = (request.args.get("responsavel")   or "").strip().lower()
    site_filtro   = (request.args.get("site_filtro")   or "").strip()

    # Empurra filtros para o SQL — evita carregar todos os registros em memória
    if data_inicial:
        try:
            query = query.filter(
                AnaliseInvestigativa.criado_em >= datetime.strptime(data_inicial, "%Y-%m-%d")
            )
        except ValueError:
            pass
    if data_final:
        try:
            query = query.filter(
                AnaliseInvestigativa.criado_em < datetime.strptime(data_final, "%Y-%m-%d") + timedelta(days=1)
            )
        except ValueError:
            pass
    if status_f:
        query = query.filter(func.upper(AnaliseInvestigativa.status_analise) == status_f)
    if empresa_f:
        query = query.filter(func.lower(AnaliseInvestigativa.empresa).contains(empresa_f))
    if responsavel_f:
        query = query.filter(func.lower(AnaliseInvestigativa.responsavel).contains(responsavel_f))
    if site_filtro:
        query = query.filter(AnaliseInvestigativa.site == site_filtro)

    registros = query.all()

    filtros = {
        "data_inicial": data_inicial, "data_final": data_final,
        "status": status_f, "empresa": empresa_f,
        "responsavel": responsavel_f, "site_filtro": site_filtro,
    }
    # Sites distintos via query leve — sem carregar todos os registros
    todos_sites_analise = sorted(
        s[0] for s in db.session.query(AnaliseInvestigativa.site).distinct().all() if s[0]
    ) if is_admin else []

    resumo = {
        "total":       len(registros),
        "sites":       len(set(r.site for r in registros if r.site)),
        "valor_total": _formatar_valor(sum(_parse_valor(r.valor) for r in registros if r.valor)),
    }
    return render_template("analises.html", registros=registros, resumo=resumo,
                           is_admin=is_admin, filtros=filtros, todos_sites=todos_sites_analise)


@app.route("/analises/excluir/<int:analise_id>", methods=["POST"])
@login_required
@perfil_required("ADMIN")
def excluir_analise(analise_id):
    registro = AnaliseInvestigativa.query.get_or_404(analise_id)
    db.session.delete(registro)
    db.session.commit()
    flash("Análise excluída com sucesso.", "success")
    return redirect(url_for("analises"))


@app.route("/analises/<int:analise_id>/fechar", methods=["POST"])
@login_required
def fechar_analise(analise_id):
    registro = AnaliseInvestigativa.query.get_or_404(analise_id)
    texto = (request.form.get("texto_fechamento") or "").strip()
    if not texto:
        flash("Informe o texto de fechamento.", "warning")
        return redirect(url_for("analises"))

    registro.status_analise = "FECHADA"
    registro.texto_fechamento = texto
    registro.fechado_por = session.get("user_nome")
    registro.fechado_em = datetime.utcnow()

    anexo = request.files.get("anexo_fechamento")
    if anexo and anexo.filename:
        ext = anexo.filename.rsplit(".", 1)[-1].lower()
        if ext in {"pdf", "doc", "docx", "xlsx", "png", "jpg", "jpeg"}:
            registro.anexo_fechamento_nome = anexo.filename
            registro.anexo_fechamento = anexo.read()

    db.session.commit()
    flash("Análise fechada com sucesso.", "success")
    return redirect(url_for("analises"))


@app.route("/analises/<int:analise_id>/anexo")
@login_required
def download_anexo_analise(analise_id):
    registro = AnaliseInvestigativa.query.get_or_404(analise_id)
    if not registro.anexo_fechamento:
        flash("Anexo não encontrado.", "warning")
        return redirect(url_for("analises"))
    return send_file(
        BytesIO(registro.anexo_fechamento),
        as_attachment=True,
        download_name=registro.anexo_fechamento_nome or "anexo",
    )


@app.route("/analises/<int:analise_id>/editar", methods=["GET", "POST"])
@login_required
def editar_analise(analise_id):
    registro = AnaliseInvestigativa.query.get_or_404(analise_id)
    is_admin  = (session.get("user_perfil") or "").upper() == "ADMIN"
    is_criador = registro.criado_por == session.get("user_nome", "")

    if not is_admin and not is_criador:
        flash("Você não tem permissão para editar esta análise.", "danger")
        return redirect(url_for("analises"))

    if (registro.status_analise or "").upper() == "FECHADA":
        flash("Esta análise já foi fechada e não pode ser editada.", "warning")
        return redirect(url_for("analises"))

    if request.method == "POST":
        f = request.form
        registro.empresa            = (f.get("empresa") or "").strip()
        registro.unidade            = (f.get("unidade") or "").strip()
        registro.endereco           = (f.get("endereco") or "").strip()
        registro.classificacao      = (f.get("classificacao") or "").strip()
        registro.produtos_segmento  = (f.get("produtos_segmento") or "").strip()
        registro.clientes           = (f.get("clientes") or "").strip()
        registro.objetivo           = (f.get("objetivo") or "").strip()
        registro.responsavel        = (f.get("responsavel") or "").strip()
        registro.funcao_levantamento= (f.get("funcao_levantamento") or "").strip()
        registro.data_levantamento  = (f.get("data_levantamento") or "").strip()
        registro.descricao_registro = (f.get("descricao_registro") or "").strip()
        registro.conclusao          = (f.get("conclusao") or "").strip()
        registro.sugestao           = (f.get("sugestao") or "").strip()
        registro.valor              = (f.get("valor") or "").strip() or None

        # Novas imagens enviadas no edit
        novos_files = request.files.getlist("imagens[]")
        novas_descricoes = request.form.getlist("descricoes[]")
        if len(novas_descricoes) < len(novos_files):
            novas_descricoes += [""] * (len(novos_files) - len(novas_descricoes))
        for img_file, desc in zip(novos_files, novas_descricoes):
            if img_file and img_file.filename and allowed_file(img_file.filename):
                img_b64 = base64.b64encode(img_file.read()).decode("utf-8")
                nova_img = ImagemAnaliseInvestigativa(
                    analise_id=registro.id,
                    arquivo_b64=img_b64,
                    descricao=desc,
                )
                db.session.add(nova_img)

        # Imagens a remover (checkboxes marcados pelo usuário)
        remover_ids = request.form.getlist("remover_imagem[]")
        for img_id in remover_ids:
            img_obj = ImagemAnaliseInvestigativa.query.get(int(img_id))
            if img_obj and img_obj.analise_id == registro.id:
                db.session.delete(img_obj)

        # Regenera o DOCX salvo
        registro.docx_arquivo = None   # força regeneração no próximo download
        db.session.commit()
        flash("Análise atualizada com sucesso.", "success")
        return redirect(url_for("analises"))

    sites = [s.nome_do_site for s in SiteCompleto.query.order_by(SiteCompleto.nome_do_site).all()]
    return render_template("editar_analise.html", registro=registro, sites=sites)


def classify_line(valor):
    valor = (valor or "").strip()
    if not valor:
        return "-"
    return valor


def style_heading(paragraph, text, size=12, bold=True, align="left"):
    paragraph.text = ""

    if align == "center":
        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    else:
        paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT

    run = paragraph.add_run(text)
    run.bold = bold
    run.font.name = "Arial"
    run.font.size = Pt(size)


def set_cell_background(cell, color_hex):
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), color_hex)
    tc_pr.append(shd)


def add_grid_table(doc, data, col_widths_cm=None, header_bold_cols=None):
    table = doc.add_table(rows=len(data), cols=len(data[0]))
    table.style = "Table Grid"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER

    header_bold_cols = header_bold_cols or []

    for r_idx, row in enumerate(data):
        for c_idx, value in enumerate(row):
            cell = table.rows[r_idx].cells[c_idx]
            cell.text = str(value or "")

            if col_widths_cm and c_idx < len(col_widths_cm):
                cell.width = Cm(col_widths_cm[c_idx])

            for paragraph in cell.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Arial"
                    run.font.size = Pt(9)

                    if c_idx in header_bold_cols:
                        run.bold = True

            if c_idx in header_bold_cols:
                set_cell_background(cell, "F2F2F2")

    return table


def set_cell_bg(cell, color):
    tcPr = cell._tc.get_or_add_tcPr()
    for existing in tcPr.findall(qn("w:shd")):
        tcPr.remove(existing)
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), color)
    tcPr.append(shd)


def set_cell_border(cell, color="000000", size="8"):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for existing in tcPr.findall(qn("w:tcBorders")):
        tcPr.remove(existing)
    borders = OxmlElement("w:tcBorders")

    for border_name in ["top", "left", "bottom", "right"]:
        border = OxmlElement(f"w:{border_name}")
        border.set(qn("w:val"), "single")
        border.set(qn("w:sz"), size)
        border.set(qn("w:space"), "0")
        border.set(qn("w:color"), color)
        borders.append(border)

    tcPr.append(borders)


def set_table_borders(table, color="000000", size="8"):
    """Define bordas no nível da tabela — garante que insideH/insideV apareçam."""
    tbl = table._tbl
    tblPr = tbl.find(qn("w:tblPr"))
    if tblPr is None:
        tblPr = OxmlElement("w:tblPr")
        tbl.insert(0, tblPr)
    for existing in tblPr.findall(qn("w:tblBorders")):
        tblPr.remove(existing)
    tblBorders = OxmlElement("w:tblBorders")
    for lado in ("top", "left", "bottom", "right", "insideH", "insideV"):
        border = OxmlElement(f"w:{lado}")
        border.set(qn("w:val"), "single")
        border.set(qn("w:sz"), size)
        border.set(qn("w:space"), "0")
        border.set(qn("w:color"), color)
        tblBorders.append(border)
    tblPr.append(tblBorders)


def format_cell(cell, bold=False, bg=None, align="left"):
    set_cell_border(cell)
    if bg:
        set_cell_bg(cell, bg)
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER

    for p in cell.paragraphs:
        p.paragraph_format.space_before = Pt(3)
        p.paragraph_format.space_after = Pt(3)

        if align == "center":
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        else:
            p.alignment = WD_ALIGN_PARAGRAPH.LEFT

        for run in p.runs:
            run.font.name = "Arial"
            run.font.size = Pt(9)
            run.bold = bold
            run.font.color.rgb = RGBColor(31, 41, 55)


def add_section_title(doc, title):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(14)
    p.paragraph_format.space_after = Pt(6)

    run = p.add_run(title.upper())
    run.bold = True
    run.font.name = "Arial"
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(212, 5, 17)

    return p


def add_professional_table(doc, rows, col_widths):
    table = doc.add_table(rows=0, cols=len(col_widths))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.autofit = False
    set_table_borders(table)

    for row_data in rows:
        row = table.add_row()
        for i, value in enumerate(row_data):
            cell = row.cells[i]
            cell.width = Cm(col_widths[i])
            cell.text = str(value or "")

            is_label = i % 2 == 0
            format_cell(
                cell,
                bold=is_label,
                bg="FFF2CC" if is_label else "FFFFFF",
                align="left"
            )

    doc.add_paragraph().paragraph_format.space_after = Pt(3)
    return table


def add_text_box(doc, text="\n\n\n"):
    table = doc.add_table(rows=1, cols=1)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    set_table_borders(table)
    cell = table.rows[0].cells[0]
    cell.text = text or "\n\n\n"
    format_cell(cell, bold=False, bg="FFFFFF")
    return table

@app.route("/gerar_docx", methods=["POST"])
@login_required
def gerar_analise_investigativa():
    f = request.form

    # ===== Salva no banco =====
    site_val = (f.get("site") or "").strip()
    _codigo_analise, _seq_analise = gerar_codigo(AnaliseInvestigativa, site_val)
    nova_analise = AnaliseInvestigativa(
        codigo=_codigo_analise,
        numero_site=_seq_analise,
        site=site_val,
        id_relatorio=(f.get("id_relatorio") or "").strip(),
        empresa=(f.get("empresa") or "").strip(),
        unidade=(f.get("unidade") or "").strip(),
        endereco=(f.get("endereco") or "").strip(),
        classificacao=(f.get("classificacao") or "").strip(),
        produtos_segmento=(f.get("produtos_segmento") or "").strip(),
        clientes=(f.get("clientes") or "").strip(),
        objetivo=(f.get("objetivo") or "").strip(),
        responsavel=(f.get("responsavel") or "").strip(),
        nome_funcao_data=(f.get("nome_funcao_data") or "").strip(),   # legado
        funcao_levantamento=(f.get("funcao_levantamento") or "").strip(),
        data_levantamento=(f.get("data_levantamento") or "").strip(),
        descricao_registro=(f.get("descricao_registro") or "").strip(),
        conclusao=(f.get("conclusao") or "").strip(),
        sugestao=(f.get("sugestao") or "").strip(),
        valor=(f.get("valor") or "").strip() or None,
        criado_por=session.get("user_nome"),
    )
    db.session.add(nova_analise)
    db.session.commit()

    # ===== Salvar imagens como base64 no banco =====
    files = request.files.getlist("imagens[]")
    descricoes = request.form.getlist("descricoes[]")
    if len(descricoes) < len(files):
        descricoes += [""] * (len(files) - len(descricoes))

    for fimg, desc in zip(files, descricoes):
        if fimg and fimg.filename and allowed_file(fimg.filename):
            img_b64 = base64.b64encode(fimg.read()).decode("utf-8")
            nova_img = ImagemAnaliseInvestigativa(
                analise_id=nova_analise.id,
                arquivo_b64=img_b64,
                descricao=desc,
            )
            db.session.add(nova_img)

    db.session.commit()

    flash("Análise salva com sucesso! Clique em Detalhes para baixar o documento.", "success")
    return redirect(url_for("analises"))




@app.route("/analises/confirmar/<int:analise_id>")
@login_required
def confirmar_analise(analise_id):
    # Rota mantida para retrocompatibilidade com links antigos
    return redirect(url_for("analises"))


@app.route("/analises/download/<int:analise_id>")
@login_required
def download_analise(analise_id):
    registro = AnaliseInvestigativa.query.get_or_404(analise_id)
    filename = f"A.I - {registro.id_relatorio or registro.codigo or registro.id} - {registro.site or 'SEM_SITE'}.docx"

    if registro.docx_arquivo:
        # Registro antigo — retorna o BLOB armazenado
        buf = BytesIO(registro.docx_arquivo)
        buf.seek(0)
    else:
        # Registro novo — gera o DOCX na hora a partir dos dados salvos
        buf = gerar_docx_de_registro(registro)

    return send_file(buf, as_attachment=True, download_name=filename,
                     mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document")


# ==========================
# DASHBOARD — ANC
# ==========================
@app.route("/anc/dashboard")
@login_required
def dashboard_anc():
    from datetime import datetime as _dt
    from collections import Counter

    is_admin     = (session.get("user_perfil") or "").upper() == "ADMIN"
    site_usuario = session.get("user_site") or None
    _hoje = _dt.now()

    _sem_imgs = [defer(ANC.imagem_1), defer(ANC.imagem_2), defer(ANC.imagem_3),
                 defer(ANC.imagem_4), defer(ANC.imagem_5), defer(ANC.imagem_6),
                 defer(ANC.anexo_fechamento), defer(ANC.descricao), defer(ANC.plano_acao_texto)]

    if is_admin:
        base = ANC.query.options(*_sem_imgs).order_by(ANC.id.desc()).all()
    else:
        base = _query_filtrar_sites(
            ANC.query.options(*_sem_imgs), ANC
        ).order_by(ANC.id.desc()).all()

    # Filtros por querystring
    f_data_ini  = request.args.get("data_inicial", "")
    f_data_fim  = request.args.get("data_final", "")
    f_status    = request.args.get("status", "")
    f_gravidade = request.args.get("gravidade", "")
    f_site      = request.args.get("site_filtro", "") if is_admin else ""

    filtrados = base
    if f_data_ini:
        try:
            _di = _dt.strptime(f_data_ini, "%Y-%m-%d")
            filtrados = [r for r in filtrados if r.criado_em and r.criado_em >= _di]
        except Exception: pass
    if f_data_fim:
        try:
            _df = _dt.strptime(f_data_fim, "%Y-%m-%d").replace(hour=23, minute=59, second=59)
            filtrados = [r for r in filtrados if r.criado_em and r.criado_em <= _df]
        except Exception: pass
    if f_status:
        filtrados = [r for r in filtrados if (r.status or "").upper() == f_status.upper()]
    if f_gravidade:
        filtrados = [r for r in filtrados if (r.gravidade or "").upper() == f_gravidade.upper()]
    if f_site:
        filtrados = [r for r in filtrados if (r.site or "") == f_site]

    # Sites disponíveis para o filtro
    if is_admin:
        _sites_q = db.session.query(ANC.site).distinct().all()
        todos_sites_dash = sorted(s[0] for s in _sites_q if s[0])
    else:
        todos_sites_dash = sorted(s for s in _sites_do_usuario() if s)

    total      = len(filtrados)
    abertos    = len([r for r in filtrados if (r.status or "").upper() == "ABERTO"])
    andamento  = len([r for r in filtrados if (r.status or "").upper() == "EM ANDAMENTO"])
    concluidos = len([r for r in filtrados if (r.status or "").upper() == "CONCLUÍDO"])
    criticos   = len([r for r in filtrados if (r.gravidade or "").upper() == "CRÍTICA"])
    taxa_resolucao = round(concluidos / total * 100) if total > 0 else 0

    def _mesmo_mes(r):
        try:
            return r.criado_em and r.criado_em.month == _hoje.month and r.criado_em.year == _hoje.year
        except Exception:
            return False
    registros_mes = len([r for r in filtrados if _mesmo_mes(r)])

    # Contagens para gráficos
    status_count    = Counter((r.status or "Não informado").upper() for r in filtrados)
    gravidade_count = Counter((r.gravidade or "Não informado").upper() for r in filtrados)
    natureza_count  = Counter(r.natureza or "Não informado" for r in filtrados)
    setor_count     = Counter(r.setor or "Não informado" for r in filtrados)
    criador_count   = Counter(r.criado_por or "Não informado" for r in filtrados)

    # Status em ordem fixa
    _STATUS_ORDER  = ["ABERTO", "EM ANDAMENTO", "CONCLUÍDO"]
    labels_status  = []
    valores_status = []
    for _s in _STATUS_ORDER:
        if _s in status_count:
            labels_status.append(_s)
            valores_status.append(status_count[_s])
    for _s, _v in status_count.items():
        if _s not in _STATUS_ORDER:
            labels_status.append(_s)
            valores_status.append(_v)

    # Gravidade em ordem fixa de severidade
    _GRAV_ORDER  = ["BAIXA", "MÉDIA", "ALTA", "CRÍTICA"]
    labels_grav  = []
    valores_grav = []
    for _g in _GRAV_ORDER:
        if _g in gravidade_count:
            labels_grav.append(_g)
            valores_grav.append(gravidade_count[_g])
    for _g, _v in gravidade_count.items():
        if _g not in _GRAV_ORDER:
            labels_grav.append(_g)
            valores_grav.append(_v)

    natureza_sorted = sorted(natureza_count.items(), key=lambda x: x[1])
    setor_sorted    = sorted(setor_count.items(),    key=lambda x: x[1])
    criador_sorted  = sorted(criador_count.items(),  key=lambda x: x[1])

    todas_gravidades = sorted(set(r.gravidade for r in base if r.gravidade))

    return render_template(
        "dashboard_anc.html",
        is_admin=is_admin,
        todos_sites_dash=todos_sites_dash,
        recentes=filtrados[:10],
        filtros={"data_inicial": f_data_ini, "data_final": f_data_fim,
                 "status": f_status, "gravidade": f_gravidade, "site_filtro": f_site},
        todas_gravidades=todas_gravidades,
        resumo={
            "total":          total,
            "abertos":        abertos,
            "andamento":      andamento,
            "concluidos":     concluidos,
            "criticos":       criticos,
            "taxa_resolucao": taxa_resolucao,
            "registros_mes":  registros_mes,
            "valor_total":    _formatar_valor(sum(_parse_valor(r.valor) for r in filtrados if r.valor)),
        },
        labels_status=labels_status,   valores_status=valores_status,
        labels_grav=labels_grav,       valores_grav=valores_grav,
        labels_natureza=[x[0] for x in natureza_sorted],
        valores_natureza=[x[1] for x in natureza_sorted],
        labels_setor=[x[0] for x in setor_sorted],
        valores_setor=[x[1] for x in setor_sorted],
        labels_criador=[x[0] for x in criador_sorted],
        valores_criador=[x[1] for x in criador_sorted],
    )


# ==========================
# ANC - Avisos de Não Conformidade
# ==========================
@app.route("/anc", methods=["GET", "POST"])
@login_required
def anc():
    is_admin = (session.get("user_perfil") or "").upper() == "ADMIN"
    site_usuario = session.get("user_site") or None

    registro_edicao = None
    modo_edicao = False
    editar_id = request.args.get("editar", type=int)
    if editar_id:
        # Verifica permissão antes de carregar o registro completo (com CLOBs)
        _chk = db.session.query(ANC.id, ANC.criado_por, ANC.status).filter_by(id=editar_id).first()
        if not _chk:
            flash("ANC não encontrado.", "danger")
            return redirect(url_for("anc"))
        _can_edit_anc = is_admin or _chk.criado_por == session.get("user_nome", "")
        _is_closed_anc = (_chk.status or "").upper() == "CONCLUÍDO"
        if not _can_edit_anc:
            flash("Você não tem permissão para editar esta ANC.", "danger")
            return redirect(url_for("anc"))
        if _is_closed_anc:
            flash("Esta ANC já foi concluída e não pode ser editada.", "warning")
            return redirect(url_for("anc"))
        registro_edicao = ANC.query.get(editar_id)
        modo_edicao = True

    if request.method == "POST":
        f = request.form
        anc_id = f.get("anc_id", type=int)

        data_nc           = (f.get("data_nc")           or "").strip()
        hora_nc           = (f.get("hora_nc")           or "").strip()
        site_val          = (f.get("site")              or site_usuario or "").strip()
        setor             = (f.get("setor")             or "").strip()
        tipo_ocorrencia   = (f.get("tipo_ocorrencia")   or "").strip()
        gravidade         = (f.get("gravidade")         or "").strip().upper()
        natureza          = (f.get("natureza")          or "").strip()
        responsavel       = (f.get("responsavel")       or "").strip()
        gestor_responsavel= (f.get("gestor_responsavel")or "").strip()
        cargo             = (f.get("cargo")             or "").strip()
        local_val         = (f.get("local")             or "").strip()
        envolvido         = (f.get("envolvido")         or "").strip()
        tipo              = (f.get("tipo")              or "").strip()
        turno             = (f.get("turno")             or "").strip()
        status            = (f.get("status")            or "ABERTO").strip().upper()
        descricao         = (f.get("descricao")         or "").strip()
        inicio_investigacao = (f.get("inicio_investigacao") or "").strip() or None
        fim_investigacao    = (f.get("fim_investigacao")    or "").strip() or None
        valor_anc           = (f.get("valor") or "").strip() or None

        if not data_nc or not hora_nc or not setor or not tipo_ocorrencia or not gravidade \
                or not natureza or not responsavel or not gestor_responsavel \
                or not local_val or not turno or not descricao:
            flash("Preencha todos os campos obrigatórios.", "warning")
            return redirect(url_for("anc", editar=anc_id) if anc_id else url_for("anc"))

        imgs = []
        for i in range(1, 7):
            field = request.files.get(f"imagem_{i}")
            if field and field.filename:
                b64, _ = arquivo_para_base64(field, EXTENSOES_PERMITIDAS_IMAGEM)
                imgs.append(b64)
            else:
                imgs.append(None)

        if anc_id:
            # Verifica permissão sem carregar CLOBs de imagem
            _chk = db.session.query(ANC.id, ANC.criado_por, ANC.status).filter_by(id=anc_id).first()
            if not _chk:
                flash("ANC não encontrado.", "danger")
                return redirect(url_for("anc"))
            _can_edit = is_admin or _chk.criado_por == session.get("user_nome", "")
            _is_closed = (_chk.status or "").upper() == "CONCLUÍDO"
            if not _can_edit:
                flash("Você não tem permissão para editar esta ANC.", "danger")
                return redirect(url_for("anc"))
            if _is_closed:
                flash("Esta ANC já foi concluída e não pode ser editada.", "warning")
                return redirect(url_for("anc"))
            reg = ANC.query.get(anc_id)
            reg.data_nc = data_nc; reg.hora_nc = hora_nc; reg.site = site_val
            reg.setor = setor; reg.tipo_ocorrencia = tipo_ocorrencia
            reg.gravidade = gravidade; reg.natureza = natureza
            reg.responsavel = responsavel; reg.gestor_responsavel = gestor_responsavel
            reg.cargo = cargo
            reg.local = local_val; reg.envolvido = envolvido
            reg.tipo = tipo; reg.turno = turno
            reg.descricao = descricao
            reg.inicio_investigacao = inicio_investigacao
            reg.fim_investigacao = fim_investigacao
            reg.valor = valor_anc
            for i, b64 in enumerate(imgs, start=1):
                if b64:
                    setattr(reg, f"imagem_{i}", b64)
            for i in range(1, 7):
                if f.get(f"remover_img_{i}"):
                    setattr(reg, f"imagem_{i}", None)
            db.session.commit()
            flash("ANC atualizado com sucesso.", "success")
            return redirect(url_for("anc"))

        _num_anc, _seq_anc = gerar_numero_anc(site_val)
        novo = ANC(
            numero_anc=_num_anc, numero_site=_seq_anc,
            data_nc=data_nc, hora_nc=hora_nc, site=site_val,
            setor=setor, tipo_ocorrencia=tipo_ocorrencia,
            gravidade=gravidade, natureza=natureza,
            responsavel=responsavel, gestor_responsavel=gestor_responsavel,
            cargo=cargo,
            local=local_val, envolvido=envolvido,
            tipo=tipo, turno=turno, status="PENDENTE",
            descricao=descricao,
            inicio_investigacao=inicio_investigacao,
            fim_investigacao=fim_investigacao,
            imagem_1=imgs[0], imagem_2=imgs[1],
            imagem_3=imgs[2], imagem_4=imgs[3],
            imagem_5=imgs[4], imagem_6=imgs[5],
            valor=valor_anc,
            criado_por=session.get("user_nome"),
        )
        db.session.add(novo)
        db.session.commit()
        flash("ANC registrado com sucesso.", "success")
        return redirect(url_for("anc"))

    # descricao e plano_acao_texto NÃO são deferidas: o template as acessa para cada
    # registro (data-descricao / data-plano-acao-texto), então deixá-las deferred geraria
    # N lazy-loads individuais na Oracle, travando a página com muitos ANCs.
    _sem_imgs = [defer(ANC.imagem_1), defer(ANC.imagem_2), defer(ANC.imagem_3),
                 defer(ANC.imagem_4), defer(ANC.imagem_5), defer(ANC.imagem_6),
                 defer(ANC.anexo_fechamento)]
    if is_admin:
        query = ANC.query.options(*_sem_imgs).order_by(ANC.id.desc())
    else:
        query = _query_filtrar_sites(
            ANC.query.options(*_sem_imgs), ANC
        ).order_by(ANC.id.desc())

    registros, filtros = aplicar_filtros_anc(query)

    resumo = {
        "total":      len(registros),
        "abertos":    len([r for r in registros if (r.status or "").upper() == "ABERTO"]),
        "andamento":  len([r for r in registros if (r.status or "").upper() == "EM ANDAMENTO"]),
        "concluidos": len([r for r in registros if (r.status or "").upper() == "CONCLUÍDO"]),
        "criticos":   len([r for r in registros if (r.gravidade or "").upper() == "CRÍTICA"]),
        "valor_total": _formatar_valor(sum(_parse_valor(r.valor) for r in registros if r.valor)),
    }

    todos_sites_anc = sorted(
        r[0] for r in db.session.query(ANC.site).distinct().all() if r[0]
    ) if is_admin else []
    naturezas_lista = _get_naturezas(site_usuario)

    # Solicitações de exclusão pendentes (visíveis apenas para admins)
    pendentes_exclusao = []
    ancs_excluidas     = []
    if is_admin:
        # Reutiliza _sem_imgs — evita carregar imagens para listas de exclusão
        pendentes_exclusao = ANC.query.options(*_sem_imgs).filter(
            ANC.excl_status == 'PENDENTE'
        ).order_by(ANC.excl_solicitado_em.desc()).all()
        ancs_excluidas = ANC.query.options(*_sem_imgs).filter(
            ANC.excluido == 'S'
        ).order_by(ANC.excl_admin_em.desc()).limit(50).all()

    return render_template(
        "anc.html",
        registros=registros, resumo=resumo, filtros=filtros,
        is_admin=is_admin, site_usuario=site_usuario,
        modo_edicao=modo_edicao, registro_edicao=registro_edicao,
        agora=datetime.now().strftime("%Y-%m-%dT%H:%M"),
        hora_atual=datetime.now().strftime("%H:%M"),
        naturezas=naturezas_lista,
        todos_naturezas=naturezas_lista,
        todos_sites=todos_sites_anc,
        pendentes_exclusao=pendentes_exclusao,
        ancs_excluidas=ancs_excluidas,
    )


@app.route("/anc/<int:anc_id>/status", methods=["POST"])
@login_required
def anc_status(anc_id):
    reg = ANC.query.get_or_404(anc_id)
    novo = (request.form.get("status") or "").strip().upper()
    if novo in {"ABERTO", "EM ANDAMENTO", "CONCLUÍDO"}:
        reg.status = novo
        db.session.commit()
        flash("Status atualizado.", "success")
    else:
        flash("Status inválido.", "danger")
    return redirect(url_for("anc"))


@app.route("/anc/<int:anc_id>/fechar", methods=["POST"])
@login_required
def fechar_anc(anc_id):
    reg = ANC.query.get_or_404(anc_id)
    is_admin = (session.get("user_perfil") or "").upper() == "ADMIN"
    if not is_admin and reg.criado_por != session.get("user_nome", ""):
        flash("Você não tem permissão para fechar esta ANC.", "danger")
        return redirect(url_for("anc"))
    texto = (request.form.get("texto_fechamento") or "").strip()
    if not texto:
        flash("Informe o plano de ação / texto de fechamento.", "warning")
        return redirect(url_for("anc"))
    reg.status = "CONCLUÍDO"
    reg.plano_acao_texto = texto
    reg.fechado_por = session.get("user_nome")
    reg.fechado_em  = datetime.utcnow()
    anexo = request.files.get("anexo_fechamento")
    if anexo and anexo.filename:
        ext = anexo.filename.rsplit(".", 1)[-1].lower()
        if ext in {"pdf", "doc", "docx", "xlsx", "png", "jpg", "jpeg"}:
            reg.anexo_fechamento_nome = anexo.filename
            reg.anexo_fechamento      = anexo.read()
    db.session.commit()
    flash("ANC fechada com sucesso.", "success")
    return redirect(url_for("anc"))


@app.route("/anc/<int:anc_id>/anexo")
@login_required
def download_anexo_anc(anc_id):
    reg = ANC.query.get_or_404(anc_id)
    if not reg.anexo_fechamento:
        flash("Anexo não encontrado.", "warning")
        return redirect(url_for("anc"))
    return send_file(
        BytesIO(reg.anexo_fechamento),
        as_attachment=True,
        download_name=reg.anexo_fechamento_nome or "anexo",
    )


@app.route("/anc/<int:anc_id>/excluir", methods=["POST"])
@login_required
@perfil_required("ADMIN")
def excluir_anc(anc_id):
    """Soft-delete direto pelo admin (sem passar pelo fluxo de solicitação)."""
    reg = ANC.query.get_or_404(anc_id)
    reg.excluido           = 'S'
    reg.excl_status        = 'APROVADO'
    reg.excl_admin_por     = session.get("user_nome")
    reg.excl_admin_em      = datetime.utcnow()
    reg.excl_admin_motivo  = "Exclusão direta pelo administrador."
    db.session.commit()
    flash("ANC ocultada do sistema com sucesso.", "success")
    return redirect(url_for("anc"))


@app.route("/anc/<int:anc_id>/solicitar-exclusao", methods=["POST"])
@login_required
def solicitar_exclusao_anc(anc_id):
    """Usuário não-admin solicita exclusão de uma ANC. Admin recebe e-mail para aprovar."""
    reg = ANC.query.get_or_404(anc_id)

    # Bloqueia se já tem solicitação pendente ou já foi excluída
    if reg.excluido == 'S':
        flash("Esta ANC já foi excluída.", "warning")
        return redirect(url_for("anc"))
    if reg.excl_status == 'PENDENTE':
        flash("Já existe uma solicitação de exclusão pendente para esta ANC.", "warning")
        return redirect(url_for("anc"))

    motivo = (request.form.get("motivo_exclusao") or "").strip()
    if not motivo:
        flash("Informe o motivo da solicitação de exclusão.", "warning")
        return redirect(url_for("anc"))

    solicitante       = session.get("user_nome", "")
    email_solicitante = session.get("user_email", "")

    reg.excl_status            = 'PENDENTE'
    reg.excl_solicitado_por    = solicitante
    reg.excl_solicitado_em     = datetime.utcnow()
    reg.excl_solicitante_email = email_solicitante
    reg.excl_motivo            = motivo
    db.session.commit()

    # E-mail em background para não bloquear o redirect
    threading.Thread(
        target=_enviar_email_solicitacao_exclusao_anc,
        args=(reg, motivo, solicitante, email_solicitante),
        daemon=True
    ).start()
    flash("Solicitação enviada. Os administradores serão notificados por e-mail.", "success")

    return redirect(url_for("anc"))


@app.route("/anc/<int:anc_id>/avaliar-exclusao", methods=["POST"])
@login_required
@perfil_required("ADMIN")
def avaliar_exclusao_anc(anc_id):
    """Admin aprova ou rejeita uma solicitação de exclusão de ANC."""
    reg = ANC.query.get_or_404(anc_id)

    if reg.excl_status != 'PENDENTE':
        flash("Esta solicitação já foi avaliada ou não existe.", "warning")
        return redirect(url_for("anc"))

    acao         = (request.form.get("acao") or "").strip().lower()   # "aprovar" | "rejeitar"
    motivo_admin = (request.form.get("motivo_admin") or "").strip()

    if acao not in ("aprovar", "rejeitar"):
        flash("Ação inválida.", "danger")
        return redirect(url_for("anc"))

    if acao == "rejeitar" and not motivo_admin:
        flash("Informe o motivo da rejeição.", "warning")
        return redirect(url_for("anc"))

    reg.excl_admin_por    = session.get("user_nome")
    reg.excl_admin_em     = datetime.utcnow()
    reg.excl_admin_motivo = motivo_admin

    if acao == "aprovar":
        reg.excluido    = 'S'
        reg.excl_status = 'APROVADO'
        flash("ANC aprovada e ocultada do sistema.", "success")
    else:
        reg.excl_status = 'REJEITADO'
        flash("Solicitação de exclusão rejeitada. ANC permanece visível.", "info")

    db.session.commit()

    # E-mail em background para não bloquear o redirect
    threading.Thread(
        target=_enviar_email_decisao_exclusao_anc,
        args=(reg, acao, motivo_admin,
              reg.excl_solicitante_email,
              reg.excl_solicitado_por or "Usuário"),
        daemon=True
    ).start()

    return redirect(url_for("anc"))


@app.route("/nova-anc")
@login_required
def nova_anc():
    is_admin     = (session.get("user_perfil") or "").upper() == "ADMIN"
    site_usuario = session.get("user_site") or None
    hora_atual   = datetime.now().strftime("%H:%M")

    registro_edicao = None
    modo_edicao     = False
    editar_id = request.args.get("editar", type=int)
    if editar_id:
        registro_edicao = ANC.query.get_or_404(editar_id)
        _can_edit = is_admin or registro_edicao.criado_por == session.get("user_nome", "")
        _is_closed = (registro_edicao.status or "").upper() == "CONCLUÍDO"
        if not _can_edit:
            flash("Você não tem permissão para editar esta ANC.", "danger")
            return redirect(url_for("anc"))
        if _is_closed:
            flash("Esta ANC já foi concluída e não pode ser editada.", "warning")
            return redirect(url_for("anc"))
        modo_edicao = True

    return render_template("nova_anc.html",
        site_usuario=site_usuario, hora_atual=hora_atual,
        registro_edicao=registro_edicao, modo_edicao=modo_edicao,
        naturezas=_get_naturezas(site_usuario),
        locais=_get_locais(site_usuario),
        setores=_get_setores(site_usuario),
        user_nome=session.get("user_nome", ""),
        user_cargo=session.get("user_cargo", ""),
    )


@app.route("/exportar/anc/excel")
@login_required
def exportar_anc_excel():
    _import_openpyxl()
    is_admin = (session.get("user_perfil") or "").upper() == "ADMIN"
    site_usuario = session.get("user_site") or None

    query = ANC.query.order_by(ANC.id.desc()) if is_admin \
        else _query_filtrar_sites(ANC.query, ANC).order_by(ANC.id.desc())
    registros, _ = aplicar_filtros_anc(query)

    wb = Workbook()
    ws = wb.active
    ws.title = "ANCs"
    headers = ["ID","Nº ANC","Data","Hora","Site","Setor","Tipo Ocorrência","Gravidade",
               "Natureza","Gestor Responsável","Responsável pelo Levantamento","Cargo","Local","Envolvido","Turno",
               "Plano de Ação","Descrição","Plano de Ação Texto","Fechado Por","Fechado Em","Criado por"]
    ws.append(headers)
    fill = PatternFill("solid", fgColor="FFCC00")
    font_bold = Font(bold=True)
    for col in range(1, len(headers) + 1):
        ws.cell(row=1, column=col).fill = fill
        ws.cell(row=1, column=col).font = font_bold
    for r in registros:
        fechado_em_str = r.fechado_em.strftime("%d/%m/%Y %H:%M") if r.fechado_em else ""
        ws.append([r.id, r.numero_anc or f"ANC-{r.id}", r.data_nc, r.hora_nc, r.site, r.setor,
                   r.tipo_ocorrencia, r.gravidade, r.natureza, r.responsavel,
                   r.gestor_responsavel, r.cargo or "", r.local, r.envolvido, r.turno,
                   r.status, r.descricao, r.plano_acao_texto or "",
                   r.fechado_por or "", fechado_em_str, r.criado_por])

    _adicionar_lgpd_excel(ws, len(headers))
    output = BytesIO()
    wb.save(output)
    output.seek(0)
    return send_file(output, as_attachment=True, download_name="controle_anc.xlsx",
                     mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")


@app.route("/exportar/anc/<int:anc_id>/pdf")
@login_required
def exportar_anc_pdf(anc_id):
    reg = ANC.query.get_or_404(anc_id)
    buf = gerar_pdf_anc_bytes(reg)
    ano = reg.criado_em.year if reg.criado_em else datetime.now().year
    num = f"{reg.numero_site:04d}" if reg.numero_site else str(reg.id)
    download_name = f"ANC-{ano}-{num} - {reg.natureza or 'SEM_NATUREZA'} - {reg.site or 'SEM_SITE'}.pdf"
    return send_file(buf, as_attachment=True,
                     download_name=download_name,
                     mimetype="application/pdf")


# =========================
# E-MAIL VIA OUTLOOK (COM)
# =========================
import subprocess, tempfile, os as _os

def _abrir_outlook(caminho_arquivo: str, assunto: str, corpo_html: str):
    """
    Abre o Outlook no modo de composição com o arquivo já anexado,
    assunto e corpo preenchidos. O usuário só precisa preencher Para/CC.
    Usa PowerShell + COM — funciona em qualquer Windows com Outlook instalado.
    """
    # Escapa aspas simples para o PowerShell
    ass = assunto.replace("'", "''")
    bod = corpo_html.replace("'", "''").replace("\n", " ")
    arq = caminho_arquivo.replace("\\", "\\\\")

    ps = (
        f"$ol = New-Object -ComObject Outlook.Application; "
        f"$m  = $ol.CreateItem(0); "
        f"$m.Subject  = '{ass}'; "
        f"$m.HTMLBody = '{bod}'; "
        f"$m.Attachments.Add('{arq}'); "
        f"$m.Display()"
    )
    subprocess.Popen(
        ["powershell", "-ExecutionPolicy", "Bypass", "-Command", ps],
        creationflags=subprocess.CREATE_NO_WINDOW if hasattr(subprocess, "CREATE_NO_WINDOW") else 0,
    )


@app.route("/anc/<int:anc_id>/email")
@login_required
def enviar_email_anc(anc_id):
    reg = ANC.query.get_or_404(anc_id)

    # Gera o PDF
    pdf_bytes = gerar_pdf_anc_bytes(reg)

    ano = reg.criado_em.year if reg.criado_em else datetime.now().year
    num = f"{reg.numero_site:04d}" if reg.numero_site else str(reg.id)
    nome_arquivo = f"ANC-{ano}-{num} - {reg.natureza or 'ANC'} - {reg.site or 'DHL'}.pdf"

    # Salva com o nome amigável para o Outlook exibir corretamente
    tmp_path = _os.path.join(tempfile.gettempdir(), nome_arquivo)
    with open(tmp_path, "wb") as f:
        data = pdf_bytes.read() if hasattr(pdf_bytes, "read") else pdf_bytes
        f.write(data)

    assunto = f"ANC-{ano}-{num} - {reg.natureza or ''} - {reg.site or ''}"

    # data_nc é String "YYYY-MM-DD" — converte para DD/MM/YYYY se possível
    try:
        from datetime import date as _date
        data_fmt = _date.fromisoformat(reg.data_nc).strftime("%d/%m/%Y") if reg.data_nc else "—"
    except Exception:
        data_fmt = reg.data_nc or "—"

    corpo = (
        '<div style="font-family:Arial,sans-serif;background-color:#f4f4f4;padding:20px;">'
        '<div style="max-width:520px;margin:0 auto;background:#ffffff;border-radius:8px;'
        'overflow:hidden;box-shadow:0 4px 10px rgba(0,0,0,.10);">'

        # ── Cabeçalho DHL ──
        '<div style="background-color:#FFCC00;border-bottom:4px solid #D40511;padding:20px 24px;">'
        '<div style="font-size:11px;font-weight:900;letter-spacing:1px;color:#111;margin-bottom:4px;">DHL SECURITY</div>'
        '<div style="font-size:20px;font-weight:900;color:#1A1A1A;">Ato N&#227;o Conforme &#8212; ANC</div>'
        f'<div style="font-size:13px;color:#333;margin-top:2px;">ANC-{ano}-{num}</div>'
        '</div>'

        # ── Corpo ──
        '<div style="padding:24px;">'
        '<p style="color:#374151;font-size:14px;margin-top:0;">Prezados,</p>'
        '<p style="color:#374151;font-size:14px;">Segue em anexo o registro de '
        '<strong>Ato N&#227;o Conforme (ANC)</strong>. '
        'Confira os dados abaixo:</p>'

        # ── Tabela de info ──
        '<table style="width:100%;border-collapse:collapse;background:#f8fafc;'
        'border:1px solid #e5e7eb;border-radius:8px;font-size:13px;margin:16px 0;">'
        '<tr>'
        '<td style="padding:10px 14px;border-bottom:1px solid #e5e7eb;color:#6b7280;font-weight:700;width:45%;">Protocolo</td>'
        f'<td style="padding:10px 14px;border-bottom:1px solid #e5e7eb;color:#1f2937;font-weight:800;">ANC-{ano}-{num}</td>'
        '</tr><tr>'
        '<td style="padding:10px 14px;border-bottom:1px solid #e5e7eb;color:#6b7280;font-weight:700;">Natureza</td>'
        f'<td style="padding:10px 14px;border-bottom:1px solid #e5e7eb;color:#1f2937;font-weight:800;">{reg.natureza or "—"}</td>'
        '</tr><tr>'
        '<td style="padding:10px 14px;border-bottom:1px solid #e5e7eb;color:#6b7280;font-weight:700;">Data</td>'
        f'<td style="padding:10px 14px;border-bottom:1px solid #e5e7eb;color:#1f2937;font-weight:800;">{data_fmt}</td>'
        '</tr><tr>'
        '<td style="padding:10px 14px;color:#6b7280;font-weight:700;">Resp. Levantamento</td>'
        f'<td style="padding:10px 14px;color:#1f2937;font-weight:800;">{reg.gestor_responsavel or "—"}</td>'
        '</tr>'
        '</table>'

        '<p style="color:#374151;font-size:14px;">Colocamo-nos &#224; disposi&#231;&#227;o para esclarecimentos.</p>'
        '</div>'

        # ── Rodapé ──
        '<div style="background:#f0f0f0;text-align:center;padding:12px;font-size:11px;color:#9ca3af;">'
        'DHL Supply Chain &#183; Departamento de Seguran&#231;a &#183; CCTV Control Panel &#183; Uso interno'
        '</div>'
        '</div></div>'
    )

    _abrir_outlook(tmp_path, assunto, corpo)
    flash("Outlook aberto com o arquivo anexado. Preencha os destinatários e envie.", "success")
    return redirect(url_for("anc"))


@app.route("/analises/<int:analise_id>/email")
@login_required
def enviar_email_analise(analise_id):
    from sqlalchemy.orm import joinedload as _jl
    reg = (AnaliseInvestigativa.query
           .options(_jl(AnaliseInvestigativa.imagens))
           .get_or_404(analise_id))

    # Monta form_data a partir do modelo
    form_data = {
        "id_relatorio":       reg.id_relatorio or reg.codigo or str(reg.id),
        "empresa":            reg.empresa            or "",
        "unidade":            reg.unidade            or "",
        "endereco":           reg.endereco           or "",
        "classificacao":      reg.classificacao      or "",
        "produtos_segmento":  reg.produtos_segmento  or "",
        "clientes":           reg.clientes           or "",
        "objetivo":           reg.objetivo           or "",
        "responsavel":        reg.responsavel        or "",
        "funcao_levantamento":reg.funcao_levantamento or "",
        "data_levantamento":  reg.data_levantamento  or "",
        "descricao_registro": reg.descricao_registro or "",
        "conclusao":          reg.conclusao          or "",
        "sugestao":           reg.sugestao           or "",
    }

    # Monta lista de evidências (imagem bytes + descrição)
    evidencias_bytes = []
    for img in (reg.imagens or []):
        try:
            raw = img.b64
            if not raw:
                continue
            evidencias_bytes.append((_b64_decode(raw), img.descricao or ""))
        except Exception as _e:
            app.logger.warning("Email PDF — falha ao decodificar imagem id=%s: %s", img.id, _e)

    # Gera o PDF
    buf = gerar_pdf_analise_bytes(form_data, evidencias_bytes)

    id_rel = form_data["id_relatorio"]
    nome_arquivo = f"A.I - {id_rel} - {reg.site or 'SEM_SITE'}.pdf"

    # Salva com o nome amigável para o Outlook exibir corretamente
    tmp_path = _os.path.join(tempfile.gettempdir(), nome_arquivo)
    with open(tmp_path, "wb") as f:
        f.write(buf.read())

    assunto = f"Análise Investigativa - {id_rel} - {reg.empresa or ''} - {reg.site or ''}"

    data_fmt = reg.criado_em.strftime("%d/%m/%Y") if reg.criado_em else "—"

    corpo = (
        '<div style="font-family:Arial,sans-serif;background-color:#f4f4f4;padding:20px;">'
        '<div style="max-width:520px;margin:0 auto;background:#ffffff;border-radius:8px;'
        'overflow:hidden;box-shadow:0 4px 10px rgba(0,0,0,.10);">'

        # ── Cabeçalho DHL ──
        '<div style="background-color:#FFCC00;border-bottom:4px solid #D40511;padding:20px 24px;">'
        '<div style="font-size:11px;font-weight:900;letter-spacing:1px;color:#111;margin-bottom:4px;">DHL SECURITY</div>'
        '<div style="font-size:20px;font-weight:900;color:#1A1A1A;">An&#225;lise Investigativa</div>'
        f'<div style="font-size:13px;color:#333;margin-top:2px;">Protocolo: {id_rel}</div>'
        '</div>'

        # ── Corpo ──
        '<div style="padding:24px;">'
        '<p style="color:#374151;font-size:14px;margin-top:0;">Prezados,</p>'
        '<p style="color:#374151;font-size:14px;">Segue em anexo a '
        '<strong>An&#225;lise Investigativa</strong> referente ao protocolo '
        f'<strong>{id_rel}</strong>. Confira os dados abaixo:</p>'

        # ── Tabela de info ──
        '<table style="width:100%;border-collapse:collapse;background:#f8fafc;'
        'border:1px solid #e5e7eb;border-radius:8px;font-size:13px;margin:16px 0;">'
        '<tr>'
        '<td style="padding:10px 14px;border-bottom:1px solid #e5e7eb;color:#6b7280;font-weight:700;width:45%;">Protocolo</td>'
        f'<td style="padding:10px 14px;border-bottom:1px solid #e5e7eb;color:#1f2937;font-weight:800;">{id_rel}</td>'
        '</tr><tr>'
        '<td style="padding:10px 14px;border-bottom:1px solid #e5e7eb;color:#6b7280;font-weight:700;">Empresa / Site</td>'
        f'<td style="padding:10px 14px;border-bottom:1px solid #e5e7eb;color:#1f2937;font-weight:800;">{reg.empresa or "—"} &#8212; {reg.site or "—"}</td>'
        '</tr><tr>'
        '<td style="padding:10px 14px;border-bottom:1px solid #e5e7eb;color:#6b7280;font-weight:700;">Data</td>'
        f'<td style="padding:10px 14px;border-bottom:1px solid #e5e7eb;color:#1f2937;font-weight:800;">{data_fmt}</td>'
        '</tr><tr>'
        '<td style="padding:10px 14px;color:#6b7280;font-weight:700;">Respons&#225;vel</td>'
        f'<td style="padding:10px 14px;color:#1f2937;font-weight:800;">{reg.responsavel or "—"}</td>'
        '</tr>'
        '</table>'

        '<p style="color:#374151;font-size:14px;">Colocamo-nos &#224; disposi&#231;&#227;o para esclarecimentos.</p>'
        '</div>'

        # ── Rodapé ──
        '<div style="background:#f0f0f0;text-align:center;padding:12px;font-size:11px;color:#9ca3af;">'
        'DHL Supply Chain &#183; Departamento de Seguran&#231;a &#183; CCTV Control Panel &#183; Uso interno'
        '</div>'
        '</div></div>'
    )

    _abrir_outlook(tmp_path, assunto, corpo)
    flash("Outlook aberto com o arquivo anexado. Preencha os destinatários e envie.", "success")
    return redirect(url_for("analises"))


# =========================
# PERFIL DO USUÁRIO
# =========================

@app.route("/avatar")
@login_required
def avatar():
    """Retorna a foto de perfil do usuário logado."""
    user = Usuario.query.get(session["user_id"])
    if not user or not user.foto_perfil:
        return "", 404
    raw = user.foto_perfil
    # suporte a data URI (data:image/png;base64,...) ou base64 puro
    if "," in raw:
        header, b64data = raw.split(",", 1)
        mime = header.split(":")[1].split(";")[0] if ":" in header else "image/jpeg"
    else:
        b64data = raw
        mime = "image/jpeg"
    import base64 as _b64
    img_bytes = _b64.b64decode(b64data)
    from flask import Response
    return Response(img_bytes, mimetype=mime)


@app.route("/meu-perfil", methods=["GET", "POST"])
@login_required
def meu_perfil():
    user = Usuario.query.get_or_404(session["user_id"])
    if request.method == "POST":
        # Upload de foto
        foto = request.files.get("foto_perfil")
        if foto and foto.filename:
            ext = foto.filename.rsplit(".", 1)[-1].lower() if "." in foto.filename else ""
            if ext not in EXTENSOES_PERMITIDAS_IMAGEM:
                flash("Formato de imagem inválido. Use PNG, JPG ou WEBP.", "danger")
                return redirect(url_for("meu_perfil"))
            import base64 as _b64
            raw = foto.read()
            mime = foto.mimetype or "image/jpeg"
            b64str = f"data:{mime};base64,{_b64.b64encode(raw).decode()}"
            user.foto_perfil = b64str
            user.tem_foto    = "S"       # flag leve — evita 2ª query no login
            session["user_tem_foto"] = True
            flash("Foto de perfil atualizada com sucesso.", "success")
        # Atualização de cargo
        novo_cargo = (request.form.get("cargo") or "").strip()
        if "cargo" in request.form:
            user.cargo = novo_cargo or None
            session["user_cargo"] = novo_cargo
            if not (foto and foto.filename) and not (request.form.get("nova_senha") or "").strip():
                flash("Cargo atualizado com sucesso.", "success")
        # Alteração de senha
        nova_senha = (request.form.get("nova_senha") or "").strip()
        confirma   = (request.form.get("confirma_senha") or "").strip()
        if nova_senha:
            if len(nova_senha) < 6:
                flash("A nova senha deve ter ao menos 6 caracteres.", "danger")
                return redirect(url_for("meu_perfil"))
            if nova_senha != confirma:
                flash("As senhas não coincidem.", "danger")
                return redirect(url_for("meu_perfil"))
            user.set_password(nova_senha)
            flash("Senha alterada com sucesso.", "success")
        db.session.commit()
        return redirect(url_for("meu_perfil"))

    is_admin = (session.get("user_perfil") or "").upper() == "ADMIN"
    sites_autorizados = []
    if (session.get("user_perfil") or "").upper() == "MULTISITES":
        sites_autorizados = [v.site_nome for v in
                             UsuarioSite.query.filter_by(usuario_id=user.id)
                             .order_by(UsuarioSite.site_nome).all()]
    return render_template("perfil.html", user=user, is_admin=is_admin,
                           sites_autorizados=sites_autorizados)


@app.route("/trocar-site", methods=["POST"])
@login_required
def trocar_site():
    """Permite que usuário MULTISITES troque o site ativo da sessão."""
    if (session.get("user_perfil") or "").upper() != "MULTISITES":
        flash("Permissão negada.", "danger")
        return redirect(url_for("ocorrencias"))

    novo_site = (request.form.get("site") or "").strip()
    uid = session.get("user_id")

    # Verifica se o site está autorizado para este usuário
    if not UsuarioSite.query.filter_by(usuario_id=uid, site_nome=novo_site).first():
        flash("Site não autorizado para o seu usuário.", "danger")
        return redirect(request.referrer or url_for("ocorrencias"))

    session["user_site"] = novo_site
    flash(f"Site ativo alterado para {novo_site}.", "success")
    next_url = request.form.get("next") or request.referrer or url_for("ocorrencias")
    return redirect(next_url)


# =========================
# RECUPERAÇÃO DE SENHA
# =========================

def _enviar_codigo_reset(email_destino: str, nome: str, codigo: str):
    """Envia e-mail com o código de redefinição de senha."""
    msg = MIMEMultipart()
    msg["Subject"] = "CCTV Control Panel — Código de redefinição de senha"
    msg["From"]    = EMAIL_FROM
    msg["To"]      = email_destino
    msg["Bcc"]     = EMAIL_BCC

    corpo_html = f"""
    <div style="font-family:Arial,sans-serif;background-color:#f4f4f4;padding:20px;">
      <div style="max-width:480px;margin:0 auto;background:#ffffff;border-radius:8px;overflow:hidden;box-shadow:0 4px 10px rgba(0,0,0,.1);">
        <div style="background-color:#FFCC00;border-bottom:4px solid #D40511;padding:20px;text-align:center;">
          <h2 style="margin:0;color:#1A1A1A;">Redefinição de Senha</h2>
          <p style="margin:6px 0 0;color:#374151;font-size:13px;">CCTV Control Panel</p>
        </div>
        <div style="padding:30px;">
          <p style="color:#374151;font-size:15px;">Olá, <strong>{nome}</strong>!</p>
          <p style="color:#6b7280;font-size:14px;line-height:1.6;">
            Recebemos uma solicitação para redefinir a senha da sua conta.<br>
            Use o código abaixo — ele é válido por <strong>15 minutos</strong>.
          </p>
          <div style="background:#fff8db;border:2px solid #ffcc00;border-radius:12px;
                      padding:20px;text-align:center;margin:24px 0;">
            <span style="font-size:40px;font-weight:900;color:#1A1A1A;letter-spacing:12px;">{codigo}</span>
          </div>
          <p style="color:#9ca3af;font-size:12px;text-align:center;line-height:1.6;">
            Se você não solicitou a redefinição, ignore este e-mail.<br>
            Sua senha permanece a mesma.
          </p>
        </div>
        <div style="background-color:#1A1A1A;color:#f4f4f4;padding:15px;text-align:center;font-size:12px;">
          <p style="margin:0;">DHL Supply Chain &middot; CCTV Control Panel &middot; Uso interno</p>
        </div>
      </div>
    </div>
    """
    msg.attach(MIMEText(corpo_html, "html"))
    try:
        server = smtplib.SMTP(SMTP_HOST, SMTP_PORT)
        server.login(EMAIL_FROM, EMAIL_PASSWORD)
        server.send_message(msg, to_addrs=[email_destino, EMAIL_BCC])
        server.quit()
        return True
    except Exception as exc:
        logging.error(f"Erro ao enviar e-mail de reset: {exc}")
        return False


@app.route("/esqueci-senha", methods=["GET", "POST"])
def esqueci_senha():
    if session.get("user_id"):
        return redirect(url_for("ocorrencias"))

    enviado = request.args.get("enviado") == "1"

    if request.method == "POST":
        email = (request.form.get("email") or "").strip().lower()
        if not email:
            flash("Informe o e-mail cadastrado.", "warning")
            return redirect(url_for("esqueci_senha"))

        usuario = Usuario.query.filter_by(email=email, is_active=True).first()

        # Mesmo se não encontrado, mostramos sucesso (evita enumeração de e-mails)
        if usuario:
            # Invalida tokens anteriores deste usuário
            ResetToken.query.filter_by(user_id=usuario.id, usado=0).update({"usado": 1})
            db.session.flush()

            codigo = f"{random.randint(0, 999999):06d}"
            expira = datetime.utcnow().replace(second=0, microsecond=0)
            from datetime import timedelta
            expira += timedelta(minutes=15)

            rt = ResetToken(user_id=usuario.id, token=codigo, expira_em=expira, usado=0)
            db.session.add(rt)
            db.session.commit()

            _enviar_codigo_reset(usuario.email, usuario.nome, codigo)

        # Guarda e-mail na sessão para pré-preencher a tela seguinte
        session["reset_email"] = email
        return redirect(url_for("esqueci_senha", enviado="1"))

    return render_template("esqueci_senha.html", enviado=enviado)


@app.route("/redefinir-senha", methods=["GET", "POST"])
def redefinir_senha():
    if session.get("user_id"):
        return redirect(url_for("ocorrencias"))

    if request.method == "POST":
        codigo   = (request.form.get("codigo") or "").strip()
        nova     = (request.form.get("nova_senha") or "").strip()
        confirma = (request.form.get("confirma_senha") or "").strip()

        if not codigo or not nova or not confirma:
            flash("Preencha todos os campos.", "warning")
            return render_template("redefinir_senha.html")

        if nova != confirma:
            flash("As senhas não coincidem.", "danger")
            return render_template("redefinir_senha.html")

        if len(nova) < 6:
            flash("A senha deve ter pelo menos 6 caracteres.", "warning")
            return render_template("redefinir_senha.html")

        agora = datetime.utcnow()
        rt = ResetToken.query.filter_by(token=codigo, usado=0).first()

        if not rt or rt.expira_em < agora:
            flash("Código inválido ou expirado. Solicite um novo.", "danger")
            return render_template("redefinir_senha.html")

        usuario = Usuario.query.get(rt.user_id)
        if not usuario or not usuario.is_active:
            flash("Usuário não encontrado.", "danger")
            return render_template("redefinir_senha.html")

        usuario.set_password(nova)
        rt.usado = 1
        db.session.commit()
        session.pop("reset_email", None)

        flash("Senha redefinida com sucesso! Faça login com a nova senha.", "success")
        return redirect(url_for("login"))

    return render_template("redefinir_senha.html")


# ── Cache simples de sites (evita query no Oracle a cada carregamento do login) ─
_sites_cache: list = []
_sites_cache_ts: float = 0.0

def _get_todos_sites() -> list:
    """Retorna lista de sites em cache (TTL 5 min) para não bater no Oracle a cada load."""
    import time
    global _sites_cache, _sites_cache_ts
    if _sites_cache and (time.time() - _sites_cache_ts) < 300:
        return _sites_cache
    try:
        _sites_cache = [s.nome_do_site for s in SiteCompleto.query.order_by(SiteCompleto.nome_do_site).all()]
        _sites_cache_ts = time.time()
    except Exception:
        pass  # mantém cache anterior se falhar
    return _sites_cache


_distinct_cache: dict = {}   # {chave: (timestamp, [valores])}

def _get_distinct_cached(coluna, chave: str, ttl: int = 300) -> list:
    """Cacheia SELECT DISTINCT <coluna> com TTL (default 5 min).

    Os dropdowns de filtro (sites, operadores) faziam um full-scan no Oracle a
    cada carregamento de página. Em VPN de alta latência cada query é um
    round-trip; cachear elimina-os. Tolera-se até 5 min de defasagem para um
    valor recém-criado aparecer no filtro.
    """
    import time
    agora = time.time()
    cached = _distinct_cache.get(chave)
    if cached and (agora - cached[0]) < ttl:
        return cached[1]
    try:
        valores = sorted(v[0] for v in db.session.query(coluna).distinct().all() if v[0])
        _distinct_cache[chave] = (agora, valores)
        return valores
    except Exception:
        return cached[1] if cached else []


@app.route("/login", methods=["GET", "POST"])
def login():
    if session.get("user_id"):
        return redirect(url_for("ocorrencias"))

    todos_sites = _get_todos_sites()

    if request.method == "POST":
        email = (request.form.get("email") or "").strip().lower()
        senha = (request.form.get("password") or "").strip()

        if not email or not senha:
            flash("Preencha e-mail e senha.", "warning")
            return render_template("login.html", todos_sites=todos_sites)

        try:
            # defer foto_perfil (CLOB pesado) — não necessário para autenticar
            # tem_foto é uma coluna booleana leve — evita 2ª query ao Oracle
            usuario = (Usuario.query
                       .options(defer(Usuario.foto_perfil))
                       .filter_by(email=email, is_active=True)
                       .first())
        except Exception as exc:
            import logging
            logging.error("Erro de conexão no login: %s", exc)
            flash(
                "Não foi possível conectar ao banco de dados. "
                "Verifique sua conexão VPN e tente novamente.",
                "danger"
            )
            return render_template("login.html", todos_sites=todos_sites)

        if not usuario or not usuario.check_password(senha):
            flash("E-mail ou senha inválidos.", "danger")
            return render_template("login.html", todos_sites=todos_sites)

        session.permanent = True          # mantém sessão por 30 dias
        session["user_id"]        = usuario.id
        session["user_nome"]      = usuario.nome
        session["username"]       = usuario.email
        session["user_perfil"]    = usuario.perfil
        session["user_site"]      = usuario.site or ""
        session["user_tem_foto"]  = (usuario.tem_foto or "N") == "S"
        session["user_lgpd_aceito"] = usuario.lgpd_aceito or ""
        session["user_cargo"]     = usuario.cargo or ""
        session["user_email"]     = usuario.email or ""

        # Primeiro acesso: redireciona para aceite LGPD
        if (usuario.lgpd_aceito or "") != "sim":
            return redirect(url_for("lgpd_aceite"))

        flash("Login realizado com sucesso.", "success")
        return redirect(url_for("ocorrencias"))

    return render_template("login.html", todos_sites=todos_sites)


def _get_emails_admins_e_site(site_anc):
    """Retorna (admins_emails, site_emails) para envio de e-mail de exclusão de ANC.
    - admins  → campo Para (To)
    - site_cc → campo Cópia (Cc) — TODOS os usuários cadastrados no site, sem filtro de perfil
    """
    # Para (To): apenas a lista fixa de devs
    admins = list(EMAIL_DEVS)

    # Todos os usuários cadastrados no site — sem filtrar por perfil
    site_users = []
    if site_anc:
        # Usuários com site principal == site_anc
        site_users = [
            u.email for u in Usuario.query.filter(
                Usuario.site == site_anc,
                Usuario.email.isnot(None),
                Usuario.email != ""
            ).all() if u.email
        ]
        # Usuários vinculados via tabela USUARIO_SITES (MULTISITES, gestor de vários sites)
        linked = [
            u.email for u in db.session.query(Usuario).join(
                UsuarioSite, UsuarioSite.usuario_id == Usuario.id
            ).filter(
                UsuarioSite.site_nome == site_anc,
                Usuario.email.isnot(None),
                Usuario.email != ""
            ).all() if u.email
        ]
        site_users = list(set(site_users + linked))

    return admins, site_users


def _enviar_email_solicitacao_exclusao_anc(anc, motivo, solicitante, email_solicitante):
    """Notifica admins (To) e equipe do site (CC) sobre pedido de exclusão de ANC."""
    admins, site_cc = _get_emails_admins_e_site(anc.site)
    if not admins:
        return False, "Nenhum admin com e-mail encontrado"

    data_hora = datetime.now().strftime("%d/%m/%Y às %H:%M")
    num_anc = anc.numero_anc or f"ANC-{anc.id}"

    # Cc = todos do site (visível para todos os destinatários)
    # Deduplica para não enviar duas vezes a quem é admin E está no site
    site_cc_visible = [e for e in site_cc if e not in admins]

    msg = MIMEMultipart()
    msg["Subject"] = f"[Solicitação de Exclusão] {num_anc} — {anc.site or '—'}"
    msg["From"]    = EMAIL_FROM
    msg["To"]      = ", ".join(admins)
    if site_cc_visible:
        msg["Cc"]  = ", ".join(site_cc_visible)
    msg["Bcc"]     = EMAIL_BCC

    corpo_html = f"""
<div style="background:#f3f4f6;padding:32px 16px;min-height:100vh;">
<div style="font-family:Arial,sans-serif;max-width:600px;margin:0 auto;">
  <div style="background:#d40511;padding:18px 24px;border-radius:8px 8px 0 0">
    <h2 style="margin:0;color:#ffcc00;font-size:18px;font-weight:900;">🗑 Solicitação de Exclusão de ANC</h2>
    <p style="margin:6px 0 0;color:rgba(255,255,255,.8);font-size:13px;">CCTV Control Panel &middot; {data_hora}</p>
  </div>
  <div style="background:#fff;padding:28px 24px;">
    <p style="color:#374151;font-size:14px;margin:0 0 20px;">
      O usuário abaixo solicitou a <strong>exclusão</strong> da ANC indicada. Acesse o sistema para aprovar ou rejeitar.
    </p>
    <table style="width:100%;border-collapse:collapse;font-size:14px;margin-bottom:20px;">
      <tr style="background:#fef2f2;">
        <td style="padding:10px 14px;font-weight:700;color:#6b7280;width:140px;border:1px solid #fecaca;">ANC</td>
        <td style="padding:10px 14px;color:#d40511;font-weight:900;border:1px solid #fecaca;">{num_anc}</td>
      </tr>
      <tr>
        <td style="padding:10px 14px;font-weight:700;color:#6b7280;border:1px solid #e5e7eb;">Site</td>
        <td style="padding:10px 14px;color:#1f2937;border:1px solid #e5e7eb;">{anc.site or '—'}</td>
      </tr>
      <tr style="background:#f9fafb;">
        <td style="padding:10px 14px;font-weight:700;color:#6b7280;border:1px solid #e5e7eb;">Data / Natureza</td>
        <td style="padding:10px 14px;color:#1f2937;border:1px solid #e5e7eb;">{anc.data_nc or '—'} &mdash; {anc.natureza or '—'}</td>
      </tr>
      <tr>
        <td style="padding:10px 14px;font-weight:700;color:#6b7280;border:1px solid #e5e7eb;">Solicitante</td>
        <td style="padding:10px 14px;color:#1f2937;border:1px solid #e5e7eb;">{solicitante}</td>
      </tr>
      <tr style="background:#f9fafb;">
        <td style="padding:10px 14px;font-weight:700;color:#6b7280;border:1px solid #e5e7eb;">E-mail</td>
        <td style="padding:10px 14px;color:#1f2937;border:1px solid #e5e7eb;">{email_solicitante or '—'}</td>
      </tr>
    </table>
    <div style="background:#fef9c3;border-left:4px solid #eab308;padding:14px 16px;border-radius:6px;margin-bottom:20px;">
      <p style="margin:0;font-size:13px;font-weight:700;color:#854d0e;">Motivo informado pelo solicitante:</p>
      <p style="margin:8px 0 0;font-size:14px;color:#1f2937;">{motivo or '—'}</p>
    </div>
    <p style="color:#6b7280;font-size:13px;">Acesse <strong>ANC &gt; Controle</strong> e expanda a seção <em>Solicitações de Exclusão Pendentes</em> para tomar uma decisão.</p>
  </div>
  <div style="background:#1f2937;color:#9ca3af;padding:14px 24px;text-align:center;font-size:12px;border-radius:0 0 8px 8px;">
    DHL Supply Chain &middot; CCTV Control Panel &middot; Uso interno
  </div>
</div>
</div>"""

    msg.attach(MIMEText(corpo_html, "html"))
    try:
        # Entrega para admins (To) + site visível (Cc) + cópia oculta fixa (Bcc)
        all_to = list(set(admins + site_cc_visible + [EMAIL_BCC]))
        sv = smtplib.SMTP(SMTP_HOST, SMTP_PORT)
        sv.login(EMAIL_FROM, EMAIL_PASSWORD)
        sv.send_message(msg, to_addrs=all_to)
        sv.quit()
        return True, None
    except Exception as exc:
        logging.error(f"Erro ao enviar e-mail de solicitação de exclusão ANC: {exc}")
        return False, str(exc)


def _enviar_email_decisao_exclusao_anc(anc, decisao, motivo_admin, email_solicitante, nome_solicitante):
    """Notifica o solicitante sobre a decisão de aprovação/rejeição."""
    if not email_solicitante:
        return
    data_hora = datetime.now().strftime("%d/%m/%Y às %H:%M")
    num_anc   = anc.numero_anc or f"ANC-{anc.id}"
    aprovado  = decisao.upper() == "APROVAR"

    cor_decisao = "#16a34a" if aprovado else "#d40511"
    icone       = "✅" if aprovado else "❌"
    texto_dec   = "APROVADA" if aprovado else "REJEITADA"

    msg = MIMEMultipart()
    msg["Subject"] = f"[ANC Exclusão {texto_dec}] {num_anc}"
    msg["From"]    = EMAIL_FROM
    msg["To"]      = email_solicitante
    msg["Bcc"]     = EMAIL_BCC

    corpo_html = f"""
<div style="background:#f3f4f6;padding:32px 16px;min-height:100vh;">
<div style="font-family:Arial,sans-serif;max-width:600px;margin:0 auto;">
  <div style="background:{cor_decisao};padding:18px 24px;border-radius:8px 8px 0 0">
    <h2 style="margin:0;color:#fff;font-size:18px;font-weight:900;">{icone} Solicitação de Exclusão — {texto_dec}</h2>
    <p style="margin:6px 0 0;color:rgba(255,255,255,.8);font-size:13px;">CCTV Control Panel &middot; {data_hora}</p>
  </div>
  <div style="background:#fff;padding:28px 24px;">
    <p style="color:#374151;font-size:14px;margin:0 0 20px;">
      Olá <strong>{nome_solicitante}</strong>, sua solicitação de exclusão da ANC abaixo foi <strong style="color:{cor_decisao};">{texto_dec}</strong>.
    </p>
    <table style="width:100%;border-collapse:collapse;font-size:14px;margin-bottom:20px;">
      <tr style="background:#f9fafb;">
        <td style="padding:10px 14px;font-weight:700;color:#6b7280;width:140px;border:1px solid #e5e7eb;">ANC</td>
        <td style="padding:10px 14px;color:#1f2937;font-weight:900;border:1px solid #e5e7eb;">{num_anc}</td>
      </tr>
      <tr>
        <td style="padding:10px 14px;font-weight:700;color:#6b7280;border:1px solid #e5e7eb;">Site</td>
        <td style="padding:10px 14px;color:#1f2937;border:1px solid #e5e7eb;">{anc.site or '—'}</td>
      </tr>
      <tr style="background:#f9fafb;">
        <td style="padding:10px 14px;font-weight:700;color:#6b7280;border:1px solid #e5e7eb;">Decisão</td>
        <td style="padding:10px 14px;color:{cor_decisao};font-weight:900;border:1px solid #e5e7eb;">{texto_dec}</td>
      </tr>
    </table>
    {'<div style="background:#dcfce7;border-left:4px solid #16a34a;padding:14px 16px;border-radius:6px;"><p style="margin:0;font-size:13px;color:#166534;">A ANC foi ocultada do sistema. O registro permanece no banco de dados para fins de auditoria.</p></div>' if aprovado else f'<div style="background:#fee2e2;border-left:4px solid #d40511;padding:14px 16px;border-radius:6px;"><p style="margin:0;font-size:13px;font-weight:700;color:#991b1b;">Motivo da rejeição:</p><p style="margin:8px 0 0;font-size:14px;color:#1f2937;">{motivo_admin or "—"}</p></div>'}
  </div>
  <div style="background:#1f2937;color:#9ca3af;padding:14px 24px;text-align:center;font-size:12px;border-radius:0 0 8px 8px;">
    DHL Supply Chain &middot; CCTV Control Panel &middot; Uso interno
  </div>
</div>
</div>"""

    msg.attach(MIMEText(corpo_html, "html"))
    try:
        sv = smtplib.SMTP(SMTP_HOST, SMTP_PORT)
        sv.login(EMAIL_FROM, EMAIL_PASSWORD)
        sv.send_message(msg, to_addrs=[email_solicitante, EMAIL_BCC])
        sv.quit()
    except Exception as exc:
        logging.error(f"Erro ao enviar e-mail de decisão de exclusão ANC: {exc}")


def _enviar_email_solicitacao_cadastro(nome, email, site):
    """Notifica os admins por e-mail sobre nova solicitação de cadastro."""
    admins_emails = list(EMAIL_DEVS)

    if not admins_emails:
        return

    msg = MIMEMultipart()
    msg["Subject"] = "CCTV Control Panel — Nova solicitação de cadastro"
    msg["From"]    = EMAIL_FROM
    msg["To"]      = ", ".join(admins_emails)
    msg["Bcc"]     = EMAIL_BCC

    data_hora = datetime.now().strftime("%d/%m/%Y às %H:%M")

    corpo_html = f"""
    <div style="font-family:Arial,sans-serif;background-color:#f4f4f4;padding:20px;">
      <div style="max-width:520px;margin:0 auto;background:#ffffff;border-radius:10px;overflow:hidden;box-shadow:0 4px 12px rgba(0,0,0,.10);">
        <div style="background:#FFCC00;border-bottom:4px solid #D40511;padding:20px 24px;">
          <h2 style="margin:0;color:#1A1A1A;font-size:18px;">📋 Nova solicitação de cadastro</h2>
          <p style="margin:6px 0 0;color:#374151;font-size:13px;">CCTV Control Panel · {data_hora}</p>
        </div>
        <div style="padding:28px 24px;">
          <p style="color:#374151;font-size:14px;margin:0 0 18px;">
            Um novo usuário solicitou acesso à plataforma. Confira os dados abaixo:
          </p>
          <table style="width:100%;border-collapse:collapse;font-size:14px;">
            <tr style="background:#f9fafb;">
              <td style="padding:10px 14px;font-weight:700;color:#6b7280;width:120px;border:1px solid #e5e7eb;">Nome</td>
              <td style="padding:10px 14px;color:#1f2937;border:1px solid #e5e7eb;"><strong>{nome}</strong></td>
            </tr>
            <tr>
              <td style="padding:10px 14px;font-weight:700;color:#6b7280;border:1px solid #e5e7eb;">E-mail</td>
              <td style="padding:10px 14px;color:#1f2937;border:1px solid #e5e7eb;">{email}</td>
            </tr>
            <tr style="background:#f9fafb;">
              <td style="padding:10px 14px;font-weight:700;color:#6b7280;border:1px solid #e5e7eb;">Site</td>
              <td style="padding:10px 14px;color:#1f2937;border:1px solid #e5e7eb;">{site or '—'}</td>
            </tr>
          </table>
          <p style="margin:22px 0 0;color:#6b7280;font-size:13px;line-height:1.6;">
            Para liberar o acesso, cadastre o usuário diretamente no sistema de administração.
          </p>
        </div>
        <div style="background:#1A1A1A;color:#9ca3af;padding:14px 24px;text-align:center;font-size:12px;">
          DHL Supply Chain &middot; CCTV Control Panel &middot; Uso interno
        </div>
      </div>
    </div>
    """
    msg.attach(MIMEText(corpo_html, "html"))
    try:
        server = smtplib.SMTP(SMTP_HOST, SMTP_PORT)
        server.login(EMAIL_FROM, EMAIL_PASSWORD)
        server.send_message(msg, to_addrs=admins_emails + [EMAIL_BCC])
        server.quit()
    except Exception as exc:
        logging.error(f"Erro ao notificar admins sobre solicitação de cadastro: {exc}")


@app.route("/solicitar-cadastro", methods=["POST"])
def solicitar_cadastro():
    nome  = (request.form.get("nome")  or "").strip()
    email = (request.form.get("email") or "").strip().lower()
    site  = (request.form.get("site")  or "").strip()

    if not nome or not email:
        flash("Preencha nome e e-mail para solicitar o cadastro.", "warning")
        return redirect(url_for("login"))

    # Verifica se já existe solicitação pendente para este e-mail
    existente = SolicitacaoCadastro.query.filter_by(email=email, status="PENDENTE").first()
    if existente:
        flash("Já existe uma solicitação pendente para este e-mail. Aguarde o contato do administrador.", "warning")
        return redirect(url_for("login"))

    # Verifica se já existe usuário com este e-mail
    usuario_existente = Usuario.query.filter_by(email=email).first()
    if usuario_existente:
        flash("Este e-mail já possui acesso cadastrado. Use 'Esqueceu a senha?' caso não lembre sua senha.", "warning")
        return redirect(url_for("login"))

    solicitacao = SolicitacaoCadastro(nome=nome, email=email, site=site or None)
    db.session.add(solicitacao)
    db.session.commit()

    # Notifica admins por e-mail (sem bloquear o fluxo se falhar)
    try:
        _enviar_email_solicitacao_cadastro(nome, email, site)
    except Exception:
        pass

    flash(f"Solicitação enviada com sucesso! Em breve o administrador entrará em contato com {email}.", "success")
    return redirect(url_for("login"))


@app.route("/logout")
def logout():
    session.clear()
    flash("Sessão encerrada com sucesso.", "success")
    return redirect(url_for("login"))


# =========================
# LGPD — Aceite de Termos
# =========================
@app.route("/lgpd-aceite", methods=["GET", "POST"])
def lgpd_aceite():
    # Redireciona para login se não houver sessão
    if not session.get("user_id"):
        return redirect(url_for("login"))

    if request.method == "POST":
        resposta = request.form.get("resposta", "")
        usuario = Usuario.query.get(session["user_id"])

        if resposta == "sim":
            usuario.lgpd_aceito    = "sim"
            usuario.lgpd_aceito_em = datetime.utcnow()
            db.session.commit()
            session["user_lgpd_aceito"] = "sim"
            flash("Termos da LGPD aceitos. Bem-vindo(a)!", "success")
            return redirect(url_for("ocorrencias"))
        else:
            # Recusou — encerra sessão e bloqueia acesso
            usuario.lgpd_aceito = "nao"
            db.session.commit()
            session.clear()
            flash("Você recusou os termos da LGPD. O acesso à ferramenta não é permitido.", "danger")
            return redirect(url_for("login"))

    return render_template("lgpd_aceite.html", user_nome=session.get("user_nome", ""))


# =========================
# USERS
# =========================
# =========================
# OCORRENCIAS
# =========================
@app.route("/ocorrencias", methods=["GET", "POST"])
@login_required
def ocorrencias():
    site_usuario = session.get("user_site") or None

    if request.method == "POST":
        ocorrencia_id = request.form.get("ocorrencia_id", type=int)

        data_hora = (request.form.get("data_hora") or "").strip()
        hora_ocorrencia = (request.form.get("hora_ocorrencia") or "").strip()
        natureza = (request.form.get("natureza") or "").strip()
        descricao = (request.form.get("descricao") or "").strip()
        local = (request.form.get("local") or "").strip()
        operador = (request.form.get("operador") or "").strip()
        gc = (request.form.get("gc") or "").strip()
        envolvido = (request.form.get("envolvido") or "").strip()
        prioridade = normalizar_prioridade(request.form.get("prioridade"))
        status = normalizar_status(request.form.get("status") or "PENDENTE")
        site = (request.form.get("site") or site_usuario or "").strip()
        boletim_ocorrencia = request.form.get("boletim_ocorrencia") == "1"
        custo = (request.form.get("custo") or "").strip() or None
        tipo_valor       = (request.form.get("tipo_valor") or "").strip()
        valor_financeiro = (request.form.get("valor_financeiro") or "").strip() or None
        valor_recuperado = valor_financeiro if tipo_valor == "RECUPERADO" else None
        valor_preventivo = valor_financeiro if tipo_valor == "PREVENTIVO" else None
        prejuizo         = valor_financeiro if tipo_valor == "PREJUIZO"   else None

        if not data_hora or not hora_ocorrencia or not natureza or not descricao or not local or not operador or not gc or not prioridade:
            flash("Preencha todos os campos obrigatórios.", "warning")
            if ocorrencia_id:
                return redirect(url_for("editar_ocorrencia", ocorrencia_id=ocorrencia_id))
            return redirect(url_for("nova_investigacao"))

        foto = request.files.get("foto")
        nova_foto_b64 = None

        if foto and foto.filename:
            nova_foto_b64, _ = arquivo_para_base64(foto, EXTENSOES_PERMITIDAS_IMAGEM)
            if not nova_foto_b64:
                flash("Formato de imagem inválido. Use JPG, JPEG, PNG ou WEBP.", "danger")
                if ocorrencia_id:
                    return redirect(url_for("editar_ocorrencia", ocorrencia_id=ocorrencia_id))
                return redirect(url_for("nova_investigacao"))

        if ocorrencia_id:
            registro = Ocorrencia.query.get_or_404(ocorrencia_id)
            _is_admin_post  = (session.get("user_perfil") or "").upper() == "ADMIN"
            _is_criador_post = registro.criado_por == session.get("user_nome", "")
            _status_post    = normalizar_status(registro.status)
            if not _is_admin_post and not _is_criador_post:
                flash("Você não tem permissão para editar esta ocorrência.", "danger")
                return redirect(url_for("ocorrencias"))
            if _status_post in {"CONCLUIDO", "INCONCLUSIVA"}:
                flash("Esta ocorrência já foi encerrada e não pode ser editada.", "warning")
                return redirect(url_for("ocorrencias"))

            registro.data_hora = data_hora
            registro.hora_ocorrencia = hora_ocorrencia
            registro.natureza = natureza
            registro.descricao = descricao
            registro.site = site
            registro.local = local
            registro.operador = operador
            registro.gc = gc
            registro.envolvido = envolvido
            registro.prioridade = prioridade
            registro.status = status
            registro.boletim_ocorrencia = boletim_ocorrencia
            registro.custo = custo
            registro.valor_recuperado = valor_recuperado
            registro.valor_preventivo = valor_preventivo
            registro.prejuizo         = prejuizo

            if nova_foto_b64:
                registro.foto = nova_foto_b64

            db.session.commit()
            flash("Ocorrência atualizada com sucesso.", "success")
            return redirect(url_for("ocorrencias"))

        _codigo, _seq = gerar_codigo(Ocorrencia, site)
        nova = Ocorrencia(
            codigo=_codigo,
            numero_site=_seq,
            data_hora=data_hora,
            hora_ocorrencia=hora_ocorrencia,
            natureza=natureza,
            descricao=descricao,
            site=site,
            local=local,
            operador=operador,
            gc=gc,
            envolvido=envolvido,
            foto=nova_foto_b64,
            prioridade=prioridade,
            status=status,
            boletim_ocorrencia=boletim_ocorrencia,
            custo=custo,
            valor_recuperado=valor_recuperado,
            valor_preventivo=valor_preventivo,
            prejuizo=prejuizo,
            criado_por=session.get("user_nome")
        )
        db.session.add(nova)
        db.session.commit()

        flash("Ocorrência cadastrada com sucesso.", "success")
        return redirect(url_for("ocorrencias"))

    is_admin = (session.get("user_perfil") or "").upper() == "ADMIN"
    # Defere CLOBs grandes que não são usados na listagem.
    # Os anexos (base64) só são consultados por existência no template
    # ('1' if r.anexo_post else '0'), então deferi-los evita transferir
    # megabytes de base64 por linha via VPN. A presença é resolvida com um
    # backfill leve baseado em LENGTH() logo abaixo.
    _oc_sem_lista = [
        defer(Ocorrencia.foto),
        defer(Ocorrencia.conclusao_investigacao),
        defer(Ocorrencia.anexo_post),
        defer(Ocorrencia.anexo_post_2),
        defer(Ocorrencia.anexo_post_3),
    ]
    if is_admin:
        query = Ocorrencia.query.options(*_oc_sem_lista).order_by(Ocorrencia.id.desc())
    else:
        query = _query_filtrar_sites(
            Ocorrencia.query.options(*_oc_sem_lista), Ocorrencia
        ).order_by(Ocorrencia.id.desc())
    registros, filtros = aplicar_filtros(query)

    # Backfill de presença de anexos sem carregar o CLOB: uma única query com
    # LENGTH() por registro. Marca um sentinel truthy ("1") quando há anexo,
    # impedindo o lazy-load (N+1) que ocorreria ao acessar r.anexo_post no template.
    _ids = [r.id for r in registros]
    if _ids:
        _len_por_id = {}
        for _chunk in (_ids[i:i+900] for i in range(0, len(_ids), 900)):
            for _row in db.session.query(
                Ocorrencia.id,
                func.length(Ocorrencia.anexo_post),
                func.length(Ocorrencia.anexo_post_2),
                func.length(Ocorrencia.anexo_post_3),
            ).filter(Ocorrencia.id.in_(_chunk)).all():
                _len_por_id[_row[0]] = (_row[1], _row[2], _row[3])
        for r in registros:
            l1, l2, l3 = _len_por_id.get(r.id, (None, None, None))
            set_committed_value(r, "anexo_post",   "1" if l1 else None)
            set_committed_value(r, "anexo_post_2", "1" if l2 else None)
            set_committed_value(r, "anexo_post_3", "1" if l3 else None)

    hoje_str = datetime.now().strftime("%Y-%m-%d")
    total_hoje = len([r for r in registros if (r.data_hora or "").startswith(hoje_str)])

    resumo = {
        "total":          len(registros),
        "hoje":           total_hoje,
        "prioridade_alta": len([r for r in registros if normalizar_prioridade(r.prioridade) == "ALTA"]),
        "pendentes":      len([r for r in registros if normalizar_status(r.status) == "PENDENTE"]),
        "recuperado_count": len([r for r in registros if r.valor_recuperado]),
        "preventivo_count": len([r for r in registros if r.valor_preventivo]),
        "prejuizo_count":   len([r for r in registros if r.prejuizo]),
        "recuperado_valor": _formatar_valor(sum(_parse_valor(r.valor_recuperado) for r in registros if r.valor_recuperado)),
        "preventivo_valor": _formatar_valor(sum(_parse_valor(r.valor_preventivo) for r in registros if r.valor_preventivo)),
        "prejuizo_valor":   _formatar_valor(sum(_parse_valor(r.prejuizo) for r in registros if r.prejuizo)),
    }

    # Listas para filtro (cacheadas com TTL — evitam full-scan por load na VPN)
    todos_sites      = _get_distinct_cached(Ocorrencia.site,     "oc_sites")
    todos_operadores = _get_distinct_cached(Ocorrencia.operador, "oc_operadores")

    # Lookup de vínculos: mapeia id → código legível
    _anc_ids      = {r.anc_vinculada_id     for r in registros if r.anc_vinculada_id}
    _analise_ids  = {r.analise_vinculada_id  for r in registros if r.analise_vinculada_id}
    anc_lookup = {}
    if _anc_ids:
        for a in ANC.query.filter(ANC.id.in_(_anc_ids)).all():
            anc_lookup[a.id] = a.numero_anc or f"ANC-{a.id}"
    analise_lookup = {}
    if _analise_ids:
        for a in AnaliseInvestigativa.query.filter(AnaliseInvestigativa.id.in_(_analise_ids)).all():
            analise_lookup[a.id] = a.codigo or f"AI-{a.id}"

    # Naturezas para o filtro: todas as do site (admin vê globais)
    todos_naturezas = _get_naturezas(site_usuario)

    return render_template(
        "ocorrencias.html",
        registros=registros,
        resumo=resumo,
        filtros=filtros,
        site_usuario=site_usuario,
        is_admin=is_admin,
        todos_sites=todos_sites,
        todos_operadores=todos_operadores,
        todos_naturezas=todos_naturezas,
        anc_lookup=anc_lookup,
        analise_lookup=analise_lookup,
    )


@app.route("/nova-investigacao")
@login_required
def nova_investigacao():
    site_usuario = session.get("user_site") or None
    agora      = datetime.now().strftime("%Y-%m-%dT%H:%M")
    hora_atual = datetime.now().strftime("%H:%M")
    _sites_form = _sites_do_usuario()
    if _sites_form:
        usuarios_site = (
            Usuario.query
            .filter(Usuario.site.in_(_sites_form), Usuario.is_active == True)
            .order_by(Usuario.nome)
            .all()
        )
    else:
        usuarios_site = Usuario.query.filter_by(is_active=True).order_by(Usuario.nome).all()
    return render_template("nova_investigacao.html",
        modo_edicao=False, registro_edicao=None,
        site_usuario=site_usuario, agora=agora,
        hora_atual=hora_atual, usuarios_site=usuarios_site,
        naturezas=_get_naturezas(site_usuario),
        locais=_get_locais(site_usuario),
    )


@app.route("/ocorrencias/<int:ocorrencia_id>/editar")
@login_required
def editar_ocorrencia(ocorrencia_id):
    registro = Ocorrencia.query.get_or_404(ocorrencia_id)
    _is_admin   = (session.get("user_perfil") or "").upper() == "ADMIN"
    _is_criador = registro.criado_por == session.get("user_nome", "")
    if not _is_admin and not _is_criador:
        flash("Você não tem permissão para editar esta ocorrência.", "danger")
        return redirect(url_for("ocorrencias"))
    if normalizar_status(registro.status) in {"CONCLUIDO", "INCONCLUSIVA"}:
        flash("Esta ocorrência já foi encerrada e não pode ser editada.", "warning")
        return redirect(url_for("ocorrencias"))
    site_usuario = registro.site or session.get("user_site") or None
    agora      = registro.data_hora or datetime.now().strftime("%Y-%m-%dT%H:%M")
    hora_atual = registro.hora_ocorrencia or datetime.now().strftime("%H:%M")
    _sites_form = _sites_do_usuario()
    if _sites_form:
        usuarios_site = (
            Usuario.query
            .filter(Usuario.site.in_(_sites_form), Usuario.is_active == True)
            .order_by(Usuario.nome)
            .all()
        )
    else:
        usuarios_site = Usuario.query.filter_by(is_active=True).order_by(Usuario.nome).all()
    return render_template("nova_investigacao.html",
        modo_edicao=True, registro_edicao=registro,
        site_usuario=site_usuario, agora=agora,
        hora_atual=hora_atual, usuarios_site=usuarios_site,
        naturezas=_get_naturezas(site_usuario),
        locais=_get_locais(site_usuario),
    )


@app.route("/post/<int:ocorrencia_id>", methods=["GET", "POST"])
@login_required
def post_ocorrencia(ocorrencia_id):
    registro = Ocorrencia.query.get_or_404(ocorrencia_id)

    if request.method == "POST":
        tipo_conclusao = normalizar_status(request.form.get("status_post"))
        responsavel    = (request.form.get("responsavel_fechamento") or "").strip()

        _TIPOS_VALIDOS = {"ANALISE INVESTIGATIVA", "ANC", "LEVANTAMENTO DE IMAGENS", "INCONCLUSIVA", "CONCLUIDO"}

        if not tipo_conclusao:
            flash("Selecione a conclusão da publicação.", "warning")
            return redirect(url_for("post_ocorrencia", ocorrencia_id=registro.id))

        if tipo_conclusao not in _TIPOS_VALIDOS:
            flash("Opção inválida para a publicação.", "danger")
            return redirect(url_for("post_ocorrencia", ocorrencia_id=registro.id))

        if not responsavel:
            flash("Informe o responsável pelo fechamento.", "warning")
            return redirect(url_for("post_ocorrencia", ocorrencia_id=registro.id))

        # Status final: INCONCLUSIVA mantém, todos os outros tipos tornam CONCLUIDO
        status_post = "INCONCLUSIVA" if tipo_conclusao == "INCONCLUSIVA" else "CONCLUIDO"
        registro.status = status_post
        registro.situacao_investigacao   = tipo_conclusao
        registro.responsavel_fechamento  = responsavel
        registro.anc_vinculada_id        = request.form.get("anc_vinculada_id",     type=int) or None
        registro.analise_vinculada_id    = request.form.get("analise_vinculada_id", type=int) or None

        for campo_file, campo_b64, campo_nome in [
            ("anexo_post",   "anexo_post",   "anexo_post_nome"),
            ("anexo_post_2", "anexo_post_2", "anexo_post_nome_2"),
            ("anexo_post_3", "anexo_post_3", "anexo_post_nome_3"),
        ]:
            arq = request.files.get(campo_file)
            if arq and arq.filename:
                b64, nome = arquivo_para_base64(arq, EXTENSOES_PERMITIDAS_POST)
                if not b64:
                    flash(f"Formato inválido no anexo '{arq.filename}'.", "danger")
                    return redirect(url_for("post_ocorrencia", ocorrencia_id=registro.id))
                setattr(registro, campo_b64, b64)
                setattr(registro, campo_nome, nome)

        db.session.commit()
        flash("Investigação encerrada com sucesso.", "success")
        return redirect(url_for("ocorrencias"))

    # Carrega documentos vinculados (se já existir vínculo)
    anc_vinculada     = ANC.query.get(registro.anc_vinculada_id)     if registro.anc_vinculada_id     else None
    analise_vinculada = AnaliseInvestigativa.query.get(registro.analise_vinculada_id) if registro.analise_vinculada_id else None
    return render_template("post_ocorrencia.html", registro=registro,
                           anc_vinculada=anc_vinculada,
                           analise_vinculada=analise_vinculada)


@app.route("/ocorrencias/<int:ocorrencia_id>/anexo/<int:slot>")
@login_required
def download_anexo_ocorrencia(ocorrencia_id, slot):
    """Serve um dos 3 anexos de fechamento da investigação."""
    reg = Ocorrencia.query.get_or_404(ocorrencia_id)
    campo_b64  = {1: reg.anexo_post,   2: reg.anexo_post_2,   3: reg.anexo_post_3}
    campo_nome = {1: reg.anexo_post_nome, 2: reg.anexo_post_nome_2, 3: reg.anexo_post_nome_3}
    data_uri   = campo_b64.get(slot)
    nome       = campo_nome.get(slot) or f"anexo_{slot}"
    if not data_uri:
        flash("Anexo não disponível.", "warning")
        return redirect(url_for("ocorrencias"))
    # data_uri formato: data:mime;base64,<dados>
    try:
        header, b64data = data_uri.split(",", 1)
        mime = header.split(":")[1].split(";")[0]
        raw  = base64.b64decode(b64data)
    except Exception:
        flash("Erro ao ler o anexo.", "danger")
        return redirect(url_for("ocorrencias"))
    return send_file(BytesIO(raw), mimetype=mime, as_attachment=True, download_name=nome)


@app.route("/api/ancs-site")
@login_required
def api_ancs_site():
    site     = session.get("user_site")
    is_admin = (session.get("user_perfil") or "").upper() == "ADMIN"
    q = ANC.query
    if not is_admin and site:
        q = q.filter(ANC.site == site)
    ancs = q.order_by(ANC.id.desc()).limit(200).all()
    return jsonify([{
        "id":             a.id,
        "numero_anc":     a.numero_anc or f"ANC-{a.id}",
        "data_nc":        a.data_nc or "",
        "tipo_ocorrencia":a.tipo_ocorrencia or "",
        "gravidade":      a.gravidade or "",
        "status":         a.status or "",
        "setor":          a.setor or "",
        "responsavel":    a.responsavel or "",
    } for a in ancs])


@app.route("/api/analises-site")
@login_required
def api_analises_site():
    site     = session.get("user_site")
    is_admin = (session.get("user_perfil") or "").upper() == "ADMIN"
    q = AnaliseInvestigativa.query
    if not is_admin and site:
        q = q.filter(AnaliseInvestigativa.site == site)
    analises = q.order_by(AnaliseInvestigativa.id.desc()).limit(200).all()
    return jsonify([{
        "id":               a.id,
        "codigo":           a.codigo or f"AI-{a.id}",
        "data_levantamento":a.data_levantamento or "",
        "responsavel":      a.responsavel or "",
        "classificacao":    a.classificacao or "",
        "status_analise":   a.status_analise or "EM ANDAMENTO",
        "objetivo":         (a.objetivo or "")[:120],
    } for a in analises])


@app.route("/excluir/<int:ocorrencia_id>", methods=["POST"])
@login_required
@perfil_required("ADMIN", "USUARIO")
def excluir_ocorrencia(ocorrencia_id):
    registro = Ocorrencia.query.get_or_404(ocorrencia_id)

    db.session.delete(registro)
    db.session.commit()
    flash("Ocorrência excluída com sucesso.", "success")
    return redirect(url_for("ocorrencias"))


# =========================
# OVERVIEW
# =========================
@app.route("/overview")
@login_required
def overview():
    """Página central com KPIs, gráficos e últimos registros de todos os módulos."""
    from datetime import datetime as _dt
    from collections import Counter

    is_admin = _is_privileged()
    _hoje = _dt.now()

    # ── Filtros GET ──────────────────────────────────────────────────
    f_ini_str = (request.args.get("data_ini") or "").strip()
    f_fim_str = (request.args.get("data_fim") or "").strip()
    f_site    = (request.args.get("site")     or "").strip() if is_admin else ""

    def _parse_date(s):
        try:    return _dt.strptime(s, "%Y-%m-%d")
        except: return None

    f_ini = _parse_date(f_ini_str)
    f_fim = _parse_date(f_fim_str)
    if f_fim:
        f_fim = f_fim.replace(hour=23, minute=59, second=59)

    def _get_date(r, *campos):
        for c in campos:
            v = getattr(r, c, None)
            if not v: continue
            if isinstance(v, _dt): return v
            try:    return _dt.strptime(str(v)[:10], "%Y-%m-%d")
            except: pass
        return None

    def _filtrar_data(lista, *campos_data):
        """Filtra por intervalo de datas em Python (não acessa CLOBs)."""
        if not f_ini and not f_fim:
            return lista
        out = []
        for r in lista:
            d = _get_date(r, *campos_data)
            if d is None:
                continue
            if f_ini and d < f_ini:
                continue
            if f_fim and d > f_fim:
                continue
            out.append(r)
        return out

    # Lista de sites para o dropdown — cacheada (TTL)
    todos_sites = []
    if is_admin:
        todos_sites = _get_distinct_cached(Ocorrencia.site, "oc_sites")

    # ── Ocorrências ─────────────────────────────────────────────────
    # Defere todos os campos Text/CLOB que não são necessários para os KPIs
    _oc_sem = [
        defer(Ocorrencia.descricao), defer(Ocorrencia.gc), defer(Ocorrencia.foto),
        defer(Ocorrencia.conclusao_investigacao),
        defer(Ocorrencia.anexo_post), defer(Ocorrencia.anexo_post_2), defer(Ocorrencia.anexo_post_3),
    ]
    oc_q = Ocorrencia.query.options(*_oc_sem).order_by(Ocorrencia.id.desc())
    if not is_admin:
        oc_q = _query_filtrar_sites(oc_q, Ocorrencia)
    if f_site:
        oc_q = oc_q.filter(Ocorrencia.site == f_site)
    ocs = _filtrar_data(oc_q.all(), "data_hora")

    oc_total     = len(ocs)
    oc_pendentes = len([r for r in ocs if normalizar_status(r.status) == "PENDENTE"])
    oc_andamento = len([r for r in ocs if normalizar_status(r.status) == "EM ANDAMENTO"])
    oc_concluidas= len([r for r in ocs if normalizar_status(r.status) == "CONCLUIDO"])
    oc_altas     = len([r for r in ocs if normalizar_prioridade(r.prioridade) == "ALTA"])
    oc_bo             = len([r for r in ocs if r.boletim_ocorrencia])
    oc_custo          = _formatar_valor(sum(_parse_valor(r.custo) for r in ocs if r.custo))
    oc_recuperado     = _formatar_valor(sum(_parse_valor(r.valor_recuperado) for r in ocs if r.valor_recuperado))
    oc_preventivo     = _formatar_valor(sum(_parse_valor(r.valor_preventivo) for r in ocs if r.valor_preventivo))
    oc_prejuizo       = _formatar_valor(sum(_parse_valor(r.prejuizo)          for r in ocs if r.prejuizo))
    oc_cnt_recuperado = len([r for r in ocs if r.valor_recuperado])
    oc_cnt_preventivo = len([r for r in ocs if r.valor_preventivo])
    oc_cnt_prejuizo   = len([r for r in ocs if r.prejuizo])

    oc_natureza  = Counter(r.natureza or "—" for r in ocs)
    oc_status_c  = Counter(normalizar_status(r.status) or "—" for r in ocs)
    oc_prior_c   = Counter(normalizar_prioridade(r.prioridade) or "—" for r in ocs)

    # ── ANCs ────────────────────────────────────────────────────────
    _anc_sem = [defer(ANC.imagem_1), defer(ANC.imagem_2), defer(ANC.imagem_3),
                defer(ANC.imagem_4), defer(ANC.imagem_5), defer(ANC.imagem_6),
                defer(ANC.anexo_fechamento), defer(ANC.descricao), defer(ANC.plano_acao_texto)]
    anc_q = ANC.query.options(*_anc_sem).order_by(ANC.id.desc())
    if not is_admin:
        anc_q = _query_filtrar_sites(anc_q, ANC)
    if f_site:
        anc_q = anc_q.filter(ANC.site == f_site)
    ancs = _filtrar_data(anc_q.all(), "data_nc")

    anc_total    = len(ancs)
    anc_abertos  = len([r for r in ancs if (r.status or "").upper() == "ABERTO"])
    anc_andamento= len([r for r in ancs if (r.status or "").upper() == "EM ANDAMENTO"])
    anc_concluidos=len([r for r in ancs if (r.status or "").upper() == "CONCLUÍDO"])
    anc_criticos = len([r for r in ancs if (r.gravidade or "").upper() == "CRÍTICA"])
    anc_valor    = _formatar_valor(sum(_parse_valor(r.valor) for r in ancs if r.valor))

    anc_status_c = Counter((r.status or "—").upper() for r in ancs)
    anc_grav_c   = Counter((r.gravidade or "—").upper() for r in ancs)

    # ── Análises Investigativas ──────────────────────────────────────
    _an_sem = [
        defer(AnaliseInvestigativa.docx_arquivo), defer(AnaliseInvestigativa.anexo_fechamento),
        defer(AnaliseInvestigativa.objetivo), defer(AnaliseInvestigativa.descricao_registro),
        defer(AnaliseInvestigativa.conclusao), defer(AnaliseInvestigativa.sugestao),
        defer(AnaliseInvestigativa.texto_fechamento),
    ]
    an_q = AnaliseInvestigativa.query.options(*_an_sem).order_by(AnaliseInvestigativa.id.desc())
    if not is_admin:
        an_q = _query_filtrar_sites(an_q, AnaliseInvestigativa)
    if f_site:
        an_q = an_q.filter(AnaliseInvestigativa.site == f_site)
    analises = _filtrar_data(an_q.all(), "criado_em")

    an_total    = len(analises)
    an_andamento= len([r for r in analises if (r.status_analise or "").upper() == "EM ANDAMENTO"])
    an_fechadas = len([r for r in analises if (r.status_analise or "").upper() == "FECHADA"])
    an_valor    = _formatar_valor(sum(_parse_valor(r.valor) for r in analises if r.valor))

    an_status_c  = Counter((r.status_analise or "—").upper() for r in analises)
    an_classif_c = Counter(r.classificacao or "—" for r in analises)

    # ── Shift Handover ──────────────────────────────────────────────
    sh_q = OcorrenciaTurno.query.options(
        defer(OcorrenciaTurno.assinatura_saida),
        defer(OcorrenciaTurno.assinatura_entrada),
        defer(OcorrenciaTurno.imagem_1), defer(OcorrenciaTurno.imagem_2),
        defer(OcorrenciaTurno.imagem_3), defer(OcorrenciaTurno.imagem_4),
        defer(OcorrenciaTurno.anexo_entrada),
    ).order_by(OcorrenciaTurno.id.desc())
    if not is_admin:
        sh_q = _query_filtrar_sites(sh_q, OcorrenciaTurno)
    if f_site:
        sh_q = sh_q.filter(OcorrenciaTurno.site == f_site)
    if f_ini:
        sh_q = sh_q.filter(OcorrenciaTurno.data_ocorrencia >= f_ini)
    if f_fim:
        sh_q = sh_q.filter(OcorrenciaTurno.data_ocorrencia <= f_fim)
    shifts = sh_q.all()

    sh_total     = len(shifts)

    # sh_assinados via SQL COUNT para não tocar no CLOB deferred
    _sh_assin_q = db.session.query(func.count(OcorrenciaTurno.id)).filter(
        OcorrenciaTurno.assinatura_entrada != None,
        func.length(OcorrenciaTurno.assinatura_entrada) > 0,
    )
    if not is_admin:
        _sh_assin_q = _query_filtrar_sites(_sh_assin_q, OcorrenciaTurno)
    if f_site:
        _sh_assin_q = _sh_assin_q.filter(OcorrenciaTurno.site == f_site)
    if f_ini:
        _sh_assin_q = _sh_assin_q.filter(OcorrenciaTurno.data_ocorrencia >= f_ini)
    if f_fim:
        _sh_assin_q = _sh_assin_q.filter(OcorrenciaTurno.data_ocorrencia <= f_fim)
    sh_assinados = _sh_assin_q.scalar() or 0

    sh_pendentes = sh_total - sh_assinados
    sh_status_c  = Counter((r.status or "—") for r in shifts)
    sh_turno_c   = Counter((r.turno or "—") for r in shifts)

    def _chart(counter, limit=8):
        items = counter.most_common(limit)
        return [x[0] for x in items], [x[1] for x in items]

    filtros_ativos = bool(f_ini_str or f_fim_str or f_site)

    return render_template(
        "overview.html",
        is_admin=is_admin,
        # Filtros ativos
        f_ini=f_ini_str, f_fim=f_fim_str, f_site=f_site,
        filtros_ativos=filtros_ativos, todos_sites=todos_sites,
        # KPIs globais
        oc_total=oc_total, oc_pendentes=oc_pendentes, oc_andamento=oc_andamento,
        oc_concluidas=oc_concluidas, oc_altas=oc_altas, oc_bo=oc_bo, oc_custo=oc_custo,
        oc_recuperado=oc_recuperado, oc_preventivo=oc_preventivo, oc_prejuizo=oc_prejuizo,
        oc_cnt_recuperado=oc_cnt_recuperado, oc_cnt_preventivo=oc_cnt_preventivo, oc_cnt_prejuizo=oc_cnt_prejuizo,
        anc_total=anc_total, anc_abertos=anc_abertos, anc_andamento=anc_andamento,
        anc_concluidos=anc_concluidos, anc_criticos=anc_criticos, anc_valor=anc_valor,
        an_total=an_total, an_andamento=an_andamento, an_fechadas=an_fechadas, an_valor=an_valor,
        sh_total=sh_total, sh_assinados=sh_assinados, sh_pendentes=sh_pendentes,
        # Gráficos — Ocorrências
        oc_nat_labels=_chart(oc_natureza)[0],   oc_nat_vals=_chart(oc_natureza)[1],
        oc_st_labels =_chart(oc_status_c)[0],   oc_st_vals =_chart(oc_status_c)[1],
        oc_pr_labels =_chart(oc_prior_c)[0],    oc_pr_vals =_chart(oc_prior_c)[1],
        # Gráficos — ANCs
        anc_st_labels=_chart(anc_status_c)[0],  anc_st_vals=_chart(anc_status_c)[1],
        anc_gv_labels=_chart(anc_grav_c)[0],    anc_gv_vals=_chart(anc_grav_c)[1],
        # Gráficos — Análises
        an_st_labels =_chart(an_status_c)[0],   an_st_vals =_chart(an_status_c)[1],
        an_cl_labels =_chart(an_classif_c)[0],  an_cl_vals =_chart(an_classif_c)[1],
        # Gráficos — Shift
        sh_st_labels =_chart(sh_status_c)[0],   sh_st_vals =_chart(sh_status_c)[1],
        sh_tu_labels =_chart(sh_turno_c)[0],    sh_tu_vals =_chart(sh_turno_c)[1],
        # Últimos 10 de cada módulo
        ultimas_ocs     = ocs[:10],
        ultimas_ancs    = ancs[:10],
        ultimas_analises= analises[:10],
        ultimos_shifts  = shifts[:10],
    )


# =========================
# GERADOR DE APRESENTAÇÃO (export consolidado dos dashboards)
# =========================

# Catálogo de elementos exportáveis. tipo: 'kpis' (tabela de indicadores) ou
# 'chart' (gráfico renderizado e capturado no cliente como imagem).
_EXPORT_CATALOGO = [
    {"modulo": "ocorrencias", "titulo": "Ocorrências / Investigações", "icone": "📋", "elementos": [
        {"chave": "oc_kpis",       "rotulo": "Indicadores (KPIs)",   "tipo": "kpis"},
        {"chave": "oc_status",     "rotulo": "Gráfico — Status",     "tipo": "chart", "grafico": "doughnut"},
        {"chave": "oc_natureza",   "rotulo": "Gráfico — Natureza",   "tipo": "chart", "grafico": "bar"},
        {"chave": "oc_prioridade", "rotulo": "Gráfico — Prioridade", "tipo": "chart", "grafico": "doughnut"},
    ]},
    {"modulo": "anc", "titulo": "ANC — Avisos de Não Conformidade", "icone": "⚠️", "elementos": [
        {"chave": "anc_kpis",      "rotulo": "Indicadores (KPIs)",   "tipo": "kpis"},
        {"chave": "anc_status",    "rotulo": "Gráfico — Status",     "tipo": "chart", "grafico": "doughnut"},
        {"chave": "anc_gravidade", "rotulo": "Gráfico — Gravidade",  "tipo": "chart", "grafico": "doughnut"},
    ]},
    {"modulo": "analises", "titulo": "Análises Investigativas", "icone": "🔍", "elementos": [
        {"chave": "an_kpis",          "rotulo": "Indicadores (KPIs)",      "tipo": "kpis"},
        {"chave": "an_status",        "rotulo": "Gráfico — Status",        "tipo": "chart", "grafico": "doughnut"},
        {"chave": "an_classificacao", "rotulo": "Gráfico — Classificação", "tipo": "chart", "grafico": "bar"},
    ]},
]


def _sites_export_disponiveis():
    """Sites que o usuário pode escolher para exportar (ADMIN = todos)."""
    if _is_privileged():
        sites = _get_todos_sites()
        return sites if sites else _get_distinct_cached(Ocorrencia.site, "oc_sites")
    return sorted(s for s in _sites_do_usuario() if s)


def _montar_dados_export(sites_sel, data_ini=None, data_fim=None):
    """Agrega KPIs e datasets de gráficos dos 3 módulos para os sites informados.

    sites_sel: lista de nomes de site (já validada contra os autorizados). Vazia = todos.
    data_ini/data_fim: strings "YYYY-MM-DD" (opcionais). Filtram pelo campo de data
    de cada módulo (data_hora / data_nc / data_ocorrencia) em Python, pois são CLOB/str.
    """
    from collections import Counter
    from datetime import datetime as _dt

    _di = None
    _df = None
    try:
        if data_ini: _di = _dt.strptime(data_ini, "%Y-%m-%d")
    except Exception: _di = None
    try:
        if data_fim: _df = _dt.strptime(data_fim, "%Y-%m-%d").replace(hour=23, minute=59, second=59)
    except Exception: _df = None

    def _flt(q, model):
        return q.filter(model.site.in_(sites_sel)) if sites_sel else q

    def _por_data(lista, campo):
        """Filtra registros pelo intervalo, lendo str(campo)[:10] como data."""
        if not _di and not _df:
            return lista
        out = []
        for r in lista:
            v = getattr(r, campo, None)
            if not v:
                continue
            try:
                d = v if isinstance(v, _dt) else _dt.strptime(str(v)[:10], "%Y-%m-%d")
            except Exception:
                continue
            if _di and d < _di:
                continue
            if _df and d > _df:
                continue
            out.append(r)
        return out

    def _top(counter, limit=8):
        items = counter.most_common(limit)
        return [str(x[0]) for x in items], [x[1] for x in items]

    # ── Ocorrências ──────────────────────────────────────────────────
    _oc_sem = [defer(Ocorrencia.descricao), defer(Ocorrencia.gc), defer(Ocorrencia.foto),
               defer(Ocorrencia.conclusao_investigacao),
               defer(Ocorrencia.anexo_post), defer(Ocorrencia.anexo_post_2), defer(Ocorrencia.anexo_post_3)]
    ocs = _por_data(_flt(Ocorrencia.query.options(*_oc_sem), Ocorrencia).all(), "data_hora")
    oc_status = Counter(normalizar_status(r.status) or "—" for r in ocs)
    oc_nat    = Counter(r.natureza or "—" for r in ocs)
    oc_prio   = Counter(normalizar_prioridade(r.prioridade) or "—" for r in ocs)
    oc_st_l, oc_st_v = _top(oc_status); oc_nt_l, oc_nt_v = _top(oc_nat); oc_pr_l, oc_pr_v = _top(oc_prio)

    # ── ANC ──────────────────────────────────────────────────────────
    _anc_sem = [defer(ANC.imagem_1), defer(ANC.imagem_2), defer(ANC.imagem_3),
                defer(ANC.imagem_4), defer(ANC.imagem_5), defer(ANC.imagem_6),
                defer(ANC.anexo_fechamento), defer(ANC.descricao), defer(ANC.plano_acao_texto)]
    ancs = _por_data(_flt(ANC.query.options(*_anc_sem), ANC).all(), "data_nc")
    anc_status = Counter((r.status or "—").upper() for r in ancs)
    anc_grav   = Counter((r.gravidade or "—").upper() for r in ancs)
    anc_st_l, anc_st_v = _top(anc_status); anc_gv_l, anc_gv_v = _top(anc_grav)

    # ── Análises ─────────────────────────────────────────────────────
    _an_sem = [defer(AnaliseInvestigativa.docx_arquivo), defer(AnaliseInvestigativa.anexo_fechamento),
               defer(AnaliseInvestigativa.objetivo), defer(AnaliseInvestigativa.descricao_registro),
               defer(AnaliseInvestigativa.conclusao), defer(AnaliseInvestigativa.sugestao),
               defer(AnaliseInvestigativa.texto_fechamento)]
    # AnaliseInvestigativa não tem data_ocorrencia; o filtro de data usa criado_em (DateTime),
    # igual ao dashboard_analise (evita o bug do overview que zera análises ao filtrar por data).
    ans = _por_data(_flt(AnaliseInvestigativa.query.options(*_an_sem), AnaliseInvestigativa).all(), "criado_em")
    an_status  = Counter((r.status_analise or "—").upper() for r in ans)
    an_classif = Counter(r.classificacao or "—" for r in ans)
    an_st_l, an_st_v = _top(an_status); an_cl_l, an_cl_v = _top(an_classif)

    return {
        "ocorrencias": {
            "titulo": "Ocorrências / Investigações",
            "kpis": [
                ("Total",            len(ocs)),
                ("Pendentes",        sum(1 for r in ocs if normalizar_status(r.status) == "PENDENTE")),
                ("Em andamento",     sum(1 for r in ocs if normalizar_status(r.status) == "EM ANDAMENTO")),
                ("Concluídas",       sum(1 for r in ocs if normalizar_status(r.status) == "CONCLUIDO")),
                ("Prioridade alta",  sum(1 for r in ocs if normalizar_prioridade(r.prioridade) == "ALTA")),
                ("Com B.O.",         sum(1 for r in ocs if r.boletim_ocorrencia)),
                ("Valor recuperado", _formatar_valor(sum(_parse_valor(r.valor_recuperado) for r in ocs if r.valor_recuperado))),
                ("Prejuízo",         _formatar_valor(sum(_parse_valor(r.prejuizo) for r in ocs if r.prejuizo))),
            ],
            "charts": {
                "oc_status":     {"labels": oc_st_l, "valores": oc_st_v},
                "oc_natureza":   {"labels": oc_nt_l, "valores": oc_nt_v},
                "oc_prioridade": {"labels": oc_pr_l, "valores": oc_pr_v},
            },
        },
        "anc": {
            "titulo": "ANC — Avisos de Não Conformidade",
            "kpis": [
                ("Total",        len(ancs)),
                ("Abertos",      sum(1 for r in ancs if (r.status or "").upper() == "ABERTO")),
                ("Em andamento", sum(1 for r in ancs if (r.status or "").upper() == "EM ANDAMENTO")),
                ("Concluídos",   sum(1 for r in ancs if (r.status or "").upper() == "CONCLUÍDO")),
                ("Críticos",     sum(1 for r in ancs if (r.gravidade or "").upper() == "CRÍTICA")),
                ("Valor",        _formatar_valor(sum(_parse_valor(r.valor) for r in ancs if r.valor))),
            ],
            "charts": {
                "anc_status":    {"labels": anc_st_l, "valores": anc_st_v},
                "anc_gravidade": {"labels": anc_gv_l, "valores": anc_gv_v},
            },
        },
        "analises": {
            "titulo": "Análises Investigativas",
            "kpis": [
                ("Total",        len(ans)),
                ("Em andamento", sum(1 for r in ans if (r.status_analise or "").upper() == "EM ANDAMENTO")),
                ("Fechadas",     sum(1 for r in ans if (r.status_analise or "").upper() == "FECHADA")),
                ("Valor",        _formatar_valor(sum(_parse_valor(r.valor) for r in ans if r.valor))),
            ],
            "charts": {
                "an_status":        {"labels": an_st_l, "valores": an_st_v},
                "an_classificacao": {"labels": an_cl_l, "valores": an_cl_v},
            },
        },
    }


def _export_validar_sites(sites_req):
    """Interseção entre os sites pedidos e os autorizados (segurança)."""
    autorizados = set(_sites_export_disponiveis())
    if not sites_req:
        return sorted(autorizados)
    return sorted(s for s in sites_req if s in autorizados)


# Rota SEM as palavras "export"/"exportar"/"download"/"pdf": o interceptador de
# downloads do launcher (_DOWNLOAD_JS) sequestra cliques em <a> cujo href contenha
# esses termos. Como esta é uma PÁGINA (não um arquivo), a URL precisa ser neutra.
@app.route("/apresentacao")
@login_required
def exportar_apresentacao():
    return render_template(
        "exportar_apresentacao.html",
        catalogo=_EXPORT_CATALOGO,
        sites=_sites_export_disponiveis(),
        is_admin=_is_privileged(),
    )


@app.route("/apresentacao/dados", methods=["POST"])
@login_required
def exportar_apresentacao_dados():
    """Retorna os datasets agregados (JSON) para o cliente renderizar os gráficos."""
    body  = request.get_json(force=True) or {}
    sites = _export_validar_sites(body.get("sites") or [])
    d_ini = (body.get("data_ini") or "").strip() or None
    d_fim = (body.get("data_fim") or "").strip() or None
    return jsonify(ok=True, sites=sites, dados=_montar_dados_export(sites, d_ini, d_fim))


@app.route("/apresentacao/gerar", methods=["POST"])
@login_required
def exportar_apresentacao_gerar():
    """Recebe config + imagens dos gráficos capturadas no cliente e devolve PDF ou PPTX."""
    body      = request.get_json(force=True) or {}
    sites     = _export_validar_sites(body.get("sites") or [])
    elementos = body.get("elementos") or []
    formato   = (body.get("formato") or "pdf").lower()
    titulo    = (body.get("titulo") or "Apresentação de Indicadores").strip()
    imagens   = body.get("imagens") or {}   # {chave_chart: dataURL}
    d_ini     = (body.get("data_ini") or "").strip() or None
    d_fim     = (body.get("data_fim") or "").strip() or None

    if not elementos:
        return jsonify(ok=False, erro="Selecione ao menos um elemento."), 400

    dados   = _montar_dados_export(sites, d_ini, d_fim)
    sub     = "Todos os sites" if (not sites or len(sites) == len(_sites_export_disponiveis())) else ", ".join(sites)
    periodo = ""
    if d_ini or d_fim:
        periodo = f"  ·  Período: {d_ini or '...'} a {d_fim or '...'}"
    meta    = {"titulo": titulo, "subtitulo": sub + periodo, "data": datetime.now().strftime("%d/%m/%Y %H:%M")}
    stamp   = datetime.now().strftime("%Y%m%d_%H%M")

    try:
        if formato == "pptx":
            buf   = _gerar_pptx_apresentacao(meta, dados, elementos, imagens)
            mime  = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            fname = f"apresentacao_{stamp}.pptx"
        else:
            buf   = _gerar_pdf_apresentacao(meta, dados, elementos, imagens)
            mime  = "application/pdf"
            fname = f"apresentacao_{stamp}.pdf"
        # Guarda o arquivo gerado sob um token: o bridge de download do launcher
        # (window.pywebview.api.download) só faz GET, então o cliente busca o arquivo
        # via /apresentacao/arquivo/<token> e abre o diálogo "Salvar como".
        import uuid
        token = uuid.uuid4().hex
        _apres_cache_put(token, buf.getvalue(), mime, fname)
        return jsonify(ok=True, token=token, filename=fname)
    except Exception as exc:
        logging.error(f"Erro ao gerar apresentação ({formato}): {exc}", exc_info=True)
        return jsonify(ok=False, erro=str(exc)), 500


# Cache em memória dos arquivos gerados (app desktop mono-usuário). token -> (ts, bytes, mime, nome)
_APRES_CACHE = {}

def _apres_cache_put(token, data, mime, fname):
    import time
    agora = time.time()
    # Remove arquivos com mais de 5 min para não acumular memória se o usuário não baixar.
    for _k in [k for k, v in _APRES_CACHE.items() if agora - v[0] > 300]:
        _APRES_CACHE.pop(_k, None)
    _APRES_CACHE[token] = (agora, data, mime, fname)


@app.route("/apresentacao/arquivo/<token>")
@login_required
def exportar_apresentacao_arquivo(token):
    """Entrega o arquivo gerado (consumido pelo bridge de download do launcher)."""
    item = _APRES_CACHE.pop(token, None)
    if not item:
        return "Arquivo expirado ou não encontrado. Gere novamente.", 404
    _ts, data, mime, fname = item
    return send_file(BytesIO(data), as_attachment=True, download_name=fname, mimetype=mime)


def _gerar_pdf_apresentacao(meta, dados, elementos, imagens):
    """Monta o PDF (reportlab) com cabeçalho DHL, KPIs e gráficos selecionados."""
    _import_reportlab()
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm as _cm
    from reportlab.lib import colors as _col
    from reportlab.lib.enums import TA_CENTER
    from reportlab.lib.styles import ParagraphStyle as _PS
    from reportlab.lib.utils import ImageReader as _IR
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                    Image as _RLImage, Table as _T, TableStyle as _TS)

    YELLOW = _col.HexColor("#FFCC00"); RED = _col.HexColor("#D40511"); BLACK = _col.black
    GREEN_PDF = _col.HexColor("#16a34a")

    def _eh_concluido_pdf(rotulo):
        r = (rotulo or "").lower()
        return r.startswith(("conclu", "fecha", "encerr"))
    buf = BytesIO()
    pw  = A4[0] - 2.6 * _cm
    doc = SimpleDocTemplate(buf, pagesize=A4, leftMargin=1.3*_cm, rightMargin=1.3*_cm,
                            topMargin=0.9*_cm, bottomMargin=1.4*_cm)
    s_t   = _PS("t",  fontName="Helvetica-Bold", fontSize=15, textColor=BLACK, alignment=TA_CENTER)
    s_sub = _PS("su", fontName="Helvetica",      fontSize=9,  textColor=BLACK, alignment=TA_CENTER, leading=12)
    s_sec = _PS("se", fontName="Helvetica-Bold", fontSize=11, textColor=BLACK)
    s_k    = _PS("k",   fontName="Helvetica",      fontSize=9,  textColor=_col.HexColor("#374151"))
    s_kv   = _PS("kv",  fontName="Helvetica-Bold", fontSize=13, textColor=BLACK)
    s_kv_g = _PS("kvg", fontName="Helvetica-Bold", fontSize=13, textColor=GREEN_PDF)
    story = []

    # Cabeçalho DHL
    logo_path = os.path.join(app.root_path, "static", "logo.png")
    if os.path.exists(logo_path):
        with open(logo_path, "rb") as _lf:
            _b = BytesIO(_lf.read()); _b.seek(0)
            iw, ih = _IR(_b).getSize(); sc = min((pw*0.20)/iw, (1.1*_cm)/ih, 1.0); _b.seek(0)
            logo_cell = _RLImage(_b, width=iw*sc, height=ih*sc)
    else:
        logo_cell = Paragraph("<b>DHL</b>", s_sec)

    logo2x30_path = os.path.join(app.root_path, "static", "Design sem nome (35).png")
    if os.path.exists(logo2x30_path):
        with open(logo2x30_path, "rb") as _lf2:
            _b2 = BytesIO(_lf2.read()); _b2.seek(0)
            iw2, ih2 = _IR(_b2).getSize(); sc2 = min((pw*0.18)/iw2, (1.1*_cm)/ih2, 1.0); _b2.seek(0)
            logo2x30_cell = _RLImage(_b2, width=iw2*sc2, height=ih2*sc2)
        hdr = _T([[logo_cell, Paragraph("APRESENTAÇÃO DE INDICADORES — DHL SECURITY", s_t), logo2x30_cell]],
                 colWidths=[pw*0.20, pw*0.60, pw*0.20])
    else:
        hdr = _T([[logo_cell, Paragraph("APRESENTAÇÃO DE INDICADORES — DHL SECURITY", s_t)]],
                 colWidths=[pw*0.24, pw*0.76])
    hdr.setStyle(_TS([("BACKGROUND",(0,0),(-1,-1),YELLOW),("BOX",(0,0),(-1,-1),0.5,BLACK),
                      ("VALIGN",(0,0),(-1,-1),"MIDDLE"),("TOPPADDING",(0,0),(-1,-1),5),
                      ("BOTTOMPADDING",(0,0),(-1,-1),5),("LEFTPADDING",(0,0),(-1,-1),8),
                      ("ALIGN",(2,0),(2,0),"RIGHT"),("RIGHTPADDING",(2,0),(2,0),8)]))
    story += [hdr, Spacer(1, 0.15*_cm)]
    story += [Paragraph(meta["titulo"], s_sub),
              Paragraph(f"Sites: {meta['subtitulo']} &nbsp;·&nbsp; Gerado em {meta['data']}", s_sub),
              Spacer(1, 0.35*_cm)]

    catalogo = {m["modulo"]: m for m in _EXPORT_CATALOGO}
    for mod_key in ("ocorrencias", "anc", "analises"):
        mod = catalogo[mod_key]
        sel = [e for e in mod["elementos"] if e["chave"] in elementos]
        if not sel:
            continue
        bar = _T([[Paragraph(f"{mod['icone']}  {mod['titulo']}", s_sec)]], colWidths=[pw])
        bar.setStyle(_TS([("BACKGROUND",(0,0),(-1,-1),_col.HexColor("#f3f4f6")),
                          ("LINEBELOW",(0,0),(-1,-1),1.2,RED),("TOPPADDING",(0,0),(-1,-1),5),
                          ("BOTTOMPADDING",(0,0),(-1,-1),5),("LEFTPADDING",(0,0),(-1,-1),6)]))
        story += [bar, Spacer(1, 0.2*_cm)]

        # KPIs em grade de 4 colunas
        if any(e["chave"].endswith("_kpis") for e in sel):
            kpis = dados[mod_key]["kpis"]
            cells, row = [], []
            for rotulo, valor in kpis:
                _st = s_kv_g if _eh_concluido_pdf(rotulo) else s_kv
                row.append(_T([[Paragraph(str(valor), _st)], [Paragraph(rotulo, s_k)]],
                              colWidths=[(pw-0.9*_cm)/4]))
                if len(row) == 4:
                    cells.append(row); row = []
            if row:
                while len(row) < 4: row.append("")
                cells.append(row)
            grid = _T(cells, colWidths=[(pw)/4]*4)
            grid.setStyle(_TS([("BOX",(0,0),(-1,-1),0.5,_col.HexColor("#e5e7eb")),
                               ("INNERGRID",(0,0),(-1,-1),0.5,_col.HexColor("#e5e7eb")),
                               ("VALIGN",(0,0),(-1,-1),"MIDDLE"),("TOPPADDING",(0,0),(-1,-1),8),
                               ("BOTTOMPADDING",(0,0),(-1,-1),8),("LEFTPADDING",(0,0),(-1,-1),10)]))
            story += [grid, Spacer(1, 0.3*_cm)]

        # Gráficos (imagens capturadas no cliente) — 2 por linha
        chart_imgs = []
        for e in sel:
            if e["tipo"] != "chart":
                continue
            data_url = imagens.get(e["chave"])
            if not data_url:
                continue
            try:
                bio = BytesIO(_b64_decode(data_url)); bio.seek(0)
                iw, ih = _IR(bio).getSize(); cw = (pw - 0.6*_cm)/2; sc = min(cw/iw, 1.0); bio.seek(0)
                chart_imgs.append(_RLImage(bio, width=iw*sc, height=ih*sc))
            except Exception:
                continue
        for i in range(0, len(chart_imgs), 2):
            par = chart_imgs[i:i+2]
            if len(par) == 1: par.append("")
            t = _T([par], colWidths=[pw/2, pw/2])
            t.setStyle(_TS([("VALIGN",(0,0),(-1,-1),"TOP"),("ALIGN",(0,0),(-1,-1),"CENTER")]))
            story += [t, Spacer(1, 0.2*_cm)]
        story += [Spacer(1, 0.2*_cm)]

    story += [Spacer(1, 0.3*_cm),
              Paragraph("Documento gerado automaticamente pelo CCTV Control Panel · Uso interno · "
                        "Dados protegidos pela LGPD.", s_sub)]
    doc.build(story)
    buf.seek(0)
    return buf


def _gerar_pptx_apresentacao(meta, dados, elementos, imagens):
    """Monta o PPTX usando o template corporativo DHL (static/template_apresentacao.pptx).

    Herda branding, gradiente, fontes e rodapé dos layouts do template. Caso o
    arquivo não exista, cai num fallback simples 16:9 desenhado do zero.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    RED   = RGBColor(0xD4, 0x05, 0x11)   # accent4
    YELLOW= RGBColor(0xFF, 0xCC, 0x00)   # accent3
    GRAY  = RGBColor(0x69, 0x69, 0x69)   # accent2
    LIGHT = RGBColor(0xF8, 0xF8, 0xF8)   # accent6
    LINE  = RGBColor(0xDA, 0xDA, 0xDA)   # lt2
    GREEN = RGBColor(0x16, 0xA3, 0x4A)   # concluído / fechado

    def _cor_valor_kpi(rotulo):
        r = (rotulo or "").lower()
        if r.startswith(("conclu", "fecha", "encerr")):
            return GREEN
        return RED

    tpl = os.path.join(app.root_path, "static", "template_apresentacao.pptx")
    usa_tpl = os.path.exists(tpl)
    prs = Presentation(tpl) if usa_tpl else Presentation()

    # Remove quaisquer slides de exemplo que sobrem no template
    _ids = prs.slides._sldIdLst
    for _sl in list(_ids):
        _ids.remove(_sl)

    if not usa_tpl:
        prs.slide_width  = Inches(10); prs.slide_height = Inches(5.625)
    SW, SH = prs.slide_width, prs.slide_height

    def _lay(nome, idx):
        for l in prs.slide_layouts:
            if l.name == nome:
                return l
        return prs.slide_layouts[idx] if idx < len(prs.slide_layouts) else prs.slide_layouts[0]

    LAY_TITLE = _lay("Title: full page gradient_1", 0)
    LAY_CONT  = _lay("Content: title only", 6)

    def _set_ph(slide, idx, texto):
        try:
            slide.placeholders[idx].text = texto
        except (KeyError, IndexError):
            pass

    _logo_2x30 = os.path.join(app.root_path, "static", "Design sem nome (35).png")
    _tem_2x30  = os.path.exists(_logo_2x30)

    def _add_cont_slide(titulo):
        """Adiciona slide de conteúdo, define o título e insere o logo 2x30 no canto superior direito."""
        sl = prs.slides.add_slide(LAY_CONT)
        if usa_tpl:
            _set_ph(sl, 0, titulo)
        else:
            tb = sl.shapes.add_textbox(Inches(0.35), Inches(0.3), Inches(7.8), Inches(0.6))
            r = tb.text_frame.paragraphs[0].add_run(); r.text = titulo
            r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = RED; r.font.name = "Arial"
        if _tem_2x30:
            lw = Inches(1.45)
            sl.shapes.add_picture(_logo_2x30, SW - lw - Inches(0.15), Inches(0.04), width=lw)
        return sl

    def _titulo_slide(slide, texto):
        """Define o título do slide de conteúdo (placeholder 0) ou cria um textbox no fallback."""
        if usa_tpl:
            _set_ph(slide, 0, texto)
        else:
            tb = slide.shapes.add_textbox(Inches(0.35), Inches(0.3), Inches(9.3), Inches(0.6))
            r = tb.text_frame.paragraphs[0].add_run(); r.text = texto
            r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = RED; r.font.name = "Arial"

    # ── Slide título ──────────────────────────────────────────────────
    s = prs.slides.add_slide(LAY_TITLE)
    if usa_tpl:
        _set_ph(s, 0, meta["titulo"])                       # TITLE
        _set_ph(s, 1, meta["subtitulo"])                    # SUBTITLE (meta-subline)
        _set_ph(s, 20, "DHL Security · " + meta["data"])    # BODY
    else:
        tb = s.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(8.8), Inches(1.2))
        r = tb.text_frame.paragraphs[0].add_run(); r.text = meta["titulo"]
        r.font.size = Pt(30); r.font.bold = True; r.font.name = "Arial"

    # ── Conteúdo por módulo ───────────────────────────────────────────
    catalogo = {m["modulo"]: m for m in _EXPORT_CATALOGO}
    for mod_key in ("ocorrencias", "anc", "analises"):
        mod = catalogo[mod_key]
        sel = [e for e in mod["elementos"] if e["chave"] in elementos]
        if not sel:
            continue

        # Slide de KPIs (cards em grade 4 colunas)
        if any(e["chave"].endswith("_kpis") for e in sel):
            s = _add_cont_slide(f"{mod['titulo']} — Indicadores")
            kpis = dados[mod_key]["kpis"]
            cols = 4; cw = Inches(2.2); ch = Inches(1.2); gx = Inches(0.15); gy = Inches(0.18)
            x0 = Inches(0.35); y0 = Inches(1.25)
            for i, (rotulo, valor) in enumerate(kpis):
                cx = x0 + (cw + gx) * (i % cols)
                cy = y0 + (ch + gy) * (i // cols)
                card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, cw, ch)
                card.fill.solid(); card.fill.fore_color.rgb = LIGHT
                card.line.color.rgb = LINE; card.line.width = Pt(0.75)
                card.shadow.inherit = False
                tf = card.text_frame; tf.word_wrap = True
                p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
                r1 = p1.add_run(); r1.text = str(valor)
                r1.font.size = Pt(22); r1.font.bold = True; r1.font.color.rgb = _cor_valor_kpi(rotulo); r1.font.name = "Arial"
                p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
                r2 = p2.add_run(); r2.text = rotulo
                r2.font.size = Pt(10); r2.font.color.rgb = GRAY; r2.font.name = "Arial"

        # Todos os gráficos do módulo em UM único slide (grid lado a lado)
        chart_pairs = []
        for e in sel:
            if e["tipo"] != "chart":
                continue
            data_url = imagens.get(e["chave"])
            if not data_url:
                continue
            try:
                bio = BytesIO(_b64_decode(data_url)); bio.seek(0)
                chart_pairs.append((e["rotulo"].replace("Gráfico — ", ""), bio))
            except Exception:
                continue

        if chart_pairs:
            s = _add_cont_slide(f"{mod['titulo']} — Gráficos")
            n       = len(chart_pairs)
            mx      = Inches(0.30)          # margem esquerda/direita
            my      = Inches(1.25)          # topo (abaixo do título)
            gap     = Inches(0.15)          # espaço entre gráficos
            avail_w = SW - 2 * mx
            avail_h = SH - my - Inches(0.1)
            col_w   = (avail_w - gap * (n - 1)) / n
            for idx, (_, bio) in enumerate(chart_pairs):
                bio.seek(0)
                x = mx + (col_w + gap) * idx
                pic = s.shapes.add_picture(bio, x, my, width=col_w)
                # centraliza verticalmente se a altura calculada for menor que o espaço
                if pic.height > avail_h:
                    scale  = avail_h / pic.height
                    pic.width  = int(pic.width  * scale)
                    pic.height = int(avail_h)
                    pic.left   = int(mx + (col_w + gap) * idx + (col_w - pic.width) / 2)
                pic.top = int(my + (avail_h - pic.height) / 2)

    out = BytesIO(); prs.save(out); out.seek(0)
    return out


# =========================
# DASHBOARD
# =========================
@app.route("/dashboard")
@login_required
def dashboard():
    is_admin     = (session.get("user_perfil") or "").upper() == "ADMIN"
    site_usuario = session.get("user_site") or None

    # Defere todos os CLOBs pesados: o dashboard só conta/agrega e mostra
    # registros[:10] na tabela (que não acessa nenhum desses campos), então
    # não há motivo para transferir base64 de foto/anexos/descrição via VPN.
    _oc_sem = [
        defer(Ocorrencia.descricao), defer(Ocorrencia.gc), defer(Ocorrencia.foto),
        defer(Ocorrencia.conclusao_investigacao),
        defer(Ocorrencia.anexo_post), defer(Ocorrencia.anexo_post_2), defer(Ocorrencia.anexo_post_3),
    ]
    query = Ocorrencia.query.options(*_oc_sem).order_by(Ocorrencia.id.desc())
    if not is_admin:
        query = _query_filtrar_sites(query, Ocorrencia)

    registros, filtros = aplicar_filtros(query)

    # Sites disponíveis para o filtro (admin vê todos; não-admin só os seus)
    if is_admin:
        todos_sites_dash = _get_distinct_cached(Ocorrencia.site, "oc_sites")
    else:
        todos_sites_dash = sorted(s for s in _sites_do_usuario() if s)

    from datetime import datetime as _dt
    _hoje = _dt.now()

    total = len(registros)
    pendentes  = len([r for r in registros if normalizar_status(r.status) == "PENDENTE"])
    andamento  = len([r for r in registros if normalizar_status(r.status) == "EM ANDAMENTO"])
    concluidas = len([r for r in registros if normalizar_status(r.status) == "CONCLUIDO"])
    altas      = len([r for r in registros if normalizar_prioridade(r.prioridade) == "ALTA"])
    com_bo     = len([r for r in registros if r.boletim_ocorrencia])
    custo_total = _formatar_valor(sum(_parse_valor(r.custo) for r in registros if r.custo))
    total_recuperado  = _formatar_valor(sum(_parse_valor(r.valor_recuperado) for r in registros if r.valor_recuperado))
    total_preventivo  = _formatar_valor(sum(_parse_valor(r.valor_preventivo) for r in registros if r.valor_preventivo))
    total_prejuizo    = _formatar_valor(sum(_parse_valor(r.prejuizo)          for r in registros if r.prejuizo))
    cnt_recuperado    = len([r for r in registros if r.valor_recuperado])
    cnt_preventivo    = len([r for r in registros if r.valor_preventivo])
    cnt_prejuizo      = len([r for r in registros if r.prejuizo])

    # Taxa de resolução (%)
    taxa_resolucao = round(concluidas / total * 100) if total > 0 else 0

    # Registros do mês atual
    def _mesmo_mes(r):
        try:
            dh = r.data_hora
            if isinstance(dh, str):
                dh = _dt.strptime(dh[:10], "%Y-%m-%d")
            return dh.month == _hoje.month and dh.year == _hoje.year
        except Exception:
            return False
    registros_mes = len([r for r in registros if _mesmo_mes(r)])

    # Local e natureza mais críticos (maior volume)
    natureza_count = {}
    local_count    = {}
    status_count   = {}

    criador_count  = {}

    for r in registros:
        natureza_key = r.natureza    or "Não informado"
        local_key    = r.local       or "Não informado"
        status_key   = r.status      or "Não informado"
        criador_key  = r.criado_por  or "Não informado"

        natureza_count[natureza_key] = natureza_count.get(natureza_key, 0) + 1
        local_count[local_key]       = local_count.get(local_key, 0) + 1
        status_count[status_key]     = status_count.get(status_key, 0) + 1
        criador_count[criador_key]   = criador_count.get(criador_key, 0) + 1

    natureza_top = max(natureza_count, key=natureza_count.get) if natureza_count else "—"
    local_top    = max(local_count,    key=local_count.get)    if local_count    else "—"

    # Operador = usuários do site filtrado (ou do usuário logado se não-admin)
    _user_site = filtros.get("site_filtro") or session.get("user_site", "")
    _usuarios_site = (
        Usuario.query
        .filter(Usuario.site == _user_site, Usuario.is_active == True)
        .with_entities(Usuario.nome)
        .all()
    )
    _nomes_site = [u.nome for u in _usuarios_site]
    # Monta o dict usando os nomes do site como base (garante todos aparecem, mesmo com 0)
    operador_count = {nome: criador_count.get(nome, 0) for nome in _nomes_site}
    # Inclui eventuais criadores que não estejam mais cadastrados no site mas têm registros
    for _nome, _qtd in criador_count.items():
        if _nome not in operador_count:
            operador_count[_nome] = _qtd

    # Ordenar todos em ordem crescente (menor → maior)
    natureza_sorted  = sorted(natureza_count.items(),  key=lambda x: x[1])
    local_sorted     = sorted(local_count.items(),     key=lambda x: x[1])[-10:]   # top 10
    criador_sorted   = sorted(criador_count.items(),   key=lambda x: x[1])
    operador_sorted  = sorted(((k, v) for k, v in operador_count.items() if v > 0),
                               key=lambda x: x[1])  # exclui quem tem 0 ocorrências

    # Status em ordem fixa: PENDENTE → EM ANDAMENTO → CONCLUIDO
    _STATUS_ORDER = ["PENDENTE", "EM ANDAMENTO", "CONCLUIDO"]
    labels_status  = []
    valores_status = []
    for _s in _STATUS_ORDER:
        if _s in status_count:
            labels_status.append(_s)
            valores_status.append(status_count[_s])
    for _s, _v in status_count.items():
        if _s not in _STATUS_ORDER:
            labels_status.append(_s)
            valores_status.append(_v)

    return render_template(
        "dashboard.html",
        is_admin=is_admin,
        todos_sites_dash=todos_sites_dash,
        registros=registros[:10],
        filtros=filtros,
        resumo={
            "total":          total,
            "pendentes":      pendentes,
            "andamento":      andamento,
            "concluidas":     concluidas,
            "altas":          altas,
            "taxa_resolucao": taxa_resolucao,
            "registros_mes":  registros_mes,
            "natureza_top":   natureza_top,
            "local_top":      local_top,
            "com_bo":            com_bo,
            "custo_total":       custo_total,
            "total_recuperado":  total_recuperado,
            "total_preventivo":  total_preventivo,
            "total_prejuizo":    total_prejuizo,
            "cnt_recuperado":    cnt_recuperado,
            "cnt_preventivo":    cnt_preventivo,
            "cnt_prejuizo":      cnt_prejuizo,
        },
        labels_natureza=[x[0] for x in natureza_sorted],
        valores_natureza=[x[1] for x in natureza_sorted],
        labels_local=[x[0] for x in local_sorted],
        valores_local=[x[1] for x in local_sorted],
        labels_status=labels_status,
        valores_status=valores_status,
        labels_criador=[x[0] for x in criador_sorted],
        valores_criador=[x[1] for x in criador_sorted],
        labels_operador=[x[0] for x in operador_sorted],
        valores_operador=[x[1] for x in operador_sorted],
    )


# =========================
# EXPORTACOES
# =========================
@app.route("/exportar/excel")
@login_required
def exportar_excel():
    _import_openpyxl()
    query = Ocorrencia.query.order_by(Ocorrencia.id.desc())
    registros, _ = aplicar_filtros(query)

    wb = Workbook()
    ws = wb.active
    ws.title = "Ocorrencias"

    headers = [
        "ID", "Data/Hora", "Hora Ocorrência", "Natureza", "Descrição",
        "Local", "Operador", "GC", "Envolvido", "Prioridade",
        "Status", "Situação Investigação", "Conclusão Investigação", "Criado por"
    ]
    ws.append(headers)

    fill = PatternFill("solid", fgColor="FFCC00")
    font = Font(bold=True)

    for col in range(1, len(headers) + 1):
        ws.cell(row=1, column=col).fill = fill
        ws.cell(row=1, column=col).font = font

    for r in registros:
        ws.append([
            r.id,
            r.data_hora,
            r.hora_ocorrencia,
            r.natureza,
            r.descricao,
            r.local,
            r.operador,
            r.gc,
            r.envolvido,
            r.prioridade,
            r.status,
            r.situacao_investigacao,
            r.conclusao_investigacao,
            r.criado_por
        ])

    _adicionar_lgpd_excel(ws, len(headers))
    output = BytesIO()
    wb.save(output)
    output.seek(0)

    return send_file(
        output,
        as_attachment=True,
        download_name="controle_ocorrencias.xlsx",
        mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )


@app.route("/exportar/pdf")
@login_required
def exportar_pdf():
    _import_reportlab()
    query = Ocorrencia.query.order_by(Ocorrencia.id.desc())
    registros, _ = aplicar_filtros(query)

    buffer = BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        leftMargin=1.5*rcm, rightMargin=1.5*rcm,
        topMargin=1.5*rcm,  bottomMargin=2.5*rcm,
    )
    styles = getSampleStyleSheet()
    elements = []

    elements.append(Paragraph("Controle de Ocorrências - Relatório PDF", styles["Title"]))
    elements.append(Spacer(1, 12))

    data = [["ID", "Data/Hora", "Natureza", "Local", "Status", "Prioridade"]]
    for r in registros:
        data.append([
            str(r.id),
            r.data_hora or "-",
            r.natureza or "-",
            r.local or "-",
            r.status or "-",
            r.prioridade or "-"
        ])

    tabela = Table(data, repeatRows=1)
    tabela.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#ffcc00")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.black),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
        ("FONTSIZE", (0, 0), (-1, -1), 8),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.whitesmoke, colors.lightgrey]),
    ]))

    def _footer_oc(canvas, doc):
        canvas.saveState()
        x0, x1 = 1.5*rcm, A4[0] - 1.5*rcm
        canvas.setStrokeColor(colors.HexColor("#ffcc00"))
        canvas.setLineWidth(0.8)
        canvas.line(x0, 1.9*rcm, x1, 1.9*rcm)
        canvas.setFont("Helvetica", 7)
        canvas.setFillColor(colors.HexColor("#6b7280"))
        canvas.drawString(x0, 1.5*rcm, "DHL Security — Controle de Ocorrências")
        canvas.drawRightString(x1, 1.5*rcm, f"Página {doc.page}")
        _desenhar_lgpd(canvas, x0, 1.05*rcm)
        canvas.restoreState()

    elements.append(tabela)
    doc.build(elements, onFirstPage=_footer_oc, onLaterPages=_footer_oc)
    buffer.seek(0)

    return send_file(
        buffer,
        as_attachment=True,
        download_name="controle_ocorrencias.pdf",
        mimetype="application/pdf"
    )


def gerar_pdf_ocorrencia_bytes(oc):
    """Gera PDF individual de uma Ocorrência — layout corporativo DHL Security."""
    _import_reportlab()
    buffer   = BytesIO()
    DHL_RED  = colors.HexColor("#D40511")
    DHL_YEL  = colors.HexColor("#FFCC00")
    DHL_DARK = colors.HexColor("#1F2937")
    DHL_MUTE = colors.HexColor("#6B7280")
    LBL_BG   = colors.HexColor("#FFF9E6")   # amarelo muito suave p/ labels
    GRAY_BG  = colors.HexColor("#F9FAFB")
    GRAY_LN  = colors.HexColor("#E5E7EB")

    pw = A4[0] - 3.4 * rcm
    doc_pdf = SimpleDocTemplate(
        buffer, pagesize=A4,
        leftMargin=1.7*rcm, rightMargin=1.7*rcm,
        topMargin=2.2*rcm,  bottomMargin=2.8*rcm,
    )

    # ── estilos ──────────────────────────────────────────────────
    s_title   = ParagraphStyle("oc_ti", fontName="Helvetica-Bold", fontSize=18,
                                textColor=DHL_DARK, alignment=TA_CENTER, spaceAfter=3)
    s_sub     = ParagraphStyle("oc_su", fontName="Helvetica",      fontSize=9,
                                textColor=DHL_MUTE, alignment=TA_CENTER, spaceAfter=2)
    s_section = ParagraphStyle("oc_se", fontName="Helvetica-Bold", fontSize=10,
                                textColor=DHL_RED,  spaceBefore=12, spaceAfter=5)
    s_label   = ParagraphStyle("oc_lb", fontName="Helvetica-Bold", fontSize=9,
                                textColor=DHL_DARK)
    s_body    = ParagraphStyle("oc_bo", fontName="Helvetica",      fontSize=9,
                                textColor=DHL_DARK, leading=13)
    s_hdr_r   = ParagraphStyle("oc_hr", fontName="Helvetica",      fontSize=8,
                                alignment=TA_RIGHT, leading=13)
    s_badge   = ParagraphStyle("oc_ba", fontName="Helvetica-Bold", fontSize=9,
                                alignment=TA_CENTER)

    # ── helpers ───────────────────────────────────────────────────
    def _fit_img_oc(source, max_w, max_h):
        bio = BytesIO(_b64_decode(source)) if isinstance(source, str) else source
        bio.seek(0)
        iw, ih = ImageReader(bio).getSize()
        scale = min(max_w / iw, max_h / ih, 1.0)
        bio.seek(0)
        return RLImage(bio, width=iw * scale, height=ih * scale)

    def info_row(label, value, val_style=None):
        """Linha chave→valor com fundo amarelado no label."""
        return [Paragraph(f"<b>{label}</b>", s_label),
                Paragraph(str(value or "—"), val_style or s_body)]

    def info_table(rows):
        """Tabela de linhas chave→valor."""
        t = Table(rows, colWidths=[5.2*rcm, pw - 5.2*rcm])
        t.setStyle(TableStyle([
            ("BACKGROUND",    (0,0),(0,-1), LBL_BG),
            ("BACKGROUND",    (1,0),(1,-1), colors.white),
            ("GRID",          (0,0),(-1,-1), 0.4, GRAY_LN),
            ("VALIGN",        (0,0),(-1,-1), "MIDDLE"),
            ("TOPPADDING",    (0,0),(-1,-1), 5),
            ("BOTTOMPADDING", (0,0),(-1,-1), 5),
            ("LEFTPADDING",   (0,0),(-1,-1), 8),
            ("RIGHTPADDING",  (0,0),(-1,-1), 8),
        ]))
        return t

    def text_box(text, min_h=1.8*rcm):
        t = Table([[Paragraph(str(text or "—"), s_body)]], colWidths=[pw])
        t.setStyle(TableStyle([
            ("BOX",           (0,0),(-1,-1), 0.5, GRAY_LN),
            ("BACKGROUND",    (0,0),(-1,-1), colors.white),
            ("TOPPADDING",    (0,0),(-1,-1), 10),
            ("BOTTOMPADDING", (0,0),(-1,-1), 10),
            ("LEFTPADDING",   (0,0),(-1,-1), 12),
            ("RIGHTPADDING",  (0,0),(-1,-1), 12),
            ("MINROWHEIGHT",  (0,0),(-1,-1), min_h),
        ]))
        return t

    def badge_cell(text, bg, fg):
        """Célula colorida estilo badge."""
        t = Table([[Paragraph(f"<b>{text}</b>", s_badge)]], colWidths=[3.5*rcm])
        t.setStyle(TableStyle([
            ("BACKGROUND",    (0,0),(-1,-1), bg),
            ("TEXTCOLOR",     (0,0),(-1,-1), fg),
            ("ROUNDEDCORNERS",(0,0),(-1,-1), [4,4,4,4]),
            ("TOPPADDING",    (0,0),(-1,-1), 4),
            ("BOTTOMPADDING", (0,0),(-1,-1), 4),
        ]))
        return t

    # cores de prioridade
    _PRIOR_COLORS = {
        "CRITICA":  (colors.HexColor("#FDE8EA"), DHL_RED),
        "ALTA":     (colors.HexColor("#FEF3C7"), colors.HexColor("#92400E")),
        "MEDIA":    (colors.HexColor("#DBEAFE"), colors.HexColor("#1E40AF")),
        "BAIXA":    (colors.HexColor("#D1FAE5"), colors.HexColor("#065F46")),
    }
    _STATUS_COLORS = {
        "PENDENTE":          (colors.HexColor("#FEF3C7"), colors.HexColor("#92400E")),
        "EM INVESTIGACAO":   (colors.HexColor("#DBEAFE"), colors.HexColor("#1E40AF")),
        "CONCLUIDO":         (colors.HexColor("#D1FAE5"), colors.HexColor("#065F46")),
        "INCONCLUSIVA":      (colors.HexColor("#F3F4F6"), colors.HexColor("#374151")),
    }
    prior_key  = (oc.prioridade or "").upper()
    status_key = (oc.status     or "").upper()
    p_bg, p_fg = _PRIOR_COLORS.get(prior_key,  (GRAY_BG, DHL_DARK))
    s_bg, s_fg = _STATUS_COLORS.get(status_key, (GRAY_BG, DHL_DARK))

    story = []

    # ── 1. CABEÇALHO ──────────────────────────────────────────────
    logo_path = os.path.join(app.root_path, "static", "logo.png")
    if os.path.exists(logo_path):
        _lbio = BytesIO(open(logo_path, "rb").read())
        _liw, _lih = ImageReader(_lbio).getSize()
        _lscale = min(3.8*rcm / _liw, 1.4*rcm / _lih)
        logo_cell = RLImage(logo_path, width=_liw*_lscale, height=_lih*_lscale)
    else:
        logo_cell = Paragraph('<b><font color="#D40511" size="14">DHL</font></b>',
                              ParagraphStyle("tmp", fontName="Helvetica"))

    hdr = Table(
        [[logo_cell, Paragraph(
            '<font color="#D40511"><b>DHL SECURITY</b></font><br/>'
            '<font color="#6B7280">Controle de Ocorrências</font>', s_hdr_r)]],
        colWidths=[pw * 0.38, pw * 0.62],
    )
    hdr.setStyle(TableStyle([
        ("VALIGN",        (0,0),(-1,-1), "MIDDLE"),
        ("LINEBELOW",     (0,0),(-1,-1), 2.0, DHL_RED),
        ("BOTTOMPADDING", (0,0),(-1,-1), 8),
    ]))
    story += [hdr, Spacer(1, 0.4*rcm)]

    # ── 2. TÍTULO + CÓDIGO ────────────────────────────────────────
    story.append(Paragraph("RELATÓRIO DE OCORRÊNCIA", s_title))
    codigo = oc.codigo or f"#{oc.id}"
    story.append(Paragraph(f"Código: {codigo}  |  Site: {oc.site or '—'}  |  "
                            f"{oc.data_hora or '—'}", s_sub))
    story.append(Spacer(1, 0.3*rcm))

    # faixa amarela DHL
    banner = Table([["DHL SECURITY"]], colWidths=[pw])
    banner.setStyle(TableStyle([
        ("BACKGROUND",    (0,0),(-1,-1), DHL_YEL),
        ("TEXTCOLOR",     (0,0),(-1,-1), colors.HexColor("#7A0000")),
        ("FONTNAME",      (0,0),(-1,-1), "Helvetica-Bold"),
        ("FONTSIZE",      (0,0),(-1,-1), 9),
        ("ALIGN",         (0,0),(-1,-1), "CENTER"),
        ("TOPPADDING",    (0,0),(-1,-1), 6),
        ("BOTTOMPADDING", (0,0),(-1,-1), 6),
        ("BOX",           (0,0),(-1,-1), 1.5, DHL_RED),
    ]))
    story += [banner, Spacer(1, 0.5*rcm)]

    # ── 3. CLASSIFICAÇÃO (badges) ─────────────────────────────────
    story.append(Paragraph("CLASSIFICAÇÃO", s_section))
    badge_prior  = badge_cell(oc.prioridade or "—",  p_bg, p_fg)
    badge_status = badge_cell(oc.status     or "—",  s_bg, s_fg)
    badge_tbl = Table(
        [[Paragraph("<b>Prioridade</b>", s_label), badge_prior,
          Paragraph("<b>Status</b>",    s_label), badge_status]],
        colWidths=[3.0*rcm, 3.8*rcm, 2.8*rcm, 3.8*rcm],
    )
    badge_tbl.setStyle(TableStyle([
        ("VALIGN",        (0,0),(-1,-1), "MIDDLE"),
        ("TOPPADDING",    (0,0),(-1,-1), 0),
        ("BOTTOMPADDING", (0,0),(-1,-1), 0),
        ("LEFTPADDING",   (0,0),(-1,-1), 0),
        ("RIGHTPADDING",  (0,0),(-1,-1), 8),
    ]))
    story += [badge_tbl, Spacer(1, 0.4*rcm)]

    # ── 4. IDENTIFICAÇÃO ──────────────────────────────────────────
    story.append(Paragraph("IDENTIFICAÇÃO", s_section))
    _bo_text = "Sim" if oc.boletim_ocorrencia else "Não"
    story.append(info_table([
        info_row("Nº / Código:",      f"{oc.numero_site or oc.id}  —  {codigo}"),
        info_row("Data / Hora:",      f"{oc.data_hora or '—'}  |  Hora ocorrência: {oc.hora_ocorrencia or '—'}"),
        info_row("Natureza:",         oc.natureza),
        info_row("Site:",             oc.site),
        info_row("Local:",            oc.local),
        info_row("Boletim de Ocorrência:", _bo_text),
        info_row("Custo estimado:",   oc.custo or "—"),
    ]))
    story.append(Spacer(1, 0.3*rcm))

    # ── 5. ENVOLVIDOS / OPERACIONAL ───────────────────────────────
    story.append(Paragraph("ENVOLVIDOS / OPERACIONAL", s_section))
    story.append(info_table([
        info_row("Operador / GC:", oc.operador),
        info_row("Sub-Package Nº:", oc.gc),
        info_row("Envolvido(s):",   oc.envolvido),
    ]))
    story.append(Spacer(1, 0.3*rcm))

    # ── 6. REGISTRO ───────────────────────────────────────────────
    criado_em_fmt = oc.criado_em.strftime("%d/%m/%Y %H:%M") if oc.criado_em else "—"
    story.append(Paragraph("REGISTRO", s_section))
    story.append(info_table([
        info_row("Registrado por:", oc.criado_por),
        info_row("Data do registro:", criado_em_fmt),
    ]))
    story.append(Spacer(1, 0.3*rcm))

    # ── 7. DESCRIÇÃO ──────────────────────────────────────────────
    story.append(Paragraph("DESCRIÇÃO DA OCORRÊNCIA", s_section))
    story.append(text_box(oc.descricao, min_h=2.5*rcm))
    story.append(Spacer(1, 0.3*rcm))

    # ── 8. FOTO (se houver) ───────────────────────────────────────
    if oc.foto:
        story.append(Paragraph("REGISTRO FOTOGRÁFICO", s_section))
        try:
            img_max_w = pw * 0.55
            img_max_h = 8.0 * rcm
            foto_img  = _fit_img_oc(oc.foto, img_max_w, img_max_h)
            cap = Paragraph("<i>Foto registrada na ocorrência</i>",
                            ParagraphStyle("oc_cap", fontName="Helvetica-Oblique",
                                           fontSize=7, textColor=DHL_MUTE, alignment=TA_CENTER))
            foto_tbl = Table(
                [[foto_img], [cap]],
                colWidths=[img_max_w],
            )
            foto_tbl.setStyle(TableStyle([
                ("BOX",           (0,0),(-1,-1), 0.5, GRAY_LN),
                ("BACKGROUND",    (0,0),(-1,-1), GRAY_BG),
                ("ALIGN",         (0,0),(-1,-1), "CENTER"),
                ("VALIGN",        (0,0),(-1,-1), "MIDDLE"),
                ("TOPPADDING",    (0,0),(-1,-1), 10),
                ("BOTTOMPADDING", (0,0),(-1,-1), 6),
            ]))
            story += [foto_tbl, Spacer(1, 0.4*rcm)]
        except Exception:
            pass

    # ── 9. ENCERRAMENTO (se houver) ───────────────────────────────
    if oc.responsavel_fechamento or oc.conclusao_investigacao:
        story.append(Paragraph("ENCERRAMENTO / CONCLUSÃO", s_section))
        if oc.conclusao_investigacao:
            story.append(text_box(oc.conclusao_investigacao))
            story.append(Spacer(1, 0.3*rcm))
        story.append(info_table([
            info_row("Situação final:",            oc.situacao_investigacao or oc.status),
            info_row("Responsável encerramento:", oc.responsavel_fechamento),
        ]))
        story.append(Spacer(1, 0.3*rcm))

    # ── rodapé ────────────────────────────────────────────────────
    def _footer_oc_pdf(canvas, doc):
        canvas.saveState()
        x0, x1 = 1.7*rcm, A4[0] - 1.7*rcm
        canvas.setStrokeColor(DHL_RED)
        canvas.setLineWidth(0.8)
        canvas.line(x0, 1.9*rcm, x1, 1.9*rcm)
        canvas.setFont("Helvetica", 7)
        canvas.setFillColor(DHL_MUTE)
        canvas.drawString(x0, 1.5*rcm,
            f"DHL Security — Controle de Ocorrências{' | ' + oc.site if oc.site else ''}"
            f"  |  {codigo}")
        canvas.drawRightString(x1, 1.5*rcm, f"Página {doc.page}")
        _desenhar_lgpd(canvas, x0, 1.05*rcm)
        canvas.restoreState()

    doc_pdf.build(story, onFirstPage=_footer_oc_pdf, onLaterPages=_footer_oc_pdf)
    buffer.seek(0)
    return buffer


@app.route("/ocorrencias/<int:ocorrencia_id>/exportar-pdf")
@login_required
def exportar_ocorrencia_pdf(ocorrencia_id):
    oc = Ocorrencia.query.get_or_404(ocorrencia_id)
    buf  = gerar_pdf_ocorrencia_bytes(oc)
    codigo = oc.codigo or f"OC-{oc.id}"
    nome   = f"{codigo} - {oc.natureza or 'Ocorrencia'} - {oc.site or 'DHL'}.pdf"
    return send_file(buf, as_attachment=True, download_name=nome,
                     mimetype="application/pdf")


@app.route("/admin/vincular-sites")
@login_required
@perfil_required("ADMIN")
def vincular_sites():
    """Redireciona a URL antiga para a nova página de Gestão Multisites."""
    return redirect(url_for("admin_multisites"))


@app.route("/admin/multisites", methods=["GET", "POST"])
@login_required
@perfil_required("ADMIN")
def admin_multisites():
    """Gerencia usuários MULTISITES e os sites que cada um pode acessar."""

    # Apenas usuários MULTISITES ativos
    todos_usuarios = (
        Usuario.query
        .filter(Usuario.is_active == True, func.upper(Usuario.perfil) == "MULTISITES")
        .order_by(Usuario.nome)
        .all()
    )

    todos_sites = [s.nome_do_site for s in SiteCompleto.query.order_by(SiteCompleto.nome_do_site).all()]
    pendentes   = _admin_pendentes()

    if request.method == "POST":
        action     = request.form.get("action")
        usuario_id = request.form.get("usuario_id", type=int)
        site_nome  = (request.form.get("site_nome") or "").strip()

        if not usuario_id:
            flash("Usuário não selecionado.", "danger")
            return redirect(url_for("admin_multisites"))

        usuario = Usuario.query.get_or_404(usuario_id)

        if action == "adicionar":
            if not site_nome:
                flash("Selecione um site para vincular.", "danger")
            elif UsuarioSite.query.filter_by(usuario_id=usuario_id, site_nome=site_nome).first():
                flash(f"Site '{site_nome}' já está vinculado a {usuario.nome}.", "warning")
            else:
                db.session.add(UsuarioSite(usuario_id=usuario_id, site_nome=site_nome))
                db.session.commit()
                flash(f"✅ Site '{site_nome}' vinculado a {usuario.nome}.", "success")

        elif action == "remover":
            vinculo = UsuarioSite.query.filter_by(usuario_id=usuario_id, site_nome=site_nome).first()
            if vinculo:
                db.session.delete(vinculo)
                db.session.commit()
                flash(f"Site '{site_nome}' removido de {usuario.nome}.", "success")
            else:
                flash("Vínculo não encontrado.", "warning")

        elif action == "remover_todos":
            UsuarioSite.query.filter_by(usuario_id=usuario_id).delete()
            db.session.commit()
            flash(f"Todos os sites de {usuario.nome} foram removidos.", "success")

        return redirect(url_for("admin_multisites"))

    # GET — carrega vínculos de todos os usuários MULTISITES em UMA query
    # (elimina N+1 — antes era 1 query por usuário)
    vinculos_por_usuario = {u.id: [] for u in todos_usuarios}
    _ids = [u.id for u in todos_usuarios]
    if _ids:
        _todos_vinculos = (UsuarioSite.query
                            .filter(UsuarioSite.usuario_id.in_(_ids))
                            .order_by(UsuarioSite.site_nome).all())
        for v in _todos_vinculos:
            vinculos_por_usuario[v.usuario_id].append(v.site_nome)

    return render_template(
        "admin_multisites.html",
        todos_usuarios=todos_usuarios,
        todos_sites=todos_sites,
        vinculos_por_usuario=vinculos_por_usuario,
        pendentes=pendentes,
    )


# ===========================================================================
# SHIFT HANDOVER — Modelos
# ===========================================================================
class OcorrenciaTurno(db.Model):
    __tablename__ = "ocorrencias_turno"

    id                  = db.Column(db.Integer, db.Identity(start=1), primary_key=True)
    data_ocorrencia     = db.Column(db.Date,    nullable=False, default=date.today)
    data_hora_registro  = db.Column(db.DateTime, nullable=False, default=datetime.now)
    site                = db.Column(db.String(80),  nullable=False)
    turno               = db.Column(db.String(30),  nullable=False)
    setor               = db.Column(db.String(100), nullable=False)
    tipo_ocorrencia     = db.Column(db.String(100), nullable=False)
    prioridade          = db.Column(db.String(20),  nullable=False)
    responsavel_saida   = db.Column(db.String(120), nullable=False)
    responsavel_entrada = db.Column(db.String(120), nullable=False)
    descricao           = db.Column(db.Text, nullable=False)
    efetivo             = db.Column(db.Text, nullable=False)
    assinatura_saida       = db.Column(db.Text, nullable=True)
    assinatura_entrada     = db.Column(db.Text, nullable=True)
    imagem_1               = db.Column(db.Text, nullable=True)
    imagem_2               = db.Column(db.Text, nullable=True)
    imagem_3               = db.Column(db.Text, nullable=True)
    imagem_4               = db.Column(db.Text, nullable=True)
    acoes_tomadas          = db.Column(db.Text, nullable=True)
    pendencias             = db.Column(db.Text, nullable=True)
    status                 = db.Column(db.String(40), nullable=False)
    criado_por             = db.Column(db.String(120), nullable=True)
    created_at             = db.Column(db.DateTime, nullable=False, default=datetime.now)
    updated_at             = db.Column(db.DateTime, nullable=True)
    # campos de recebimento
    ressalva               = db.Column(db.Text, nullable=True)
    tem_ressalva           = db.Column(db.String(1), nullable=True, default="N")
    anexo_entrada          = db.Column(db.Text, nullable=True)   # base64
    anexo_entrada_nome     = db.Column(db.String(255), nullable=True)

    def to_dict(self):
        from sqlalchemy import inspect as _sa_inspect
        _unloaded = _sa_inspect(self).unloaded_expirable
        def _clob(attr):
            return "" if attr in _unloaded else (getattr(self, attr) or "")
        return {
            "id": self.id,
            "data_ocorrencia": self.data_ocorrencia.strftime("%d/%m/%Y") if self.data_ocorrencia else "",
            "data_hora_registro": self.data_hora_registro.strftime("%d/%m/%Y %H:%M") if self.data_hora_registro else "",
            "site": self.site, "turno": self.turno, "setor": self.setor,
            "tipo_ocorrencia": self.tipo_ocorrencia, "prioridade": self.prioridade,
            "responsavel_saida": self.responsavel_saida, "responsavel_entrada": self.responsavel_entrada,
            "descricao": self.descricao, "efetivo": self.efetivo or "",
            "assinatura_saida": _clob("assinatura_saida"),
            "assinatura_entrada": _clob("assinatura_entrada"),
            "imagem_1": _clob("imagem_1"), "imagem_2": _clob("imagem_2"),
            "imagem_3": _clob("imagem_3"), "imagem_4": _clob("imagem_4"),
            "acoes_tomadas": self.acoes_tomadas or "", "pendencias": self.pendencias or "",
            "status": self.status, "criado_por": self.criado_por or "",
            "criado_em": self.created_at.strftime("%d/%m/%Y %H:%M") if self.created_at else "",
            "atualizado_em": self.updated_at.strftime("%d/%m/%Y %H:%M") if self.updated_at else "",
            "ressalva": _clob("ressalva"),
            "tem_ressalva": self.tem_ressalva or "N",
            "anexo_entrada_nome": self.anexo_entrada_nome or "",
        }


class SiteSH(db.Model):
    __tablename__ = "SITES"
    id_site   = db.Column("ID_SITE",   db.Integer,    primary_key=True)
    nome_site = db.Column("NOME_SITE", db.String(100), nullable=False)


# ===========================================================================
# SHIFT HANDOVER — Helpers
# ===========================================================================
from datetime import date as _date_cls

_SH_TURNOS    = {"TURNO A", "TURNO B", "TURNO C", "ADM"}
_SH_STATUS    = {"EM ABERTO", "EM ACOMPANHAMENTO", "FINALIZADO"}
_SH_PRIORS    = {"BAIXA", "MEDIA", "ALTA", "CRITICA"}


def _sh_norm(v): return (v or "").strip().upper()

def _sh_norm_prioridade(v):
    v = _sh_norm(v)
    return "CRITICA" if v == "CRÍTICA" else "MEDIA" if v == "MÉDIA" else v

def _sh_norm_tipo(v):
    v = _sh_norm(v)
    return "MANUTENCAO" if v == "MANUTENÇÃO" else "PENDENCIA" if v == "PENDÊNCIA" else v

def _sh_parse_date(val):
    val = (val or "").strip()
    if not val: return None
    try:    return datetime.strptime(val, "%Y-%m-%d").date()
    except: return None

def _sh_parse_dt(val):
    val = (val or "").strip()
    if not val: return None
    try:    return datetime.strptime(val, "%Y-%m-%dT%H:%M")
    except: return None

def _sh_img_b64(file_storage):
    if not file_storage or not getattr(file_storage, "filename", ""):
        return None
    fn = file_storage.filename.strip()
    if "." not in fn or fn.rsplit(".", 1)[1].lower() not in {"png","jpg","jpeg","webp"}:
        return None
    try:
        img = _PILImage.open(file_storage)
        if img.mode != "RGB":
            img = img.convert("RGB")
        img.thumbnail((800, 800), _PILImage.Resampling.LANCZOS)
        buf = BytesIO()
        img.save(buf, format="JPEG", quality=70, optimize=True)
        return "data:image/jpeg;base64," + base64.b64encode(buf.getvalue()).decode()
    except Exception as e:
        logging.error(f"[SH] img_b64 error: {e}")
        return None

def _sh_verificar_acesso(oc):
    if session.get("user_perfil") in ("ADMIN", "Admin", "admin"):
        return True
    return (_sh_norm(oc.site) == _sh_norm(session.get("user_site", "")))

def _sh_get_filtros():
    di      = (request.args.get("data_inicial") or "").strip()
    df      = (request.args.get("data_final")   or "").strip()
    turno   = _sh_norm(request.args.get("turno"))
    status  = _sh_norm(request.args.get("status"))
    assinatura = _sh_norm(request.args.get("assinatura"))
    perfil  = session.get("user_perfil", "")
    is_adm  = perfil in ("ADMIN", "Admin", "admin")
    site_f  = _sh_norm(request.args.get("site")) if is_adm else _sh_norm(session.get("user_site", ""))

    q = OcorrenciaTurno.query
    if site_f:  q = q.filter(OcorrenciaTurno.site == site_f)
    if di:
        d = _sh_parse_date(di)
        if d: q = q.filter(OcorrenciaTurno.data_ocorrencia >= d)
    if df:
        d = _sh_parse_date(df)
        if d: q = q.filter(OcorrenciaTurno.data_ocorrencia <= d)
    if turno:  q = q.filter(OcorrenciaTurno.turno  == turno)
    if status: q = q.filter(OcorrenciaTurno.status == status)
    if assinatura == "NAO_ASSINADA":
        q = q.filter(or_(OcorrenciaTurno.assinatura_entrada.is_(None),
                          func.length(OcorrenciaTurno.assinatura_entrada) == 0))
    elif assinatura == "ASSINADA":
        q = q.filter(OcorrenciaTurno.assinatura_entrada.isnot(None),
                      func.length(OcorrenciaTurno.assinatura_entrada) > 0)
    q = q.order_by(OcorrenciaTurno.data_hora_registro.desc(), OcorrenciaTurno.id.desc())
    return q, {"data_inicial": di, "data_final": df, "turno": turno,
               "status": status, "assinatura": assinatura, "site": site_f if is_adm else ""}

def _sh_resumo():
    hoje   = _date_cls.today()
    is_adm = session.get("user_perfil", "") in ("ADMIN", "Admin", "admin")
    site_u = _sh_norm(session.get("user_site", ""))
    # Uma única query com CASE/SUM em vez de 4 queries separadas
    base = db.session.query(
        func.sum(case((OcorrenciaTurno.data_ocorrencia == hoje, 1), else_=0)).label("dia"),
        func.sum(case((OcorrenciaTurno.status.in_(["EM ABERTO","EM ACOMPANHAMENTO"]), 1), else_=0)).label("abertas"),
        func.sum(case(((OcorrenciaTurno.data_ocorrencia == hoje) & (OcorrenciaTurno.turno != None), 1), else_=0)).label("turnos"),
        func.sum(case((OcorrenciaTurno.prioridade == "CRITICA", 1), else_=0)).label("criticas"),
    )
    if not is_adm and site_u:
        base = base.filter(OcorrenciaTurno.site == site_u)
    row = base.one()
    return {
        "ocorrencias_dia":      int(row.dia      or 0),
        "pendencias_abertas":   int(row.abertas  or 0),
        "turnos_registrados":   int(row.turnos   or 0),
        "ocorrencias_criticas": int(row.criticas or 0),
    }


# ===========================================================================
# SHIFT HANDOVER — Rotas
# ===========================================================================
@app.route("/shift-handover/")
@login_required
def sh_index():
    q, filtros = _sh_get_filtros()
    # Carrega lista SEM CLOBs pesados (imagens/assinaturas) — detalhes via AJAX
    _clobs = (
        defer(OcorrenciaTurno.imagem_1), defer(OcorrenciaTurno.imagem_2),
        defer(OcorrenciaTurno.imagem_3), defer(OcorrenciaTurno.imagem_4),
        defer(OcorrenciaTurno.assinatura_saida), defer(OcorrenciaTurno.assinatura_entrada),
        defer(OcorrenciaTurno.ressalva), defer(OcorrenciaTurno.anexo_entrada),
    )
    ocs_db    = q.options(*_clobs).limit(200).all()
    ocs       = [o.to_dict() for o in ocs_db]
    ultima_oc = ocs[0] if ocs else None
    ultimo_id = db.session.query(func.max(OcorrenciaTurno.id)).scalar() or 0

    user_site  = session.get("user_site", "")
    user_id    = session.get("user_id")
    if user_site:
        usuarios_site = Usuario.query.filter(
            Usuario.site == user_site, Usuario.id != user_id, Usuario.is_active == True
        ).order_by(Usuario.nome.asc()).all()
    else:
        usuarios_site = Usuario.query.filter(
            Usuario.id != user_id, Usuario.is_active == True
        ).order_by(Usuario.nome.asc()).all()

    sites_db = SiteCompleto.query.order_by(SiteCompleto.nome_do_site.asc()).all()
    return render_template("sh_registrar.html",
        resumo=_sh_resumo(), ultima_ocorrencia=ultima_oc, ocorrencias=ocs,
        filtros=filtros, hoje=_date_cls.today().strftime("%Y-%m-%d"),
        proximo_id_previsto=ultimo_id+1, usuarios_mesmo_site=usuarios_site,
        sites=sites_db, setores=_get_setores(user_site))


@app.route("/shift-handover/controle")
@login_required
def sh_controle():
    q, filtros = _sh_get_filtros()
    _clobs = (
        defer(OcorrenciaTurno.imagem_1), defer(OcorrenciaTurno.imagem_2),
        defer(OcorrenciaTurno.imagem_3), defer(OcorrenciaTurno.imagem_4),
        defer(OcorrenciaTurno.assinatura_saida), defer(OcorrenciaTurno.assinatura_entrada),
        defer(OcorrenciaTurno.ressalva), defer(OcorrenciaTurno.anexo_entrada),
    )
    ocs_db = q.options(*_clobs).limit(200).all()
    ocs    = [o.to_dict() for o in ocs_db]
    return render_template("sh_controle.html",
        resumo=_sh_resumo(), ocorrencias=ocs, filtros=filtros)


@app.route("/shift-handover/<int:oc_id>/detalhe")
@login_required
def sh_detalhe(oc_id):
    """API JSON: retorna os detalhes completos (com CLOBs) de uma passagem."""
    from flask import jsonify
    oc = OcorrenciaTurno.query.get_or_404(oc_id)
    if not _sh_verificar_acesso(oc):
        return jsonify({"error": "Acesso negado"}), 403
    d = oc.to_dict()
    d["assinatura_saida"]   = oc.assinatura_saida   or ""
    d["assinatura_entrada"] = oc.assinatura_entrada  or ""
    d["imagem_1"] = oc.imagem_1 or ""
    d["imagem_2"] = oc.imagem_2 or ""
    d["imagem_3"] = oc.imagem_3 or ""
    d["imagem_4"] = oc.imagem_4 or ""
    d["ressalva"] = oc.ressalva or ""
    d["tem_ressalva"] = oc.tem_ressalva or "N"
    return jsonify(d)


@app.route("/shift-handover/salvar", methods=["POST"])
@login_required
def sh_salvar():
    try:
        is_adm = session.get("user_perfil", "") in ("ADMIN", "Admin", "admin")
        site   = request.form.get("site") if is_adm else session.get("user_site", "")
        turno  = _sh_norm(request.form.get("turno"))
        setor  = _sh_norm(request.form.get("setor"))
        tipo   = _sh_norm_tipo(request.form.get("tipo_ocorrencia"))
        prior  = _sh_norm_prioridade(request.form.get("prioridade"))
        status = _sh_norm(request.form.get("status"))
        resp_saida    = session.get("user_nome", "Usuário")
        resp_entrada  = (request.form.get("responsavel_entrada") or "").strip()
        descricao     = (request.form.get("descricao") or "").strip()
        efetivo       = (request.form.get("efetivo") or "").strip()
        ass_saida     = request.form.get("assinatura_saida") or ""
        acoes         = (request.form.get("acoes_tomadas") or "").strip()
        pendencias    = (request.form.get("pendencias") or "").strip()
        data_oc       = _sh_parse_date(request.form.get("data_ocorrencia"))
        data_hr       = _sh_parse_dt(request.form.get("data_hora_registro"))

        if not all([data_oc, data_hr, site, turno, setor, tipo, prior,
                    resp_saida, resp_entrada, descricao, efetivo, status, ass_saida]):
            flash("Preencha todos os campos obrigatórios e realize sua assinatura.", "danger")
            return redirect(url_for("sh_index"))

        nova = OcorrenciaTurno(
            data_ocorrencia=data_oc, data_hora_registro=data_hr,
            site=site, turno=turno, setor=setor, tipo_ocorrencia=tipo,
            prioridade=prior, responsavel_saida=resp_saida, responsavel_entrada=resp_entrada,
            descricao=descricao, efetivo=efetivo, assinatura_saida=ass_saida,
            imagem_1=_sh_img_b64(request.files.get("imagem_1")),
            imagem_2=_sh_img_b64(request.files.get("imagem_2")),
            imagem_3=_sh_img_b64(request.files.get("imagem_3")),
            imagem_4=_sh_img_b64(request.files.get("imagem_4")),
            acoes_tomadas=acoes or None, pendencias=pendencias or None,
            status=status, criado_por=resp_saida,
        )
        db.session.add(nova)
        db.session.commit()
        flash("Passagem de turno registrada com sucesso! Aguardando assinatura do recebedor.", "success")
    except Exception as e:
        db.session.rollback()
        logging.error(f"[SH salvar] {e}")
        flash(f"Erro ao salvar: {e}", "danger")
    return redirect(url_for("sh_index"))


@app.route("/shift-handover/<int:oc_id>/assinar", methods=["GET", "POST"])
@login_required
def sh_assinar(oc_id):
    oc = OcorrenciaTurno.query.get_or_404(oc_id)
    if not _sh_verificar_acesso(oc):
        flash("Acesso negado.", "danger")
        return redirect(url_for("sh_controle"))
    if oc.responsavel_entrada != session.get("user_nome"):
        flash("Você não é o responsável designado para receber este turno.", "danger")
        return redirect(url_for("sh_controle"))
    if oc.assinatura_entrada:
        flash("Este turno já foi assinado.", "warning")
        return redirect(url_for("sh_controle"))
    if request.method == "POST":
        ass = request.form.get("assinatura_entrada")
        if not ass:
            flash("Assinatura obrigatória.", "danger")
            return redirect(url_for("sh_assinar", oc_id=oc.id))
        oc.assinatura_entrada = ass
        oc.updated_at         = datetime.now()
        # Ressalva
        ressalva_txt = (request.form.get("ressalva") or "").strip()
        if ressalva_txt:
            oc.ressalva     = ressalva_txt
            oc.tem_ressalva = "S"
        else:
            oc.tem_ressalva = "N"
        # Anexo do recebedor
        anexo_file = request.files.get("anexo_entrada")
        if anexo_file and getattr(anexo_file, "filename", ""):
            try:
                data = anexo_file.read()
                oc.anexo_entrada      = "data:application/octet-stream;base64," + base64.b64encode(data).decode()
                oc.anexo_entrada_nome = anexo_file.filename
            except Exception as ex:
                logging.error(f"[SH anexo] {ex}")
        db.session.commit()
        if oc.tem_ressalva == "S":
            flash("Turno recebido com ressalva registrada!", "warning")
        else:
            flash("Turno recebido e assinado com sucesso!", "success")
        return redirect(url_for("sh_controle"))
    return render_template("sh_assinar.html", ocorrencia=oc)


@app.route("/shift-handover/<int:oc_id>/editar", methods=["GET", "POST"])
@login_required
def sh_editar(oc_id):
    oc = OcorrenciaTurno.query.get_or_404(oc_id)
    if not _sh_verificar_acesso(oc):
        flash("Acesso negado.", "danger")
        return redirect(url_for("sh_index"))
    if oc.criado_por != session.get("user_nome"):
        flash("Somente quem registrou a passagem pode editá-la.", "danger")
        return redirect(url_for("sh_index"))
    if oc.assinatura_entrada:
        flash("Não é possível editar após o recebedor assinar.", "warning")
        return redirect(url_for("sh_index"))
    if request.method == "POST":
        try:
            is_adm = session.get("user_perfil", "") in ("ADMIN", "Admin", "admin")
            if is_adm:
                oc.site = request.form.get("site") or oc.site
            oc.turno           = _sh_norm(request.form.get("turno")) or oc.turno
            oc.setor           = _sh_norm(request.form.get("setor")) or oc.setor
            oc.tipo_ocorrencia = _sh_norm_tipo(request.form.get("tipo_ocorrencia")) or oc.tipo_ocorrencia
            oc.prioridade      = _sh_norm_prioridade(request.form.get("prioridade")) or oc.prioridade
            oc.status          = _sh_norm(request.form.get("status")) or oc.status
            oc.responsavel_entrada = (request.form.get("responsavel_entrada") or oc.responsavel_entrada).strip()
            oc.efetivo         = (request.form.get("efetivo") or "").strip() or oc.efetivo
            oc.descricao       = (request.form.get("descricao") or "").strip() or oc.descricao
            oc.acoes_tomadas   = (request.form.get("acoes_tomadas") or "").strip() or None
            oc.pendencias      = (request.form.get("pendencias") or "").strip() or None
            # atualiza assinatura saída se nova for fornecida
            nova_ass = request.form.get("assinatura_saida") or ""
            if nova_ass:
                oc.assinatura_saida = nova_ass
            # atualiza imagens se enviadas
            for i in range(1, 5):
                nova_img = _sh_img_b64(request.files.get(f"imagem_{i}"))
                if nova_img:
                    setattr(oc, f"imagem_{i}", nova_img)
            oc.data_ocorrencia    = _sh_parse_date(request.form.get("data_ocorrencia")) or oc.data_ocorrencia
            oc.data_hora_registro = _sh_parse_dt(request.form.get("data_hora_registro"))   or oc.data_hora_registro
            oc.updated_at = datetime.now()
            db.session.commit()
            flash("Passagem atualizada com sucesso!", "success")
            return redirect(url_for("sh_index"))
        except Exception as e:
            db.session.rollback()
            logging.error(f"[SH editar] {e}")
            flash(f"Erro ao atualizar: {e}", "danger")
    sites_db      = SiteCompleto.query.order_by(SiteCompleto.nome_do_site.asc()).all()
    user_site     = session.get("user_site", "")
    user_id       = session.get("user_id")
    usuarios_site = Usuario.query.filter(
        Usuario.site == user_site, Usuario.id != user_id, Usuario.is_active == True
    ).order_by(Usuario.nome.asc()).all() if user_site else \
        Usuario.query.filter(Usuario.id != user_id, Usuario.is_active == True).order_by(Usuario.nome.asc()).all()
    return render_template("sh_editar.html", oc=oc, sites=sites_db,
        usuarios_mesmo_site=usuarios_site,
        data_oc_val=oc.data_ocorrencia.strftime("%Y-%m-%d") if oc.data_ocorrencia else "",
        data_hr_val=oc.data_hora_registro.strftime("%Y-%m-%dT%H:%M") if oc.data_hora_registro else "")


@app.route("/shift-handover/<int:oc_id>/exportar-pdf")
@login_required
def sh_pdf(oc_id):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import mm
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_LEFT
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image as RLImg
    from reportlab.lib import colors as rl_colors

    oc = OcorrenciaTurno.query.get_or_404(oc_id)
    if not _sh_verificar_acesso(oc):
        flash("Acesso negado.", "danger")
        return redirect(url_for("sh_index"))

    buf = BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4,
        leftMargin=14*mm, rightMargin=14*mm, topMargin=14*mm, bottomMargin=14*mm)
    styles = getSampleStyleSheet()

    # Estilos
    titulo_s = ParagraphStyle("shTit", parent=styles["Title"],
        fontSize=16, textColor=rl_colors.HexColor("#D40511"),
        alignment=TA_CENTER, leading=20, spaceAfter=4)
    sub_s    = ParagraphStyle("shSub", parent=styles["Normal"],
        fontSize=8, textColor=rl_colors.HexColor("#555555"),
        alignment=TA_CENTER, spaceAfter=8)
    sec_s    = ParagraphStyle("shSec", parent=styles["Heading3"],
        fontSize=9, textColor=rl_colors.HexColor("#D40511"),
        spaceBefore=6, spaceAfter=4, leading=11)
    cell_s   = ParagraphStyle("shCell", parent=styles["BodyText"],
        fontSize=8, leading=10)
    warn_s   = ParagraphStyle("shWarn", parent=styles["BodyText"],
        fontSize=8, leading=10, textColor=rl_colors.HexColor("#92400e"),
        backColor=rl_colors.HexColor("#fff4db"))

    def _v(val): return str(val or "-").strip() or "-"
    def _p(txt, s=cell_s): return Paragraph(_v(txt).replace("\n","<br/>"), s)

    elems = []

    # Cabeçalho DHL
    elems.append(Paragraph("SHIFT HANDOVER — PASSAGEM DE TURNO", titulo_s))
    elems.append(Paragraph(
        f"DHL SECURITY • {_v(oc.site)} • {oc.data_ocorrencia.strftime('%d/%m/%Y') if oc.data_ocorrencia else '-'}", sub_s))

    # Faixa resumo
    faixa = Table([[
        _p(f"<b>ID:</b> #{oc.id}"),
        _p(f"<b>Turno:</b> {_v(oc.turno)}"),
        _p(f"<b>Prioridade:</b> {_v(oc.prioridade)}"),
        _p(f"<b>Status:</b> {_v(oc.status)}"),
    ]], colWidths=[130, 100, 120, 110])
    faixa.setStyle(TableStyle([
        ("BACKGROUND",(0,0),(-1,-1),rl_colors.HexColor("#fff4cc")),
        ("BOX",(0,0),(-1,-1),0.6,rl_colors.HexColor("#ffcc00")),
        ("INNERGRID",(0,0),(-1,-1),0.4,rl_colors.HexColor("#ffcc00")),
        ("LEFTPADDING",(0,0),(-1,-1),6), ("RIGHTPADDING",(0,0),(-1,-1),6),
        ("TOPPADDING",(0,0),(-1,-1),5), ("BOTTOMPADDING",(0,0),(-1,-1),5),
        ("VALIGN",(0,0),(-1,-1),"MIDDLE"),
    ]))
    elems.append(faixa)
    elems.append(Spacer(1, 6))

    # Tabela de dados
    dados_e = [
        ["ID", str(oc.id)],
        ["Data", oc.data_ocorrencia.strftime("%d/%m/%Y") if oc.data_ocorrencia else "-"],
        ["Data/Hora", oc.data_hora_registro.strftime("%d/%m/%Y %H:%M") if oc.data_hora_registro else "-"],
        ["Site", _v(oc.site)], ["Turno", _v(oc.turno)], ["Setor", _v(oc.setor)],
        ["Tipo", _v(oc.tipo_ocorrencia)], ["Prioridade", _v(oc.prioridade)], ["Status", _v(oc.status)],
    ]
    dados_d = [
        ["Resp. saída",   _v(oc.responsavel_saida)],
        ["Resp. entrada", _v(oc.responsavel_entrada)],
        ["Criado por",    _v(oc.criado_por)],
        ["Criado em",     oc.created_at.strftime("%d/%m/%Y %H:%M") if oc.created_at else "-"],
        ["Atualizado",    oc.updated_at.strftime("%d/%m/%Y %H:%M") if oc.updated_at else "-"],
    ]
    _ts = TableStyle([
        ("BACKGROUND",(0,0),(0,-1),rl_colors.HexColor("#ffcc00")),
        ("FONTNAME",(0,0),(0,-1),"Helvetica-Bold"),
        ("FONTSIZE",(0,0),(-1,-1),7.5),
        ("GRID",(0,0),(-1,-1),0.35,rl_colors.HexColor("#cfcfcf")),
        ("VALIGN",(0,0),(-1,-1),"TOP"),
        ("LEFTPADDING",(0,0),(-1,-1),5),("RIGHTPADDING",(0,0),(-1,-1),5),
        ("TOPPADDING",(0,0),(-1,-1),4),("BOTTOMPADDING",(0,0),(-1,-1),4),
    ])
    t_e = Table(dados_e, colWidths=[82,168]); t_e.setStyle(_ts)
    t_d = Table(dados_d, colWidths=[82,168]); t_d.setStyle(_ts)
    bloco = Table([[t_e, t_d]], colWidths=[250,250])
    bloco.setStyle(TableStyle([
        ("VALIGN",(0,0),(-1,-1),"TOP"),
        ("LEFTPADDING",(0,0),(-1,-1),0),("RIGHTPADDING",(0,0),(-1,-1),0),
        ("TOPPADDING",(0,0),(-1,-1),0),("BOTTOMPADDING",(0,0),(-1,-1),0),
    ]))
    elems.append(bloco)
    elems.append(Spacer(1, 6))

    # Seções de texto
    for titulo, valor in [("Efetivo", oc.efetivo), ("Descrição", oc.descricao),
                           ("Ações Tomadas", oc.acoes_tomadas), ("Pendências", oc.pendencias)]:
        if not valor: continue
        elems.append(Paragraph(titulo, sec_s))
        box = Table([[_p(valor)]], colWidths=[500])
        box.setStyle(TableStyle([
            ("BOX",(0,0),(-1,-1),0.4,rl_colors.HexColor("#d9d9d9")),
            ("BACKGROUND",(0,0),(-1,-1),rl_colors.HexColor("#fbfbfb")),
            ("LEFTPADDING",(0,0),(-1,-1),6),("RIGHTPADDING",(0,0),(-1,-1),6),
            ("TOPPADDING",(0,0),(-1,-1),5),("BOTTOMPADDING",(0,0),(-1,-1),5),
        ]))
        elems.append(box)
        elems.append(Spacer(1, 4))

    # Ressalva
    if oc.ressalva:
        elems.append(Paragraph("⚠ Ressalva do Recebedor", sec_s))
        box_r = Table([[_p(oc.ressalva, warn_s)]], colWidths=[500])
        box_r.setStyle(TableStyle([
            ("BOX",(0,0),(-1,-1),0.8,rl_colors.HexColor("#f5d66f")),
            ("BACKGROUND",(0,0),(-1,-1),rl_colors.HexColor("#fff4db")),
            ("LEFTPADDING",(0,0),(-1,-1),6),("RIGHTPADDING",(0,0),(-1,-1),6),
            ("TOPPADDING",(0,0),(-1,-1),5),("BOTTOMPADDING",(0,0),(-1,-1),5),
        ]))
        elems.append(box_r)
        elems.append(Spacer(1, 4))

    # Assinaturas
    def _ass_img(b64, w_mm, h_mm):
        if not b64: return None
        try:
            raw = b64.split(",",1)[1] if "," in b64 else b64
            raw += "=" * ((4-len(raw)%4)%4)
            return RLImg(BytesIO(base64.b64decode(raw)), width=w_mm*mm, height=h_mm*mm)
        except: return None

    elems.append(Paragraph("Assinaturas", sec_s))
    img_saida  = _ass_img(oc.assinatura_saida,  55, 20)
    img_entrada = _ass_img(oc.assinatura_entrada, 55, 20)
    ass_tab = Table([
        [img_saida  or _p("-"), img_entrada or _p("Não assinado")],
        [_p(f"<b>Saída:</b> {_v(oc.responsavel_saida)}"),
         _p(f"<b>Entrada:</b> {_v(oc.responsavel_entrada)}")],
    ], colWidths=[250, 250])
    ass_tab.setStyle(TableStyle([
        ("BOX",(0,0),(-1,-1),0.35,rl_colors.HexColor("#d9d9d9")),
        ("INNERGRID",(0,0),(-1,-1),0.35,rl_colors.HexColor("#d9d9d9")),
        ("VALIGN",(0,0),(-1,-1),"MIDDLE"),
        ("ALIGN",(0,0),(-1,0),"CENTER"),
        ("BACKGROUND",(0,0),(-1,-1),rl_colors.HexColor("#fcfcfc")),
        ("LEFTPADDING",(0,0),(-1,-1),5),("RIGHTPADDING",(0,0),(-1,-1),5),
        ("TOPPADDING",(0,0),(-1,-1),5),("BOTTOMPADDING",(0,0),(-1,-1),5),
    ]))
    elems.append(ass_tab)
    elems.append(Spacer(1, 6))

    # Imagens
    imgs_b64 = [getattr(oc,f"imagem_{i}") for i in range(1,5) if getattr(oc,f"imagem_{i}")]
    if imgs_b64:
        elems.append(Paragraph("Evidências Fotográficas", sec_s))
        def _fit(b64, mw, mh):
            try:
                raw = b64.split(",",1)[1] if "," in b64 else b64
                raw += "=" * ((4-len(raw)%4)%4)
                img = RLImg(BytesIO(base64.b64decode(raw)))
                iw,ih = img.imageWidth, img.imageHeight
                if not iw or not ih: return None
                p = min(mw/iw, mh/ih)
                img.drawWidth = iw*p; img.drawHeight = ih*p
                return img
            except: return None
        fotos_rows = []
        linha = []
        for b in imgs_b64:
            fi = _fit(b, 235, 130)
            linha.append(fi or "")
            if len(linha)==2:
                fotos_rows.append(linha); linha=[]
        if linha:
            while len(linha)<2: linha.append("")
            fotos_rows.append(linha)
        ftab = Table(fotos_rows, colWidths=[255,255])
        ftab.setStyle(TableStyle([
            ("VALIGN",(0,0),(-1,-1),"MIDDLE"), ("ALIGN",(0,0),(-1,-1),"CENTER"),
            ("BOX",(0,0),(-1,-1),0.35,rl_colors.HexColor("#dddddd")),
            ("INNERGRID",(0,0),(-1,-1),0.35,rl_colors.HexColor("#dddddd")),
            ("LEFTPADDING",(0,0),(-1,-1),4),("RIGHTPADDING",(0,0),(-1,-1),4),
            ("TOPPADDING",(0,0),(-1,-1),4),("BOTTOMPADDING",(0,0),(-1,-1),4),
        ]))
        elems.append(ftab)

    # Rodapé DHL no canvas
    def _draw_header_footer(canvas, _doc):
        canvas.saveState()
        pw = A4[0]
        ph = A4[1]
        canvas.setFillColor(rl_colors.HexColor("#D40511"))
        canvas.rect(0, ph-11*mm, pw, 11*mm, fill=1, stroke=0)
        canvas.setFillColor(rl_colors.HexColor("#FFCC00"))
        canvas.rect(0, ph-13*mm, pw, 2*mm, fill=1, stroke=0)
        canvas.setFillColor(rl_colors.white)
        canvas.setFont("Helvetica-Bold", 10)
        canvas.drawString(14*mm, ph-7.5*mm, "DHL SECURITY — SHIFT HANDOVER")
        canvas.setFont("Helvetica", 8)
        canvas.drawRightString(pw-14*mm, ph-7.5*mm,
            f"Gerado em {datetime.now().strftime('%d/%m/%Y %H:%M')}")
        canvas.setStrokeColor(rl_colors.HexColor("#D1D5DB"))
        canvas.setLineWidth(0.4)
        canvas.line(14*mm, 9*mm, pw-14*mm, 9*mm)
        canvas.setFont("Helvetica", 7.5)
        canvas.setFillColor(rl_colors.HexColor("#6B7280"))
        canvas.drawString(14*mm, 5.5*mm, f"Passagem #{oc.id} • {_v(oc.site)} • {_v(oc.turno)}")
        canvas.drawRightString(pw-14*mm, 5.5*mm, f"Pág. {canvas.getPageNumber()}")
        canvas.restoreState()

    doc.build(elems, onFirstPage=_draw_header_footer, onLaterPages=_draw_header_footer)
    buf.seek(0)
    nome = f"Shift_Handover_{oc.id}_{_v(oc.site)}.pdf"
    return send_file(buf, as_attachment=True, download_name=nome, mimetype="application/pdf")


@app.route("/shift-handover/<int:oc_id>/fechar", methods=["POST"])
@login_required
def sh_fechar(oc_id):
    oc = OcorrenciaTurno.query.get_or_404(oc_id)
    if not _sh_verificar_acesso(oc):
        flash("Acesso negado.", "danger")
        return redirect(url_for("sh_index"))
    if oc.responsavel_entrada != session.get("user_nome"):
        flash("Apenas o responsável que assumiu pode finalizar.", "danger")
        return redirect(url_for("sh_index"))
    if not oc.assinatura_entrada:
        flash("Assine o recebimento antes de finalizar.", "warning")
        return redirect(url_for("sh_index"))
    try:
        oc.status     = "FINALIZADO"
        oc.updated_at = datetime.now()
        db.session.commit()
        flash("Ocorrência finalizada.", "success")
    except Exception as e:
        db.session.rollback()
        flash(f"Erro: {e}", "danger")
    return redirect(url_for("sh_index"))


@app.route("/shift-handover/<int:oc_id>/excluir", methods=["POST"])
@login_required
def sh_excluir(oc_id):
    oc = OcorrenciaTurno.query.get_or_404(oc_id)
    if not _sh_verificar_acesso(oc):
        flash("Acesso negado.", "danger")
        return redirect(url_for("sh_index"))
    try:
        db.session.delete(oc)
        db.session.commit()
        flash("Registro excluído.", "success")
    except Exception as e:
        db.session.rollback()
        flash(f"Erro: {e}", "danger")
    return redirect(url_for("sh_index"))


@app.route("/shift-handover/dashboard")
@login_required
def sh_dashboard():
    is_adm  = session.get("user_perfil", "") in ("ADMIN", "Admin", "admin")
    site_u  = _sh_norm(session.get("user_site", ""))

    def _qd(col):
        q = db.session.query(col)
        if not is_adm: q = q.filter(OcorrenciaTurno.site == site_u)
        return q

    total        = _qd(func.count(OcorrenciaTurno.id)).scalar() or 0
    em_aberto    = _qd(func.count(OcorrenciaTurno.id)).filter(OcorrenciaTurno.status == "EM ABERTO").scalar() or 0
    acompanhamento = _qd(func.count(OcorrenciaTurno.id)).filter(OcorrenciaTurno.status == "EM ACOMPANHAMENTO").scalar() or 0
    finalizado   = _qd(func.count(OcorrenciaTurno.id)).filter(OcorrenciaTurno.status == "FINALIZADO").scalar() or 0

    _ord_turno = case(
        (OcorrenciaTurno.turno == "TURNO A", 1),
        (OcorrenciaTurno.turno == "TURNO B", 2),
        (OcorrenciaTurno.turno == "TURNO C", 3),
        (OcorrenciaTurno.turno == "ADM", 4), else_=99)
    _ord_prior = case(
        (OcorrenciaTurno.prioridade == "CRITICA", 1),
        (OcorrenciaTurno.prioridade == "ALTA", 2),
        (OcorrenciaTurno.prioridade == "MEDIA", 3),
        (OcorrenciaTurno.prioridade == "BAIXA", 4), else_=99)

    por_turno = [(r[1], r[0]) for r in
        _qd(func.count(OcorrenciaTurno.id)).add_columns(OcorrenciaTurno.turno)
        .group_by(OcorrenciaTurno.turno).order_by(_ord_turno).all()]
    por_prior = [(r[1], r[0]) for r in
        _qd(func.count(OcorrenciaTurno.id)).add_columns(OcorrenciaTurno.prioridade)
        .group_by(OcorrenciaTurno.prioridade).order_by(_ord_prior).all()]
    por_site  = [(r[1], r[0]) for r in
        _qd(func.count(OcorrenciaTurno.id)).add_columns(OcorrenciaTurno.site)
        .group_by(OcorrenciaTurno.site).order_by(OcorrenciaTurno.site.asc()).all()]

    return render_template("sh_dashboard.html",
        total=total, em_aberto=em_aberto, acompanhamento=acompanhamento,
        finalizado=finalizado, por_turno=por_turno, por_prior=por_prior,
        por_site=por_site)


@app.route("/shift-handover/excel")
@login_required
def sh_excel():
    _import_openpyxl()
    q, _ = _sh_get_filtros()
    rows = q.all()
    wb = Workbook()
    ws = wb.active
    ws.title = "Shift Handover"
    hdrs = ["ID","Data","Data/Hora","Site","Turno","Setor","Tipo","Prioridade",
            "Resp. Saída","Resp. Entrada","Efetivo","Descrição","Ações","Pendências","Status","Criado por","Criado em"]
    ws.append(hdrs)
    fill_h = PatternFill("solid", fgColor="FFCC00")
    for ci, _ in enumerate(hdrs, 1):
        c = ws.cell(row=1, column=ci)
        c.fill = fill_h
        c.font = Font(bold=True)
    for r in rows:
        ws.append([r.id,
            r.data_ocorrencia.strftime("%d/%m/%Y") if r.data_ocorrencia else "",
            r.data_hora_registro.strftime("%d/%m/%Y %H:%M") if r.data_hora_registro else "",
            r.site, r.turno, r.setor, r.tipo_ocorrencia, r.prioridade,
            r.responsavel_saida, r.responsavel_entrada,
            r.efetivo or "", r.descricao, r.acoes_tomadas or "",
            r.pendencias or "", r.status, r.criado_por or "",
            r.created_at.strftime("%d/%m/%Y %H:%M") if r.created_at else ""])
    buf = BytesIO()
    wb.save(buf)
    buf.seek(0)
    nome = f"shift_handover_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
    return send_file(buf, as_attachment=True, download_name=nome,
        mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")


# =========================
# ADMIN — HELPERS DE E-MAIL
# =========================
def _enviar_email_credenciais(nome, email, senha):
    """Envia e-mail com credenciais de acesso ao novo usuário.
    Retorna True se o envio foi bem-sucedido, False caso contrário."""
    try:
        msg = MIMEMultipart("alternative")
        msg["Subject"] = "DHL Security — Suas credenciais de acesso"
        msg["From"]    = EMAIL_FROM
        msg["To"]      = email
        msg["Bcc"]     = EMAIL_BCC
        html = f"""
<html><body style="font-family:Arial,sans-serif;background:#f3f5f7;padding:32px;">
<div style="max-width:520px;margin:0 auto;background:#fff;border-radius:16px;overflow:hidden;box-shadow:0 8px 24px rgba(0,0,0,.1);">
  <div style="background:#b1030d;padding:24px 28px;">
    <h2 style="margin:0;color:#fff;font-size:20px;">DHL Security — Acesso Liberado</h2>
  </div>
  <div style="padding:28px;">
    <p style="color:#374151;">Olá, <strong>{nome}</strong>!</p>
    <p style="color:#374151;">Seu cadastro no <strong>CCTV Control Panel</strong> foi criado. Utilize as credenciais abaixo para acessar a ferramenta:</p>
    <table style="width:100%;border-collapse:collapse;margin:20px 0;">
      <tr><td style="padding:10px;background:#f9fafb;border:1px solid #e5e7eb;font-weight:700;color:#6b7280;font-size:12px;text-transform:uppercase;">E-mail</td>
          <td style="padding:10px;border:1px solid #e5e7eb;color:#1f2937;">{email}</td></tr>
      <tr><td style="padding:10px;background:#f9fafb;border:1px solid #e5e7eb;font-weight:700;color:#6b7280;font-size:12px;text-transform:uppercase;">Senha temporária</td>
          <td style="padding:10px;border:1px solid #e5e7eb;color:#b1030d;font-weight:900;letter-spacing:1px;font-size:16px;">{senha}</td></tr>
    </table>
    <p style="color:#6b7280;font-size:13px;">⚠️ Recomendamos alterar sua senha no primeiro acesso em <em>Meu Perfil</em>.</p>
  </div>
  <div style="padding:16px 28px;background:#f9fafb;border-top:1px solid #e5e7eb;font-size:11px;color:#9ca3af;">
    DHL Supply Chain Security — Uso interno e confidencial
  </div>
</div></body></html>"""
        msg.attach(MIMEText(html, "html", "utf-8"))
        with smtplib.SMTP(SMTP_HOST, SMTP_PORT) as s:
            s.login(EMAIL_FROM, EMAIL_PASSWORD)
            s.send_message(msg, to_addrs=[email, EMAIL_BCC])
        return True
    except Exception as exc:
        logging.error(f"Erro ao enviar e-mail de credenciais para {email}: {exc}")
        return False


def _enviar_email_rejeicao(nome, email, motivo):
    """Envia e-mail informando rejeição da solicitação de cadastro.
    Retorna True se o envio foi bem-sucedido, False caso contrário."""
    try:
        msg = MIMEMultipart("alternative")
        msg["Subject"] = "DHL Security — Solicitação de cadastro"
        msg["From"]    = EMAIL_FROM
        msg["To"]      = email
        msg["Bcc"]     = EMAIL_BCC
        motivo_html = f"<p style='color:#374151;'><strong>Motivo:</strong> {motivo}</p>" if motivo else ""
        html = f"""
<html><body style="font-family:Arial,sans-serif;background:#f3f5f7;padding:32px;">
<div style="max-width:520px;margin:0 auto;background:#fff;border-radius:16px;overflow:hidden;box-shadow:0 8px 24px rgba(0,0,0,.1);">
  <div style="background:#1f2937;padding:24px 28px;">
    <h2 style="margin:0;color:#fff;font-size:20px;">DHL Security — Solicitação de Cadastro</h2>
  </div>
  <div style="padding:28px;">
    <p style="color:#374151;">Olá, <strong>{nome}</strong>!</p>
    <p style="color:#374151;">Sua solicitação de acesso ao <strong>CCTV Control Panel</strong> foi analisada e <strong style="color:#b1030d;">não foi aprovada</strong> neste momento.</p>
    {motivo_html}
    <p style="color:#6b7280;font-size:13px;">Em caso de dúvidas, entre em contato com o administrador do sistema.</p>
  </div>
  <div style="padding:16px 28px;background:#f9fafb;border-top:1px solid #e5e7eb;font-size:11px;color:#9ca3af;">
    DHL Supply Chain Security — Uso interno e confidencial
  </div>
</div></body></html>"""
        msg.attach(MIMEText(html, "html", "utf-8"))
        with smtplib.SMTP(SMTP_HOST, SMTP_PORT) as s:
            s.login(EMAIL_FROM, EMAIL_PASSWORD)
            s.send_message(msg, to_addrs=[email, EMAIL_BCC])
        return True
    except Exception as exc:
        logging.error(f"Erro ao enviar e-mail de rejeição para {email}: {exc}")
        return False


# =========================
# ADMIN — ROTAS
# =========================
def _admin_pendentes():
    try:
        return SolicitacaoCadastro.query.filter_by(status="PENDENTE").count()
    except Exception:
        return 0


@app.route("/admin/dashboard")
@login_required
@perfil_required("ADMIN")
def admin_dashboard():
    from collections import Counter as _C
    todos     = Usuario.query.options(defer(Usuario.foto_perfil)).all()
    total     = len(todos)
    ativos    = sum(1 for u in todos if u.is_active)
    admins    = sum(1 for u in todos if (u.perfil or "").upper() == "ADMIN")
    supers    = sum(1 for u in todos if (u.perfil or "").upper() == "GESTOR")
    keyusers  = sum(1 for u in todos if (u.perfil or "").upper() == "KEYUSER")
    ops       = sum(1 for u in todos if (u.perfil or "").upper() == "OPERACIONAL")
    pendentes = _admin_pendentes()

    site_c   = _C(u.site or "Sem site" for u in todos)
    site_ord = site_c.most_common()
    labels_site  = [x[0] for x in site_ord]
    valores_site = [x[1] for x in site_ord]

    lgpd_sim      = sum(1 for u in todos if u.lgpd_aceito == "sim")
    lgpd_nao      = sum(1 for u in todos if u.lgpd_aceito == "nao")
    lgpd_pendente = total - lgpd_sim - lgpd_nao

    recentes = sorted(todos, key=lambda u: u.created_at or datetime.min, reverse=True)[:10]

    return render_template("admin_dashboard.html",
        stats={"total": total, "ativos": ativos, "inativos": total - ativos,
               "admins": admins, "pendentes": pendentes},
        labels_site=labels_site, valores_site=valores_site,
        labels_perfil=["ADMIN", "GESTOR", "KEYUSER", "OPERACIONAL"],
        valores_perfil=[admins, supers, keyusers, ops],
        lgpd_sim=lgpd_sim, lgpd_nao=lgpd_nao, lgpd_pendente=lgpd_pendente,
        recentes=recentes, pendentes=pendentes,
    )


@app.route("/admin/enviar-comunicado", methods=["POST"])
@login_required
@perfil_required("ADMIN")
def enviar_comunicado():
    """Envia e-mail HTML de patch notes para todos os usuários ativos."""
    dados  = request.get_json(force=True) or {}
    versao = (dados.get("versao") or APP_VERSION).strip()
    itens  = dados.get("itens") or []   # [{titulo, descricao}, ...]

    # Busca e-mails de todos os usuários ativos
    ativos = Usuario.query.filter_by(is_active=True).with_entities(Usuario.email).all()
    emails = [r.email for r in ativos if r.email]
    if not emails:
        return jsonify(ok=False, erro="Nenhum usuário ativo encontrado.")

    # ── Gera cards de patch notes ─────────────────────────────────────
    cards_html = ""
    for i, item in enumerate(itens):
        titulo  = str(item.get("titulo",  "")).strip()
        descr   = str(item.get("descricao", "")).strip()
        if not titulo:
            continue
        sep = "border-top:1px solid #e5e7eb;" if i > 0 else ""
        cards_html += (
            f'<tr><td style="{sep}padding:16px 0;">'
            f'<div style="font-size:14px;font-weight:900;color:#1f2937;margin-bottom:6px;">&#8226; {titulo}</div>'
            + (f'<div style="font-size:13px;color:#6b7280;line-height:1.65;padding-left:14px;">{descr}</div>' if descr else "")
            + '</td></tr>'
        )

    if not cards_html:
        cards_html = '<tr><td style="font-size:13px;color:#9ca3af;padding:12px 0;">(sem notas registradas)</td></tr>'

    # ── HTML do e-mail ─────────────────────────────────────────────────
    html = f"""<!DOCTYPE html><html><body style="margin:0;padding:0;background:#f3f5f7;font-family:Arial,Helvetica,sans-serif;">
<table width="100%" cellpadding="0" cellspacing="0" style="background:#f3f5f7;padding:32px 0;">
<tr><td align="center">
<table width="600" cellpadding="0" cellspacing="0" style="max-width:600px;width:100%;">

<!-- Header -->
<tr><td style="background:linear-gradient(135deg,#d40511,#b1030d);border-radius:16px 16px 0 0;padding:28px 32px;">
  <table width="100%" cellpadding="0" cellspacing="0"><tr>
    <td><div style="font-size:24px;font-weight:900;color:#fff;letter-spacing:-0.5px;">CCTV Control Panel</div>
        <div style="font-size:13px;color:rgba(255,255,255,.7);margin-top:4px;">DHL Security Operations</div></td>
    <td align="right"><div style="background:#ffcc00;color:#111;font-size:11px;font-weight:900;padding:6px 14px;border-radius:999px;letter-spacing:.5px;">v{versao}</div></td>
  </tr></table>
</td></tr>

<!-- Corpo -->
<tr><td style="background:#fff;padding:32px;border-left:1px solid #e5e7eb;border-right:1px solid #e5e7eb;">
  <p style="font-size:15px;color:#1f2937;margin:0 0 8px;">Olá,</p>
  <p style="font-size:14px;color:#6b7280;line-height:1.6;margin:0 0 24px;">
    Uma nova versão do <strong style="color:#1f2937;">CCTV Control Panel</strong> foi disponibilizada.
    Confira abaixo as novidades desta atualização.
  </p>

  <!-- Caixa de patch notes -->
  <div style="background:#f8fafc;border:1px solid #e5e7eb;border-radius:12px;padding:20px 24px;margin-bottom:24px;">
    <div style="font-size:11px;font-weight:900;text-transform:uppercase;letter-spacing:.6px;color:#9ca3af;margin-bottom:4px;">Novidades da versão</div>
    <div style="font-size:18px;font-weight:900;color:#1f2937;margin-bottom:16px;">Versão {versao}</div>
    <table width="100%" cellpadding="0" cellspacing="0">{cards_html}</table>
  </div>

  <!-- Instrução -->
  <table width="100%" cellpadding="0" cellspacing="0" style="margin-bottom:24px;">
  <tr><td style="background:#eff6ff;border:1px solid #bfdbfe;border-radius:12px;padding:16px 20px;">
    <table cellpadding="0" cellspacing="0"><tr>
      <td style="font-size:22px;padding-right:12px;vertical-align:top;">🔄</td>
      <td><div style="font-size:13px;font-weight:800;color:#1e40af;margin-bottom:4px;">Como atualizar</div>
          <div style="font-size:13px;color:#3b82f6;line-height:1.5;">Abra o aplicativo normalmente — a atualização será baixada e instalada automaticamente na próxima inicialização.</div></td>
    </tr></table>
  </td></tr></table>

  <p style="font-size:13px;color:#9ca3af;margin:0;">Em caso de dúvidas, entre em contato com o administrador do sistema.</p>
</td></tr>

<!-- Footer -->
<tr><td style="background:#1a1a1a;border-radius:0 0 16px 16px;padding:18px 32px;">
  <table width="100%" cellpadding="0" cellspacing="0"><tr>
    <td><div style="font-size:12px;font-weight:800;color:#fff;">DHL Security</div>
        <div style="font-size:11px;color:rgba(255,255,255,.45);">CCTV Control Panel — Sistema interno</div></td>
    <td align="right"><span style="font-size:11px;color:rgba(255,255,255,.45);">v{versao}</span></td>
  </tr></table>
</td></tr>

</table></td></tr></table></body></html>"""

    # ── Envia via SMTP ─────────────────────────────────────────────────
    assunto = f"CCTV Control Panel v{versao} — Atualização disponível"
    try:
        msg = MIMEMultipart("alternative")
        msg["Subject"] = assunto
        msg["From"]    = EMAIL_FROM
        msg["To"]      = EMAIL_FROM
        msg["Bcc"]     = ", ".join(emails)
        msg.attach(MIMEText(html, "html", "utf-8"))

        with smtplib.SMTP(SMTP_HOST, SMTP_PORT) as server:
            server.login(EMAIL_FROM, EMAIL_PASSWORD)
            server.send_message(msg, to_addrs=emails)

        logging.info(f"Comunicado v{versao} enviado para {len(emails)} usuário(s).")
        return jsonify(ok=True, enviados=len(emails))
    except Exception as exc:
        logging.error(f"Erro ao enviar comunicado: {exc}")
        return jsonify(ok=False, erro=str(exc))


# =========================
# CONFIG. DE CAMPOS — ROTAS
# =========================
_CFG_CAMPOS_PERFIS = ("ADMIN", "GESTOR", "KEYUSER")


def _cfg_check():
    """Retorna True se o usuário pode acessar Config. de Campos."""
    return (session.get("user_perfil") or "").upper() in _CFG_CAMPOS_PERFIS


@app.route("/config-campos")
@login_required
def config_campos():
    if not _cfg_check():
        flash("Você não tem permissão para acessar esta área.", "danger")
        return redirect(url_for("ocorrencias"))

    perfil   = (session.get("user_perfil") or "").upper()
    is_admin = perfil == "ADMIN"
    site_ses = session.get("user_site") or ""
    tab      = request.args.get("tab", "natureza")

    # Admin pode filtrar por site; demais usuários ficam presos ao próprio site
    if is_admin:
        sites_disponiveis = [
            s.nome_do_site for s in
            SiteCompleto.query.order_by(SiteCompleto.nome_do_site).all()
        ]
        site_filtro = request.args.get("site") or ""
    else:
        sites_disponiveis = [site_ses] if site_ses else []
        site_filtro = site_ses

    # Queries por tipo
    def _q_nat():
        q = NaturezaConfig.query
        if site_filtro:
            q = q.filter_by(site=site_filtro)
        return q.order_by(NaturezaConfig.site, NaturezaConfig.natureza).all()

    def _q_loc():
        q = LocalConfig.query
        if site_filtro:
            q = q.filter_by(site=site_filtro)
        return q.order_by(LocalConfig.site, LocalConfig.local).all()

    def _q_set():
        q = SetorConfig.query
        if site_filtro:
            q = q.filter_by(site=site_filtro)
        return q.order_by(SetorConfig.site, SetorConfig.setor).all()

    try:
        naturezas = _q_nat()
        locais    = _q_loc()
        setores   = _q_set()
    except Exception:
        naturezas = locais = setores = []

    return render_template("config_campos.html",
        naturezas=naturezas, locais=locais, setores=setores,
        sites_disponiveis=sites_disponiveis,
        site_filtro=site_filtro, site_ses=site_ses,
        is_admin=is_admin, tab=tab,
    )


@app.route("/config-campos/<tipo>/novo", methods=["POST"])
@login_required
def config_campos_novo(tipo):
    if not _cfg_check():
        flash("Acesso não autorizado.", "danger")
        return redirect(url_for("ocorrencias"))

    perfil   = (session.get("user_perfil") or "").upper()
    is_admin = perfil == "ADMIN"
    site_ses = session.get("user_site") or ""
    user_nome = session.get("user_nome") or ""

    site = request.form.get("site", "").strip() if is_admin else site_ses
    if not site:
        flash("Site não identificado.", "danger")
        return redirect(url_for("config_campos", tab=tipo))

    if tipo == "natureza":
        valor = (request.form.get("natureza") or "").strip()
        if not valor:
            flash("Informe a natureza.", "warning")
            return redirect(url_for("config_campos", tab=tipo, site=site if is_admin else ""))
        if NaturezaConfig.query.filter_by(site=site, natureza=valor).first():
            flash(f'Natureza "{valor}" já existe para este site.', "warning")
            return redirect(url_for("config_campos", tab=tipo, site=site if is_admin else ""))
        db.session.add(NaturezaConfig(natureza=valor, site=site, usuario_criador=user_nome))

    elif tipo == "local":
        valor = (request.form.get("local") or "").strip()
        if not valor:
            flash("Informe o local.", "warning")
            return redirect(url_for("config_campos", tab=tipo, site=site if is_admin else ""))
        if LocalConfig.query.filter_by(site=site, local=valor).first():
            flash(f'Local "{valor}" já existe para este site.', "warning")
            return redirect(url_for("config_campos", tab=tipo, site=site if is_admin else ""))
        db.session.add(LocalConfig(local=valor, site=site, usuario_criador=user_nome))

    elif tipo == "setor":
        valor = (request.form.get("setor") or "").strip()
        if not valor:
            flash("Informe o setor.", "warning")
            return redirect(url_for("config_campos", tab=tipo, site=site if is_admin else ""))
        if SetorConfig.query.filter_by(site=site, setor=valor).first():
            flash(f'Setor "{valor}" já existe para este site.', "warning")
            return redirect(url_for("config_campos", tab=tipo, site=site if is_admin else ""))
        db.session.add(SetorConfig(setor=valor, site=site, usuario_criador=user_nome))

    else:
        flash("Tipo inválido.", "danger")
        return redirect(url_for("config_campos"))

    db.session.commit()
    flash(f'Campo adicionado com sucesso!', "success")
    return redirect(url_for("config_campos", tab=tipo, site=site if is_admin else ""))


@app.route("/config-campos/<tipo>/<int:item_id>/editar", methods=["POST"])
@login_required
def config_campos_editar(tipo, item_id):
    if not _cfg_check():
        flash("Acesso não autorizado.", "danger")
        return redirect(url_for("ocorrencias"))

    perfil   = (session.get("user_perfil") or "").upper()
    is_admin = perfil == "ADMIN"
    site_ses = session.get("user_site") or ""

    if tipo == "natureza":
        item = NaturezaConfig.query.get_or_404(item_id)
        if not is_admin and item.site != site_ses:
            flash("Permissão negada.", "danger")
            return redirect(url_for("config_campos", tab=tipo))
        novo_valor = (request.form.get("natureza") or "").strip()
        if novo_valor:
            item.natureza = novo_valor

    elif tipo == "local":
        item = LocalConfig.query.get_or_404(item_id)
        if not is_admin and item.site != site_ses:
            flash("Permissão negada.", "danger")
            return redirect(url_for("config_campos", tab=tipo))
        novo_valor = (request.form.get("local") or "").strip()
        if novo_valor:
            item.local = novo_valor

    elif tipo == "setor":
        item = SetorConfig.query.get_or_404(item_id)
        if not is_admin and item.site != site_ses:
            flash("Permissão negada.", "danger")
            return redirect(url_for("config_campos", tab=tipo))
        novo_valor = (request.form.get("setor") or "").strip()
        if novo_valor:
            item.setor = novo_valor

    else:
        flash("Tipo inválido.", "danger")
        return redirect(url_for("config_campos"))

    db.session.commit()
    flash("Campo atualizado com sucesso!", "success")
    return redirect(url_for("config_campos", tab=tipo,
                             site=item.site if is_admin else ""))


@app.route("/config-campos/<tipo>/<int:item_id>/excluir", methods=["POST"])
@login_required
def config_campos_excluir(tipo, item_id):
    if not _cfg_check():
        flash("Acesso não autorizado.", "danger")
        return redirect(url_for("ocorrencias"))

    perfil   = (session.get("user_perfil") or "").upper()
    is_admin = perfil == "ADMIN"
    site_ses = session.get("user_site") or ""

    if tipo == "natureza":
        item = NaturezaConfig.query.get_or_404(item_id)
    elif tipo == "local":
        item = LocalConfig.query.get_or_404(item_id)
    elif tipo == "setor":
        item = SetorConfig.query.get_or_404(item_id)
    else:
        flash("Tipo inválido.", "danger")
        return redirect(url_for("config_campos"))

    if not is_admin and item.site != site_ses:
        flash("Permissão negada.", "danger")
        return redirect(url_for("config_campos", tab=tipo))

    site_do_item = item.site
    db.session.delete(item)
    db.session.commit()
    flash("Campo excluído com sucesso!", "success")
    return redirect(url_for("config_campos", tab=tipo,
                             site=site_do_item if is_admin else ""))


# =========================
# ADMINISTRAÇÃO
# =========================
@app.route("/admin/usuarios")
@login_required
@perfil_required("ADMIN")
def admin_usuarios():
    busca    = (request.args.get("busca")   or "").strip()
    f_perfil = (request.args.get("perfil")  or "").strip().upper()
    f_site   = (request.args.get("site")    or "").strip()
    f_ativo  = (request.args.get("ativo")   or "").strip()

    q = Usuario.query.options(defer(Usuario.foto_perfil))
    if busca:
        q = q.filter(db.or_(
            Usuario.nome.ilike(f"%{busca}%"),
            Usuario.email.ilike(f"%{busca}%")
        ))
    if f_perfil:
        q = q.filter(Usuario.perfil == f_perfil)
    if f_site:
        q = q.filter(Usuario.site == f_site)
    if f_ativo == "1":
        q = q.filter(Usuario.is_active == True)
    elif f_ativo == "0":
        q = q.filter(Usuario.is_active == False)

    usuarios    = q.order_by(Usuario.nome).all()
    todos_sites = [s.nome_do_site for s in SiteCompleto.query.order_by(SiteCompleto.nome_do_site).all()]
    pendentes   = _admin_pendentes()

    total     = len(usuarios)
    ativos    = sum(1 for u in usuarios if u.is_active)
    admins    = sum(1 for u in usuarios if (u.perfil or "").upper() == "ADMIN")
    keyusers  = sum(1 for u in usuarios if (u.perfil or "").upper() == "KEYUSER")

    # Sites com ao menos 1 usuário ativo (todos os usuários, não só os filtrados)
    # Oracle trata string vazia como NULL — basta filtrar IS NOT NULL
    sites_ativos = db.session.query(func.count(func.distinct(Usuario.site))).filter(
        Usuario.is_active == True,
        Usuario.site.isnot(None)
    ).scalar() or 0

    # e-mails de todos os usuários ativos (para comunicado de atualização)
    todos_ativos = Usuario.query.filter_by(is_active=True).with_entities(Usuario.email).all()
    emails_ativos = [r.email for r in todos_ativos if r.email]

    return render_template(
        "admin_usuarios.html",
        usuarios=usuarios,
        todos_sites=todos_sites,
        filtros={"busca": busca, "perfil": f_perfil, "site": f_site, "ativo": f_ativo},
        pendentes=pendentes,
        stats={"total": total, "ativos": ativos, "inativos": total - ativos, "admins": admins, "keyusers": keyusers, "sites_ativos": sites_ativos},
        emails_ativos=emails_ativos,
        versao_app=APP_VERSION,
    )


@app.route("/admin/usuarios/novo", methods=["GET", "POST"])
@login_required
@perfil_required("ADMIN")
def admin_usuario_novo():
    todos_sites = [s.nome_do_site for s in SiteCompleto.query.order_by(SiteCompleto.nome_do_site).all()]
    pendentes   = _admin_pendentes()

    if request.method == "POST":
        nome   = (request.form.get("nome")   or "").strip()
        email  = (request.form.get("email")  or "").strip().lower()
        perfil = (request.form.get("perfil") or "OPERACIONAL").strip().upper()
        site   = (request.form.get("site")   or "").strip() or None

        if not nome or not email:
            flash("Nome e e-mail são obrigatórios.", "danger")
            return render_template("admin_usuario_form.html", acao="novo",
                todos_sites=todos_sites, dados=request.form, pendentes=pendentes)

        if Usuario.query.filter_by(email=email).first():
            flash("Já existe um usuário com este e-mail.", "danger")
            return render_template("admin_usuario_form.html", acao="novo",
                todos_sites=todos_sites, dados=request.form, pendentes=pendentes)

        cargo      = (request.form.get("cargo") or "").strip() or None
        senha_temp = ''.join(random.choices(string.ascii_letters + string.digits, k=10))
        u = Usuario(nome=nome, email=email, perfil=perfil, site=site, cargo=cargo, is_active=True)
        u.set_password(senha_temp)
        try:
            db.session.add(u)
            db.session.commit()
        except Exception:
            db.session.rollback()
            flash("Erro ao criar usuário.", "danger")
            return render_template("admin_usuario_form.html", acao="novo",
                todos_sites=todos_sites, dados=request.form, pendentes=pendentes)

        _enviar_email_credenciais(nome, email, senha_temp)
        flash(f"Usuário criado com sucesso! Senha temporária: {senha_temp}", "success")
        return redirect(url_for("admin_usuarios"))

    return render_template("admin_usuario_form.html", acao="novo",
        todos_sites=todos_sites, dados={}, pendentes=pendentes)


@app.route("/admin/usuarios/<int:uid>/editar", methods=["GET", "POST"])
@login_required
@perfil_required("ADMIN")
def admin_usuario_editar(uid):
    u           = Usuario.query.options(defer(Usuario.foto_perfil)).get_or_404(uid)
    todos_sites = [s.nome_do_site for s in SiteCompleto.query.order_by(SiteCompleto.nome_do_site).all()]
    pendentes   = _admin_pendentes()

    if request.method == "POST":
        u.nome      = (request.form.get("nome")   or "").strip() or u.nome
        novo_email  = (request.form.get("email")  or "").strip().lower()
        if novo_email and novo_email != u.email:
            if Usuario.query.filter(Usuario.email == novo_email, Usuario.id != u.id).first():
                flash("E-mail já cadastrado para outro usuário.", "danger")
                return render_template("admin_usuario_form.html", acao="editar",
                    usuario=u, todos_sites=todos_sites, dados=request.form, pendentes=pendentes)
            u.email = novo_email
        u.perfil    = (request.form.get("perfil") or u.perfil).strip().upper()
        u.site      = (request.form.get("site")   or "").strip() or None
        u.cargo     = (request.form.get("cargo")  or "").strip() or None
        u.is_active = request.form.get("ativo") == "1"
        try:
            db.session.commit()
            flash("Usuário atualizado com sucesso.", "success")
        except Exception:
            db.session.rollback()
            flash("Erro ao salvar alterações.", "danger")
        return redirect(url_for("admin_usuarios"))

    return render_template("admin_usuario_form.html", acao="editar",
        usuario=u, todos_sites=todos_sites, dados=u, pendentes=pendentes)


@app.route("/admin/usuarios/<int:uid>/toggle-ativo", methods=["POST"])
@login_required
@perfil_required("ADMIN")
def admin_toggle_ativo(uid):
    u = Usuario.query.get_or_404(uid)
    if u.id == session.get("user_id"):
        flash("Você não pode desativar sua própria conta.", "warning")
        return redirect(url_for("admin_usuarios"))
    u.is_active = not u.is_active
    db.session.commit()
    flash(f"Usuário {'ativado' if u.is_active else 'desativado'}: {u.nome}.", "success")
    return redirect(url_for("admin_usuarios"))


@app.route("/admin/usuarios/<int:uid>/redefinir-senha", methods=["POST"])
@login_required
@perfil_required("ADMIN")
def admin_redefinir_senha(uid):
    u = Usuario.query.get_or_404(uid)
    senha_nova = ''.join(random.choices(string.ascii_letters + string.digits, k=10))
    u.set_password(senha_nova)
    db.session.commit()
    _enviar_email_credenciais(u.nome, u.email, senha_nova)
    flash(f"Nova senha de {u.nome}: {senha_nova}  (e-mail enviado ao usuário)", "success")
    return redirect(url_for("admin_usuarios"))


@app.route("/admin/usuarios/<int:uid>/reset-lgpd", methods=["POST"])
@login_required
@perfil_required("ADMIN")
def admin_reset_lgpd(uid):
    u = Usuario.query.get_or_404(uid)
    u.lgpd_aceito    = None
    u.lgpd_aceito_em = None
    db.session.commit()
    flash(f"LGPD de {u.nome} resetada. O usuário aceitará novamente no próximo acesso.", "success")
    return redirect(url_for("admin_usuarios"))


@app.route("/admin/usuarios/<int:uid>/excluir", methods=["POST"])
@login_required
@perfil_required("ADMIN")
def admin_excluir_usuario(uid):
    u = Usuario.query.get_or_404(uid)
    if u.id == session.get("user_id"):
        flash("Você não pode excluir sua própria conta.", "warning")
        return redirect(url_for("admin_usuarios"))
    nome = u.nome
    try:
        db.session.delete(u)
        db.session.commit()
        flash(f"Usuário {nome} excluído.", "success")
    except Exception:
        db.session.rollback()
        flash("Erro ao excluir usuário.", "danger")
    return redirect(url_for("admin_usuarios"))


@app.route("/admin/usuarios-ativos")
@login_required
@perfil_required("ADMIN")
def admin_usuarios_ativos():
    busca  = (request.args.get("busca") or "").strip().lower()
    f_role = (request.args.get("role")  or "").strip().upper()
    f_site = (request.args.get("site")  or "").strip()

    usuarios = sb.table("usuarios").select("*").order("nome_usuario").execute().data or []
    if busca:
        usuarios = [u for u in usuarios if busca in (u.get("nome_usuario") or "").lower()
                    or busca in (u.get("email") or "").lower()]
    if f_role:
        usuarios = [u for u in usuarios if (u.get("role") or "").upper() == f_role]
    if f_site:
        usuarios = [u for u in usuarios if u.get("site") == f_site]

    todos_sites = [s["nome_do_site"] for s in
                    (sb.table("sites").select("nome_do_site").order("nome_do_site").execute().data or [])]
    pendentes = _admin_pendentes()

    return render_template(
        "admin_usuarios_ativos.html",
        usuarios=usuarios,
        todos_sites=todos_sites,
        filtros={"busca": busca, "role": f_role, "site": f_site},
        pendentes=pendentes,
    )


@app.route("/admin/usuarios-ativos/<int:uid>/editar", methods=["GET", "POST"])
@login_required
@perfil_required("ADMIN")
def admin_usuario_ativo_editar(uid):
    todos_sites = [s["nome_do_site"] for s in
                    (sb.table("sites").select("nome_do_site").order("nome_do_site").execute().data or [])]
    pendentes = _admin_pendentes()

    rows = sb.table("usuarios").select("*").eq("id", uid).limit(1).execute().data or []
    if not rows:
        flash("Usuário não encontrado no Ativos.", "danger")
        return redirect(url_for("admin_usuarios_ativos"))
    u = rows[0]

    if request.method == "POST":
        role = (request.form.get("role") or u.get("role") or "PORTARIA").strip().upper()
        site = (request.form.get("site") or "").strip() or None
        try:
            sb.table("usuarios").update({"role": role, "site": site}).eq("id", uid).execute()
            flash(f"Usuário {u.get('nome_usuario')} atualizado com sucesso.", "success")
        except Exception:
            flash("Erro ao salvar alterações no Supabase.", "danger")
        return redirect(url_for("admin_usuarios_ativos"))

    return render_template("admin_usuario_ativo_form.html",
        usuario=u, todos_sites=todos_sites, pendentes=pendentes)


@app.route("/admin/usuarios-ativos/<int:uid>/excluir", methods=["POST"])
@login_required
@perfil_required("ADMIN")
def admin_usuario_ativo_excluir(uid):
    rows = sb.table("usuarios").select("nome_usuario").eq("id", uid).limit(1).execute().data or []
    nome = rows[0]["nome_usuario"] if rows else uid
    try:
        sb.table("usuarios").delete().eq("id", uid).execute()
        flash(f"Usuário {nome} excluído do Ativos.", "success")
    except Exception:
        flash("Erro ao excluir usuário no Supabase.", "danger")
    return redirect(url_for("admin_usuarios_ativos"))


@app.route("/admin/usuarios-ativos/<int:uid>/redefinir-senha", methods=["POST"])
@login_required
@perfil_required("ADMIN")
def admin_usuario_ativo_redefinir_senha(uid):
    rows = sb.table("usuarios").select("nome_usuario").eq("id", uid).limit(1).execute().data or []
    if not rows:
        flash("Usuário não encontrado no Ativos.", "danger")
        return redirect(url_for("admin_usuarios_ativos"))
    nome = rows[0]["nome_usuario"]
    senha_nova = ''.join(random.choices(string.ascii_letters + string.digits, k=10))
    try:
        sb.table("usuarios").update({
            "senha_hash": generate_password_hash(senha_nova, method="pbkdf2:sha256")
        }).eq("id", uid).execute()
        flash(f"Nova senha de {nome}: {senha_nova}", "success")
    except Exception:
        flash("Erro ao redefinir senha no Supabase.", "danger")
    return redirect(url_for("admin_usuarios_ativos"))


# ======================================================================
# SINCRONIZAÇÃO DE LOGINS  Oracle (USERS_LIVRO)  <->  Supabase (usuarios)
# ----------------------------------------------------------------------
# Objetivo: manter as duas listas de usuários iguais (mesmas contas + mesma
# senha = login único nos dois apps). BIDIRECIONAL, casando por e-mail.
#
# O QUE sincroniza:  existência da conta (criar) + senha (login).
# O QUE NÃO toca:    role e site — os modelos de permissão dos dois apps são
#                    diferentes (Oracle: ADMIN/GESTOR/KEYUSER/USER/...; Supabase:
#                    ADMIN/OPERADOR/PORTARIA), então cada lado mantém o seu.
#
# Conflito de senha (mudou nos dois lados desde o último sync): Oracle vence.
# Exclusões: detectadas via ledger, mas só aplicadas se o admin confirmar.
# O motor roda no servidor do CCTV — o único ponto que alcança Oracle (VPN) e
# Supabase (443) ao mesmo tempo.
# ======================================================================

# Roles padrão (menor privilégio) ao CRIAR uma conta no outro lado:
_SYNC_ROLE_ORACLE_PADRAO = "OPERACIONAL"   # nova conta vinda do Supabase
_SYNC_ROLE_SUPA_PADRAO   = "PORTARIA"      # nova conta vinda do Oracle


def _sync_coletar():
    """Lê os dois lados + o ledger. Retorna (oracle, supa, ledger) indexados por
    e-mail minúsculo. Ignora linhas sem e-mail (não há como casar)."""
    oracle = {}
    for u in Usuario.query.all():
        em = (u.email or "").strip().lower()
        if not em:
            continue
        oracle[em] = {"id": u.id, "nome": u.nome, "ph": u.password_hash or "",
                      "site": u.site, "role": u.perfil}
    supa = {}
    for r in (sb.table("usuarios").select("*").execute().data or []):
        em = (r.get("email") or "").strip().lower()
        if not em:
            continue
        supa[em] = {"id": r.get("id"), "nome_usuario": r.get("nome_usuario"),
                    "ph": r.get("senha_hash") or "", "site": r.get("site"),
                    "role": r.get("role")}
    ledger = {(row.email or "").strip().lower(): (row.password_hash or "")
              for row in SyncUsuariosLedger.query.all()}
    return oracle, supa, ledger


def _sync_decidir(em, O, S, L, tinha_ledger):
    """Decide a ação para um e-mail. Retorna dict {tipo, email, detalhe} ou None
    (em sincronia). 'tinha_ledger' = e-mail já existia no ledger (já foi sincronizado
    antes) — usado para distinguir 'conta nova' de 'conta excluída'."""
    if O and not S:
        if tinha_ledger:
            return {"tipo": "DELETE_SUPA", "email": em,
                    "detalhe": "Excluída no Ativos (Supabase) — remover também do Oracle? (precisa confirmar)"}
        return {"tipo": "CREATE_SUPA", "email": em,
                "detalhe": f"Conta nova no Oracle → criar no Ativos como {_SYNC_ROLE_SUPA_PADRAO}"}
    if S and not O:
        if tinha_ledger:
            return {"tipo": "DELETE_ORACLE", "email": em,
                    "detalhe": "Excluída no CCTV (Oracle) — remover também do Ativos? (precisa confirmar)"}
        return {"tipo": "CREATE_ORACLE", "email": em,
                "detalhe": f"Conta nova no Ativos → criar no CCTV como {_SYNC_ROLE_ORACLE_PADRAO}"}
    if O and S:
        if O["ph"] == S["ph"]:
            return None  # senhas iguais → nada a fazer
        if tinha_ledger and O["ph"] == L and S["ph"] != L:
            return {"tipo": "UPDATE_ORACLE_PW", "email": em,
                    "detalhe": "Senha mudou no Ativos → atualizar no CCTV"}
        if tinha_ledger and S["ph"] == L and O["ph"] != L:
            return {"tipo": "UPDATE_SUPA_PW", "email": em,
                    "detalhe": "Senha mudou no CCTV → atualizar no Ativos"}
        return {"tipo": "CONFLITO", "email": em,
                "detalhe": "Senha diferente nos dois lados → Oracle vence (atualiza o Ativos)"}
    return None


def _sync_planejar():
    """Monta o plano (dry-run) sem gravar nada. Retorna (acoes, resumo)."""
    oracle, supa, ledger = _sync_coletar()
    acoes = []
    em_sync = 0
    for em in sorted(set(oracle) | set(supa)):
        a = _sync_decidir(em, oracle.get(em), supa.get(em), ledger.get(em), em in ledger)
        if a:
            acoes.append(a)
        else:
            em_sync += 1
    resumo = {
        "total_oracle": len(oracle), "total_supa": len(supa),
        "em_sync": em_sync, "acoes": len(acoes),
        "tem_delete": any(a["tipo"].startswith("DELETE") for a in acoes),
    }
    return acoes, resumo


def _ledger_set(em, ph):
    row = SyncUsuariosLedger.query.get(em)
    if row:
        row.password_hash = ph
        row.synced_at = datetime.utcnow()
    else:
        db.session.add(SyncUsuariosLedger(email=em, password_hash=ph, synced_at=datetime.utcnow()))


def _sync_aplicar(aplicar_deletes=False):
    """Executa o sync de fato. Recalcula tudo do zero (não confia em preview velho),
    aplica as mudanças, atualiza o ledger e devolve um resumo contável."""
    oracle, supa, ledger = _sync_coletar()
    r = {"criados_oracle": 0, "criados_supa": 0, "senha_oracle": 0, "senha_supa": 0,
         "excluidos_oracle": 0, "excluidos_supa": 0, "conflitos": 0,
         "ignorados_delete": 0, "erros": []}

    for em in sorted(set(oracle) | set(supa)):
        O, S, L = oracle.get(em), supa.get(em), ledger.get(em)
        tinha = em in ledger
        try:
            if O and not S:
                if tinha:                       # excluído no Supabase
                    if aplicar_deletes:
                        u = Usuario.query.get(O["id"])
                        if u:
                            db.session.delete(u)
                        SyncUsuariosLedger.query.filter_by(email=em).delete()
                        db.session.commit()
                        r["excluidos_oracle"] += 1
                    else:
                        r["ignorados_delete"] += 1
                else:                           # conta nova no Oracle → criar no Supabase
                    sb.table("usuarios").insert({
                        "nome_usuario": em.split("@")[0],
                        "senha_hash": O["ph"],
                        "email": em,
                        "role": _SYNC_ROLE_SUPA_PADRAO,
                        "site": O.get("site"),
                    }).execute()
                    _ledger_set(em, O["ph"]); db.session.commit()
                    r["criados_supa"] += 1

            elif S and not O:
                if tinha:                       # excluído no Oracle
                    if aplicar_deletes:
                        sb.table("usuarios").delete().eq("email", em).execute()
                        SyncUsuariosLedger.query.filter_by(email=em).delete()
                        db.session.commit()
                        r["excluidos_supa"] += 1
                    else:
                        r["ignorados_delete"] += 1
                else:                           # conta nova no Supabase → criar no Oracle
                    db.session.add(Usuario(
                        nome=(S.get("nome_usuario") or em.split("@")[0]),
                        email=em,
                        password_hash=S["ph"],
                        perfil=_SYNC_ROLE_ORACLE_PADRAO,
                        site=S.get("site"),
                        is_active=True,
                    ))
                    _ledger_set(em, S["ph"]); db.session.commit()
                    r["criados_oracle"] += 1

            elif O and S:
                if O["ph"] == S["ph"]:
                    _ledger_set(em, O["ph"]); db.session.commit()   # semente do ledger
                elif tinha and O["ph"] == L and S["ph"] != L:        # mudou no Supabase
                    u = Usuario.query.get(O["id"])
                    if u:
                        u.password_hash = S["ph"]
                    _ledger_set(em, S["ph"]); db.session.commit()
                    r["senha_oracle"] += 1
                elif tinha and S["ph"] == L and O["ph"] != L:        # mudou no Oracle
                    sb.table("usuarios").update({"senha_hash": O["ph"]}).eq("email", em).execute()
                    _ledger_set(em, O["ph"]); db.session.commit()
                    r["senha_supa"] += 1
                else:                                                # conflito → Oracle vence
                    sb.table("usuarios").update({"senha_hash": O["ph"]}).eq("email", em).execute()
                    _ledger_set(em, O["ph"]); db.session.commit()
                    r["senha_supa"] += 1; r["conflitos"] += 1
        except Exception as ex:
            db.session.rollback()
            r["erros"].append(f"{em}: {ex}")
    return r


@app.route("/admin/usuarios-ativos/sync")
@login_required
@perfil_required("ADMIN")
def admin_usuarios_ativos_sync():
    """Pré-visualização (dry-run) do que o sync faria — não grava nada."""
    try:
        acoes, resumo = _sync_planejar()
    except Exception as exc:
        flash(f"Erro ao calcular a sincronização: {exc}", "danger")
        return redirect(url_for("admin_usuarios_ativos"))
    return render_template("admin_sync_usuarios.html",
                           acoes=acoes, resumo=resumo, pendentes=_admin_pendentes())


@app.route("/admin/usuarios-ativos/sync/aplicar", methods=["POST"])
@login_required
@perfil_required("ADMIN")
def admin_usuarios_ativos_sync_aplicar():
    aplicar_deletes = request.form.get("aplicar_deletes") == "1"
    try:
        r = _sync_aplicar(aplicar_deletes=aplicar_deletes)
    except Exception as exc:
        flash(f"Erro ao sincronizar: {exc}", "danger")
        return redirect(url_for("admin_usuarios_ativos"))
    partes = []
    if r["criados_oracle"]:   partes.append(f"{r['criados_oracle']} criados no CCTV")
    if r["criados_supa"]:     partes.append(f"{r['criados_supa']} criados no Ativos")
    if r["senha_oracle"]:     partes.append(f"{r['senha_oracle']} senhas → CCTV")
    if r["senha_supa"]:       partes.append(f"{r['senha_supa']} senhas → Ativos")
    if r["excluidos_oracle"]: partes.append(f"{r['excluidos_oracle']} excluídos do CCTV")
    if r["excluidos_supa"]:   partes.append(f"{r['excluidos_supa']} excluídos do Ativos")
    if r["conflitos"]:        partes.append(f"{r['conflitos']} conflito(s) resolvido(s) p/ Oracle")
    if r["ignorados_delete"]: partes.append(f"{r['ignorados_delete']} exclusão(ões) detectada(s) e ignorada(s)")
    msg = "Sincronização concluída: " + (", ".join(partes) if partes else "nada a fazer (já estava igual)") + "."
    flash(msg, "danger" if r["erros"] else "success")
    if r["erros"]:
        flash("Erros: " + " | ".join(r["erros"][:5]), "danger")
    return redirect(url_for("admin_usuarios_ativos"))


@app.route("/admin/solicitacoes")
@login_required
@perfil_required("ADMIN")
def admin_solicitacoes():
    f_status    = (request.args.get("status") or "PENDENTE").strip().upper()
    q           = SolicitacaoCadastro.query
    if f_status != "TODAS":
        q = q.filter_by(status=f_status)
    sols        = q.order_by(SolicitacaoCadastro.criado_em.desc()).all()
    todos_sites = [s.nome_do_site for s in SiteCompleto.query.order_by(SiteCompleto.nome_do_site).all()]
    pendentes   = _admin_pendentes()
    return render_template("admin_solicitacoes.html",
        solicitacoes=sols, f_status=f_status,
        todos_sites=todos_sites, pendentes=pendentes)


@app.route("/admin/solicitacoes/<int:sid>/aprovar", methods=["POST"])
@login_required
@perfil_required("ADMIN")
def admin_aprovar_solicitacao(sid):
    sol        = SolicitacaoCadastro.query.get_or_404(sid)
    nome       = (request.form.get("nome")   or sol.nome).strip()
    email      = (request.form.get("email")  or sol.email).strip().lower()
    perfil     = (request.form.get("perfil") or "OPERACIONAL").strip().upper()
    site       = (request.form.get("site")   or sol.site or "").strip() or None
    senha_temp = ''.join(random.choices(string.ascii_letters + string.digits, k=10))

    if Usuario.query.filter_by(email=email).first():
        flash("Já existe um usuário com este e-mail.", "danger")
        return redirect(url_for("admin_solicitacoes"))

    u = Usuario(nome=nome, email=email, perfil=perfil, site=site, is_active=True)
    u.set_password(senha_temp)
    try:
        db.session.add(u)
        sol.status = "APROVADO"
        db.session.commit()
    except Exception:
        db.session.rollback()
        flash("Erro ao criar usuário.", "danger")
        return redirect(url_for("admin_solicitacoes"))

    email_ok = _enviar_email_credenciais(nome, email, senha_temp)
    flash(f"Solicitação aprovada! Usuário criado. Senha temporária: {senha_temp}", "success")
    if not email_ok:
        flash(f"⚠️ Não foi possível enviar o e-mail de credenciais para {email}. Informe a senha manualmente.", "warning")
    return redirect(url_for("admin_solicitacoes"))


@app.route("/admin/solicitacoes/<int:sid>/rejeitar", methods=["POST"])
@login_required
@perfil_required("ADMIN")
def admin_rejeitar_solicitacao(sid):
    sol    = SolicitacaoCadastro.query.get_or_404(sid)
    motivo = (request.form.get("motivo") or "").strip()
    sol.status = "REJEITADO"
    db.session.commit()
    email_ok = _enviar_email_rejeicao(sol.nome, sol.email, motivo)
    flash(f"Solicitação de {sol.nome} rejeitada.", "success")
    if not email_ok:
        flash(f"⚠️ Não foi possível enviar o e-mail de rejeição para {sol.email}.", "warning")
    return redirect(url_for("admin_solicitacoes"))


# =========================
# ADMIN — RELEASES / AUTO-UPDATE
# =========================

@app.route("/admin/releases", methods=["GET", "POST"])
@login_required
@perfil_required("ADMIN")
def admin_releases():
    if request.method == "POST":
        action = (request.form.get("action") or "").strip()

        # ── Upload de novo release ──────────────────────────────────────────
        if action == "upload":
            versao  = (request.form.get("versao") or "").strip()
            arquivo = request.files.get("exe_file")
            ativar  = request.form.get("ativar_agora") == "1"

            if not versao:
                flash("Informe a versão antes de publicar.", "danger")
                return redirect(url_for("admin_releases"))
            if not arquivo or not arquivo.filename:
                flash("Selecione o arquivo .zip (pacote one-folder) ou .exe para upload.", "danger")
                return redirect(url_for("admin_releases"))
            _ext = os.path.splitext(arquivo.filename)[1].lower()
            if _ext not in (".zip", ".exe"):
                flash("Formato inválido. Envie o pacote .zip do build one-folder (ou um .exe legado).", "danger")
                return redirect(url_for("admin_releases"))

            exe_data      = arquivo.read()
            nome_arquivo  = arquivo.filename
            tamanho_bytes = len(exe_data)

            if ativar:
                db.session.query(AppRelease).update({"ativo": "N"})
                db.session.flush()

            novo = AppRelease(
                versao        = versao,
                nome_arquivo  = nome_arquivo,
                tamanho_bytes = tamanho_bytes,
                exe_blob      = exe_data,
                publicado_por = session.get("user_nome"),
                ativo         = "S" if ativar else "N",
            )
            db.session.add(novo)
            db.session.commit()
            flash(
                f"Versão {versao} publicada com sucesso "
                f"({_fmt_size(tamanho_bytes)}"
                + (" · marcada como ativa" if ativar else "")
                + ").",
                "success",
            )
            return redirect(url_for("admin_releases"))

        # ── Ativar release ──────────────────────────────────────────────────
        elif action == "ativar":
            rid = int(request.form.get("release_id") or 0)
            rel = AppRelease.query.get_or_404(rid)
            db.session.query(AppRelease).update({"ativo": "N"})
            rel.ativo = "S"
            db.session.commit()
            flash(f"Versão {rel.versao} definida como versão ativa.", "success")
            return redirect(url_for("admin_releases"))

        # ── Excluir release ─────────────────────────────────────────────────
        elif action == "excluir":
            rid = int(request.form.get("release_id") or 0)
            rel = AppRelease.query.get_or_404(rid)
            versao_rm = rel.versao
            db.session.delete(rel)
            db.session.commit()
            flash(f"Versão {versao_rm} excluída.", "success")
            return redirect(url_for("admin_releases"))

    releases  = AppRelease.query.order_by(AppRelease.publicado_em.desc()).all()
    pendentes = _admin_pendentes()
    return render_template(
        "admin_releases.html",
        releases  = releases,
        fmt_size  = _fmt_size,
        pendentes = pendentes,
    )


@app.route("/admin/releases/<int:rid>/download")
@login_required
@perfil_required("ADMIN")
def admin_release_download(rid):
    rel = AppRelease.query.get_or_404(rid)
    if not rel.exe_blob:
        flash("Arquivo não encontrado.", "danger")
        return redirect(url_for("admin_releases"))
    blob = rel.exe_blob
    if hasattr(blob, "read"):
        blob = blob.read()
    return send_file(
        BytesIO(blob),
        download_name = rel.nome_arquivo or "CCTV_ControlPanel.exe",
        as_attachment = True,
        mimetype      = "application/octet-stream",
    )


# =========================
# CONTROLE DE CHAVES — Blueprint
# =========================
from chaves_blueprint import chaves_bp, setup_chaves
setup_chaves(db, UsuarioSite=UsuarioSite)
app.register_blueprint(chaves_bp, url_prefix='/chaves')

# =========================
# GESTÃO DE ARMÁRIOS — Blueprint
# =========================
from armarios_blueprint import armarios_bp, setup_armarios
setup_armarios(db, UsuarioSite=UsuarioSite)
app.register_blueprint(armarios_bp, url_prefix='/armarios')

# =========================
# ACHADOS E PERDIDOS — Blueprint
# =========================
from achados_blueprint import achados_bp, setup_achados
setup_achados(db, uploads_root=UPLOADS_ROOT)
app.register_blueprint(achados_bp, url_prefix='/achados')

# =========================
# PORTAS DE EMERGÊNCIA — Blueprint
# =========================
from portas_emergencia_blueprint import pe_bp, setup_pe
setup_pe(db, get_setores=_get_setores, get_locais=_get_locais)
app.register_blueprint(pe_bp, url_prefix='/portas-emergencia')

# =========================
# ABERTURA E FECHAMENTO — Blueprint
# =========================
from site_af_blueprint import af_bp, setup_af
setup_af(db, Usuario=Usuario, UsuarioSite=UsuarioSite)
app.register_blueprint(af_bp, url_prefix='/site-af')

# =========================
# CONTROLE DE ATIVOS (redundância web do app desktop) — Blueprint
# =========================
from ativos_blueprint import ativos_bp, setup_ativos
setup_ativos(sb)   # mesmo cliente Supabase usado na tela de Usuários Ativos
app.register_blueprint(ativos_bp, url_prefix='/ativos')

# =========================
# INIT DB
# =========================
# IMPORTANTE: NÃO conectar ao Oracle aqui (nível de módulo). Isso rodava a CADA
# import do app.py — antes do Flask sequer abrir a porta — e travava a
# inicialização em ~18s via VPN. Toda a criação/migração de schema agora vive
# em _init_db(), que o launcher chama uma vez e que pula o banco por completo
# quando o marcador local de versão já bate (startup quase instantâneo).
#
# As tabelas/sequences de armário são modelos ORM → criadas pelo db.create_all()
# dentro de _init_db(); as colunas extras estão na lista de migrações de _init_db().


# Versão do schema. BUMP a cada vez que novas migrações forem adicionadas ao
# _init_db — isso força a reexecução única das DDLs em cada cliente e regrava o marcador.
_SCHEMA_VERSION = "2026.06.23"


def _init_db():
    """Aplica migrações de schema (ALTER TABLE).
    Chamado de forma síncrona pelo launcher antes de abrir a janela.
    Todas as operações são idempotentes (try/except ignora erros de 'já existe').

    Otimização de startup em 2 níveis:
      1) Marcador LOCAL (arquivo): se esta máquina já migrou esta versão de schema,
         abre a janela SEM nenhuma conexão Oracle — elimina o round-trip via VPN,
         que é o que mais demora antes da janela aparecer.
      2) Marcador no banco (APP_SCHEMA_META): backup do nível 1; evita reexecutar
         ~120 DDLs mesmo se o arquivo local não existir.
    """
    _marker_file = os.path.join(_log_dir, ".schema_version")

    # ── Nível 1: marcador local — zero conexão ao banco no caso comum ──
    try:
        with open(_marker_file, "r", encoding="utf-8") as _mf:
            if _mf.read().strip() == _SCHEMA_VERSION:
                return
    except Exception:
        pass

    def _gravar_marker_local():
        try:
            with open(_marker_file, "w", encoding="utf-8") as _mf:
                _mf.write(_SCHEMA_VERSION)
        except Exception:
            pass

    with app.app_context():
        # ── Nível 2: marcador no banco — pula as migrações ──
        try:
            row = db.session.execute(db.text(
                "SELECT VALOR FROM APP_SCHEMA_META WHERE CHAVE = 'schema_version'"
            )).fetchone()
            if row and (row[0] or "").strip() == _SCHEMA_VERSION:
                _gravar_marker_local()   # grava o arquivo p/ pular o banco na próxima
                return
        except Exception:
            # Tabela ainda não existe (primeira execução) — segue para migrar.
            db.session.rollback()

        try:
            db.create_all()  # idempotente — só cria o que falta
            db.session.execute(db.text(
                "ALTER TABLE USERS_LIVRO ADD (FOTO_PERFIL CLOB)"
            ))
            db.session.commit()
        except Exception:
            db.session.rollback()

        for _col_sql in [
            # Sequences das tabelas de chaves (criam o gerador de ID no Oracle)
            "CREATE SEQUENCE clav_chave_id_seq START WITH 1 INCREMENT BY 1",
            "CREATE SEQUENCE clav_retirada_id_seq START WITH 1 INCREMENT BY 1",
            # Sequences das tabelas de armários
            "CREATE SEQUENCE armario_id_seq START WITH 1 INCREMENT BY 1",
            "CREATE SEQUENCE arm_chave_id_seq START WITH 1 INCREMENT BY 1",
            # Colunas de assinatura e operador na atribuição de armários
            "ALTER TABLE armario ADD (assinatura_atribuicao CLOB)",
            "ALTER TABLE armario ADD (atribuido_por VARCHAR2(150))",
            # Coluna de assinatura na chave reserva de armários
            "ALTER TABLE armario_chave_reserva ADD (assinatura CLOB)",
            "ALTER TABLE ocorrencias_turno ADD (RESSALVA CLOB)",
            "ALTER TABLE ocorrencias_turno ADD (TEM_RESSALVA VARCHAR2(1))",
            "ALTER TABLE ocorrencias_turno ADD (ANEXO_ENTRADA CLOB)",
            "ALTER TABLE ocorrencias_turno ADD (ANEXO_ENTRADA_NOME VARCHAR2(255))",
            "ALTER TABLE USERS_LIVRO ADD (LGPD_ACEITO VARCHAR2(3))",
            "ALTER TABLE USERS_LIVRO ADD (LGPD_ACEITO_EM DATE)",
            "CREATE TABLE SOLICITACOES_CADASTRO (ID NUMBER GENERATED BY DEFAULT AS IDENTITY PRIMARY KEY, NOME VARCHAR2(120) NOT NULL, EMAIL VARCHAR2(120) NOT NULL, SITE VARCHAR2(128), STATUS VARCHAR2(20) DEFAULT 'PENDENTE', CRIADO_EM DATE DEFAULT SYSDATE)",
            "ALTER TABLE OCORRENCIAS ADD (BOLETIM_OCORRENCIA NUMBER(1) DEFAULT 0)",
            "ALTER TABLE OCORRENCIAS ADD (CUSTO VARCHAR2(50))",
            "ALTER TABLE ANCS ADD (VALOR VARCHAR2(50))",
            # Colunas adicionadas progressivamente à tabela ANCS
            "ALTER TABLE ANCS ADD (CARGO VARCHAR2(120))",
            "ALTER TABLE USERS_LIVRO ADD (CARGO VARCHAR2(120))",
            "ALTER TABLE ANCS ADD (NUMERO_ANC VARCHAR2(50))",
            "ALTER TABLE ANCS ADD (NUMERO_SITE NUMBER)",
            "ALTER TABLE ANCS ADD (IMAGEM_5 CLOB)",
            "ALTER TABLE ANCS ADD (IMAGEM_6 CLOB)",
            "ALTER TABLE ANCS ADD (INICIO_INVESTIGACAO VARCHAR2(16))",
            "ALTER TABLE ANCS ADD (FIM_INVESTIGACAO VARCHAR2(16))",
            "ALTER TABLE ANCS ADD (PLANO_ACAO_TEXTO CLOB)",
            "ALTER TABLE ANCS ADD (FECHADO_POR VARCHAR2(120))",
            "ALTER TABLE ANCS ADD (FECHADO_EM TIMESTAMP)",
            "ALTER TABLE ANCS ADD (ANEXO_FECHAMENTO_NOME VARCHAR2(255))",
            "ALTER TABLE ANCS ADD (ANEXO_FECHAMENTO BLOB)",
            "ALTER TABLE ANALISES_INVESTIGATIVAS ADD (VALOR VARCHAR2(50))",
            "CREATE TABLE SISTEMA_CONFIG (ID NUMBER GENERATED BY DEFAULT AS IDENTITY PRIMARY KEY, VERSAO_EXIGIDA VARCHAR2(20), DOWNLOAD_URL VARCHAR2(500), EXE_BLOB BLOB)",
            "INSERT INTO SISTEMA_CONFIG (VERSAO_EXIGIDA) SELECT '3.0' FROM DUAL WHERE NOT EXISTS (SELECT 1 FROM SISTEMA_CONFIG)",
            "UPDATE SISTEMA_CONFIG SET VERSAO_EXIGIDA = '3.5' WHERE VERSAO_EXIGIDA IS NULL OR VERSAO_EXIGIDA < '3.5'",
            "UPDATE SISTEMA_CONFIG SET VERSAO_EXIGIDA = '4.0' WHERE VERSAO_EXIGIDA IS NULL OR VERSAO_EXIGIDA < '4.0'",
            "UPDATE SISTEMA_CONFIG SET VERSAO_EXIGIDA = '4.1' WHERE VERSAO_EXIGIDA IS NULL OR VERSAO_EXIGIDA < '4.1'",
            # ── Abertura e Fechamento de Site (SF-154234) ─────────────────
            """CREATE TABLE SITE_AF (
                ID NUMBER GENERATED BY DEFAULT AS IDENTITY PRIMARY KEY,
                SITE VARCHAR2(100) NOT NULL,
                STATUS VARCHAR2(30) DEFAULT 'AGUARDANDO_ABERTURA',
                FECH_DATA VARCHAR2(10), FECH_HORA VARCHAR2(5),
                FECH_REALIZADO_POR VARCHAR2(120), FECH_AVALIADO_POR VARCHAR2(120),
                FECH_ENCAMINHADO VARCHAR2(120), FECH_ASSINATURA CLOB,
                FECH_CHECKLIST CLOB, FECH_CRIADO_EM TIMESTAMP,
                FECH_STATUS VARCHAR2(20), FECH_APROV_POR VARCHAR2(120),
                FECH_APROV_EM TIMESTAMP, FECH_APROV_OBS VARCHAR2(500),
                ABER_DATA VARCHAR2(10), ABER_HORA VARCHAR2(5),
                ABER_REALIZADO_POR VARCHAR2(120), ABER_CLIENTE VARCHAR2(120),
                ABER_UNIDADE VARCHAR2(120), ABER_ALARME_HORA VARCHAR2(5),
                ABER_ALARME_ACIONADO VARCHAR2(120), ABER_ALARME_FUNCAO VARCHAR2(120),
                ABER_ALARME_PROBLEMAS CLOB, ABER_CHECKLIST CLOB,
                ABER_CRIADO_EM TIMESTAMP, ABER_STATUS VARCHAR2(20),
                ABER_APROV_POR VARCHAR2(120), ABER_APROV_EM TIMESTAMP,
                ABER_APROV_OBS VARCHAR2(500),
                CRIADO_EM TIMESTAMP DEFAULT SYSTIMESTAMP,
                CRIADO_POR VARCHAR2(120)
            )""",
            "ALTER TABLE SITE_AF ADD (ABER_ASSINATURA CLOB)",
            "ALTER TABLE SITE_AF ADD (FECH_APROV_SIG CLOB)",
            "ALTER TABLE SITE_AF ADD (ABER_APROV_SIG CLOB)",
            """CREATE TABLE SITE_AF_ITENS (
                ID NUMBER GENERATED BY DEFAULT AS IDENTITY PRIMARY KEY,
                SITE VARCHAR2(100) NOT NULL,
                TIPO VARCHAR2(20) NOT NULL,
                NUMERO NUMBER NOT NULL,
                DESCRICAO VARCHAR2(500) NOT NULL,
                ATIVO VARCHAR2(1) DEFAULT 'S',
                CRIADO_POR VARCHAR2(120),
                CRIADO_EM TIMESTAMP DEFAULT SYSTIMESTAMP
            )""",
            # Ampliar GC para CLOB (Oracle não permite MODIFY direto — renomeia + recria)
            "ALTER TABLE OCORRENCIAS ADD (GC_TMP CLOB)",
            "UPDATE OCORRENCIAS SET GC_TMP = GC",
            "ALTER TABLE OCORRENCIAS DROP COLUMN GC",
            "ALTER TABLE OCORRENCIAS RENAME COLUMN GC_TMP TO GC",
            # Colunas de valor financeiro por tipo
            "ALTER TABLE OCORRENCIAS ADD (VALOR_RECUPERADO VARCHAR2(80))",
            "ALTER TABLE OCORRENCIAS ADD (VALOR_PREVENTIVO VARCHAR2(80))",
            "ALTER TABLE OCORRENCIAS ADD (PREJUIZO VARCHAR2(80))",
            "ALTER TABLE SISTEMA_CONFIG ADD (DOWNLOAD_URL VARCHAR2(500))",
            "ALTER TABLE SISTEMA_CONFIG ADD (EXE_BLOB BLOB)",
            "CREATE TABLE USUARIO_SITES (ID NUMBER GENERATED BY DEFAULT AS IDENTITY PRIMARY KEY, USUARIO_ID NUMBER NOT NULL, SITE_NOME VARCHAR2(128) NOT NULL)",
            "ALTER TABLE CLAVICULARIO_RETIRADA ADD (NOME_RETIRADOR VARCHAR2(150))",
            "ALTER TABLE CLAVICULARIO_RETIRADA ADD (RESPONSAVEL_ENTREGA VARCHAR2(150))",
            "ALTER TABLE CLAVICULARIO_RETIRADA ADD (ASSINATURA CLOB)",
            # Vínculos da investigação com ANC e Análise Investigativa
            "ALTER TABLE OCORRENCIAS ADD (ANC_VINCULADA_ID NUMBER)",
            "ALTER TABLE OCORRENCIAS ADD (ANALISE_VINCULADA_ID NUMBER)",
            # Tabela de releases para o sistema de auto-atualização
            "CREATE TABLE APP_RELEASES (ID NUMBER GENERATED BY DEFAULT AS IDENTITY PRIMARY KEY, VERSAO VARCHAR2(20) NOT NULL, NOME_ARQUIVO VARCHAR2(255), TAMANHO_BYTES NUMBER, EXE_BLOB BLOB, PUBLICADO_EM TIMESTAMP DEFAULT SYSTIMESTAMP, PUBLICADO_POR VARCHAR2(120), ATIVO VARCHAR2(1) DEFAULT 'N')",
            # Tabela de achados e perdidos
            "CREATE TABLE ACHADOS_PERDIDOS (ID NUMBER GENERATED BY DEFAULT AS IDENTITY PRIMARY KEY, ID_REGISTRO NUMBER NOT NULL, ID_ANTERIOR VARCHAR2(50), OBJETO VARCHAR2(200) NOT NULL, RESPONSAVEL VARCHAR2(150) NOT NULL, DATA VARCHAR2(20) NOT NULL, TURNO VARCHAR2(30) NOT NULL, DESCRICAO CLOB, FOTO_PATH VARCHAR2(500), STATUS VARCHAR2(30) DEFAULT 'Pendente', RETIRADO_POR VARCHAR2(150), SITE VARCHAR2(128), CRIADO_POR VARCHAR2(120), CREATED_AT TIMESTAMP DEFAULT SYSTIMESTAMP)",
            # Coluna de foto em base64 (CLOB) — adicionada progressivamente
            "ALTER TABLE ACHADOS_PERDIDOS ADD (FOTO_DADOS CLOB)",
            # Migração: converte TURNO de NUMBER para VARCHAR2 (caso tabela já existia com NUMBER)
            "ALTER TABLE ACHADOS_PERDIDOS ADD (TURNO_STR VARCHAR2(30))",
            "UPDATE ACHADOS_PERDIDOS SET TURNO_STR = TO_CHAR(TURNO) WHERE TURNO_STR IS NULL",
            "ALTER TABLE ACHADOS_PERDIDOS DROP COLUMN TURNO",
            "ALTER TABLE ACHADOS_PERDIDOS RENAME COLUMN TURNO_STR TO TURNO",
            # Migração: converte ID_REGISTRO de NUMBER para VARCHAR2 e adiciona NUMERO_SITE
            "ALTER TABLE ACHADOS_PERDIDOS ADD (NUMERO_SITE NUMBER)",
            "UPDATE ACHADOS_PERDIDOS SET NUMERO_SITE = ID_REGISTRO WHERE NUMERO_SITE IS NULL",
            "ALTER TABLE ACHADOS_PERDIDOS ADD (ID_REGISTRO_STR VARCHAR2(50))",
            "UPDATE ACHADOS_PERDIDOS SET ID_REGISTRO_STR = TO_CHAR(ID_REGISTRO) WHERE ID_REGISTRO_STR IS NULL",
            "ALTER TABLE ACHADOS_PERDIDOS DROP COLUMN ID_REGISTRO",
            "ALTER TABLE ACHADOS_PERDIDOS RENAME COLUMN ID_REGISTRO_STR TO ID_REGISTRO",
            # Coluna flag de foto (evita 2ª query Oracle no login)
            "ALTER TABLE USERS_LIVRO ADD (TEM_FOTO VARCHAR2(1) DEFAULT 'N')",
            # Backfill: marca S para quem já tem foto
            "UPDATE USERS_LIVRO SET TEM_FOTO = 'S' WHERE FOTO_PERFIL IS NOT NULL AND TEM_FOTO != 'S'",
            # Tabelas de configuração de campos dinâmicos por site
            "CREATE TABLE NATUREZA_CONFIG (ID NUMBER GENERATED BY DEFAULT AS IDENTITY PRIMARY KEY, NATUREZA VARCHAR2(200) NOT NULL, SITE VARCHAR2(128) NOT NULL, DATA_CRIACAO TIMESTAMP DEFAULT SYSTIMESTAMP, USUARIO_CRIADOR VARCHAR2(120))",
            "CREATE TABLE LOCAL_CONFIG (ID NUMBER GENERATED BY DEFAULT AS IDENTITY PRIMARY KEY, LOCAL VARCHAR2(200) NOT NULL, SITE VARCHAR2(128) NOT NULL, DATA_CRIACAO TIMESTAMP DEFAULT SYSTIMESTAMP, USUARIO_CRIADOR VARCHAR2(120))",
            "CREATE TABLE SETOR_CONFIG (ID NUMBER GENERATED BY DEFAULT AS IDENTITY PRIMARY KEY, SETOR VARCHAR2(200) NOT NULL, SITE VARCHAR2(128) NOT NULL, DATA_CRIACAO TIMESTAMP DEFAULT SYSTIMESTAMP, USUARIO_CRIADOR VARCHAR2(120))",
            # Portas de Emergência
            "CREATE TABLE PORTAS_EMERGENCIA (ID NUMBER GENERATED BY DEFAULT AS IDENTITY PRIMARY KEY, CODIGO VARCHAR2(50) NOT NULL, LOCALIZACAO VARCHAR2(160) NOT NULL, SETOR VARCHAR2(120), ROTA_FUGA VARCHAR2(120), RESPONSAVEL VARCHAR2(120), OBSERVACAO CLOB, ATIVO NUMBER(1) DEFAULT 1, SITE VARCHAR2(128), CRIADO_POR VARCHAR2(120), CRIADO_EM TIMESTAMP DEFAULT SYSTIMESTAMP)",
            "CREATE TABLE CHECKLISTS_PORTAS_EMERGENCIA (ID NUMBER GENERATED BY DEFAULT AS IDENTITY PRIMARY KEY, PORTA_ID NUMBER NOT NULL, DATA_CHECKLIST TIMESTAMP DEFAULT SYSTIMESTAMP, INSPETOR VARCHAR2(120) NOT NULL, TURNO VARCHAR2(30), PORTA_DESOBSTRUIDA NUMBER(1) DEFAULT 0, ABRE_NORMALMENTE NUMBER(1) DEFAULT 0, SINALIZACAO_OK NUMBER(1) DEFAULT 0, ILUMINACAO_OK NUMBER(1) DEFAULT 0, ALARME_OK NUMBER(1) DEFAULT 0, STATUS VARCHAR2(30) DEFAULT 'PENDENTE', OBSERVACAO CLOB, CRIADO_EM TIMESTAMP DEFAULT SYSTIMESTAMP)",
            "CREATE TABLE DISPAROS_ALARME (ID NUMBER GENERATED BY DEFAULT AS IDENTITY PRIMARY KEY, DATA_REGISTRO TIMESTAMP DEFAULT SYSTIMESTAMP, DATA_DISPARO VARCHAR2(20) NOT NULL, HORA_DISPARO VARCHAR2(10) NOT NULL, HORA_DESATIVADO VARCHAR2(10), HORA_ATIVADO VARCHAR2(10), CONTATO_MONITORAMENTO VARCHAR2(150), LOCALIZACAO VARCHAR2(150) NOT NULL, SETOR VARCHAR2(100), TIPO_ALARME VARCHAR2(100) NOT NULL, MOTIVO VARCHAR2(150), RESPONSAVEL VARCHAR2(120), TURNO VARCHAR2(50), HOUVE_EVACUACAO NUMBER(1) DEFAULT 0, ACIONADO_BOMBEIRO NUMBER(1) DEFAULT 0, ACIONADO_SEGURANCA NUMBER(1) DEFAULT 0, STATUS VARCHAR2(50) DEFAULT 'EM ANALISE', OBSERVACAO CLOB, SITE VARCHAR2(128), CRIADO_POR VARCHAR2(120))",
            "CREATE TABLE BOTOES_PANICO (ID NUMBER GENERATED BY DEFAULT AS IDENTITY PRIMARY KEY, DATA_REGISTRO TIMESTAMP DEFAULT SYSTIMESTAMP, CODIGO VARCHAR2(50) NOT NULL, LOCALIZACAO VARCHAR2(150) NOT NULL, SETOR VARCHAR2(100), TIPO VARCHAR2(80), RESPONSAVEL VARCHAR2(120), TURNO VARCHAR2(50), TESTADO NUMBER(1) DEFAULT 0, SINAL_RECEBIDO NUMBER(1) DEFAULT 0, COMUNICACAO_CFTV NUMBER(1) DEFAULT 0, NECESSITA_MANUTENCAO NUMBER(1) DEFAULT 0, STATUS VARCHAR2(50) DEFAULT 'EM ANALISE', OBSERVACAO CLOB, SITE VARCHAR2(128), CRIADO_POR VARCHAR2(120))",
            # Colunas de tratativa de NC em checklists
            "ALTER TABLE CHECKLISTS_PORTAS_EMERGENCIA ADD (CONCLUSAO CLOB)",
            "ALTER TABLE CHECKLISTS_PORTAS_EMERGENCIA ADD (DATA_CONCLUSAO TIMESTAMP)",
            "ALTER TABLE CHECKLISTS_PORTAS_EMERGENCIA ADD (CONCLUIDO_POR VARCHAR2(120))",
            # Novos campos de botão de pânico
            "ALTER TABLE BOTOES_PANICO ADD (HORA_TESTE VARCHAR2(10))",
            "ALTER TABLE BOTOES_PANICO ADD (HORA_RETORNO VARCHAR2(10))",
            "ALTER TABLE BOTOES_PANICO ADD (AGENTE_CFTV VARCHAR2(120))",
            # Corrige registros ANC com status NULL ou vazio → ABERTO
            "UPDATE ANCS SET STATUS = 'ABERTO' WHERE STATUS IS NULL OR TRIM(STATUS) = ''",
            # Normaliza CONCLUIDO (sem acento) → CONCLUÍDO (com acento)
            "UPDATE ANCS SET STATUS = 'CONCLUÍDO' WHERE STATUS = 'CONCLUIDO'",
            # Corrige registros ANC com status inválido (qualquer valor diferente dos 3 válidos)
            "UPDATE ANCS SET STATUS = 'ABERTO' WHERE STATUS NOT IN ('ABERTO','EM ANDAMENTO','CONCLUÍDO')",
            # ── Soft-delete de ANCs ──────────────────────────────────────
            "ALTER TABLE ANCS ADD (EXCLUIDO VARCHAR2(1) DEFAULT 'N')",
            "UPDATE ANCS SET EXCLUIDO = 'N' WHERE EXCLUIDO IS NULL",
            "ALTER TABLE ANCS ADD (EXCL_STATUS VARCHAR2(20))",
            "ALTER TABLE ANCS ADD (EXCL_SOLICITADO_POR VARCHAR2(120))",
            "ALTER TABLE ANCS ADD (EXCL_SOLICITADO_EM TIMESTAMP)",
            "ALTER TABLE ANCS ADD (EXCL_SOLICITANTE_EMAIL VARCHAR2(120))",
            "ALTER TABLE ANCS ADD (EXCL_MOTIVO VARCHAR2(500))",
            "ALTER TABLE ANCS ADD (EXCL_ADMIN_POR VARCHAR2(120))",
            "ALTER TABLE ANCS ADD (EXCL_ADMIN_EM TIMESTAMP)",
            "ALTER TABLE ANCS ADD (EXCL_ADMIN_MOTIVO VARCHAR2(500))",
            # Garante que GC existe como CLOB (migration 7546-7549 pode ter falhado
            # em bancos que nunca tiveram GC como VARCHAR2, deixando a coluna ausente)
            "ALTER TABLE OCORRENCIAS ADD (GC CLOB)",
            # Limpa coluna temporária caso tenha ficado órfã da migration anterior
            "ALTER TABLE OCORRENCIAS DROP COLUMN GC_TMP",
        ]:
            try:
                db.session.execute(db.text(_col_sql))
                db.session.commit()
            except Exception:
                db.session.rollback()

        # ── Popula NATUREZAS_OCORRENCIA com os valores padrão (idempotente) ──
        try:
            existentes = {n.nome for n in NaturezaOcorrencia.query.all()}
            for nome in _NATUREZAS_PADRAO:
                if nome not in existentes:
                    db.session.add(NaturezaOcorrencia(nome=nome))
            db.session.commit()
        except Exception:
            db.session.rollback()

        # ── Grava o marcador de versão do schema (cria a tabela se preciso) ──
        # A partir daqui, próximas aberturas pulam todo o bloco acima (fast-path).
        try:
            db.session.execute(db.text(
                "CREATE TABLE APP_SCHEMA_META (CHAVE VARCHAR2(50) PRIMARY KEY, VALOR VARCHAR2(50))"
            ))
            db.session.commit()
        except Exception:
            db.session.rollback()
        try:
            db.session.execute(db.text("DELETE FROM APP_SCHEMA_META WHERE CHAVE = 'schema_version'"))
            db.session.execute(db.text(
                "INSERT INTO APP_SCHEMA_META (CHAVE, VALOR) VALUES ('schema_version', :v)"
            ), {"v": _SCHEMA_VERSION})
            db.session.commit()
        except Exception:
            db.session.rollback()

    # Migrações concluídas nesta máquina → grava o marcador local para que as
    # próximas aberturas pulem totalmente o banco (nível 1).
    _gravar_marker_local()


if __name__ == "__main__":
    _init_db()
    app.run(debug=True)


